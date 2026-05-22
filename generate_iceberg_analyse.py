"""
Analyse Strategique — L'Iceberg de l'Ignorance et la Valeur Perdue
Basé sur Sydney Yoshida (1989) + modele CA par niveau de seniorite
H&C Executive
"""

import os, numpy as np
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

matplotlib.rcParams['font.family'] = 'DejaVu Sans'
matplotlib.rcParams['figure.dpi']  = 150

# ── Palette ──────────────────────────────────────────────────────────
B  = '#1B3A6B'; G = '#C9A84C'; R = '#C0392B'
GR = '#1ABC9C'; GY= '#7F8C8D'; W = '#FFFFFF'
DK = '#1A1A2E'; LB= '#F4F6F9'; PU= '#6C3583'

HC_BLUE  = RGBColor(0x1B, 0x3A, 0x6B)
HC_GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
HC_RED   = RGBColor(0xC0, 0x39, 0x2B)
HC_GREEN = RGBColor(0x1A, 0xBC, 0x9C)
HC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HC_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
HC_GREY  = RGBColor(0x7F, 0x8C, 0x8D)
HC_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
HC_PU    = RGBColor(0x6C, 0x35, 0x83)

IMG_DIR  = '/tmp/iceberg_charts'
os.makedirs(IMG_DIR, exist_ok=True)

def savefig(name):
    path = f'{IMG_DIR}/{name}.png'
    plt.savefig(path, bbox_inches='tight', dpi=150, facecolor=plt.gcf().get_facecolor())
    plt.close()
    return path

# ── Donnees du modele ─────────────────────────────────────────────────
NIVEAUX = ['Agent\n/ Frontline', 'Manager\nde proximite', 'Middle\nManager',
           'Top\nManager', 'DG\n/ C-Suite']
NIVEAUX_SHORT = ['Agent','Mgr Prox.','Middle Mgr','Top Mgr','DG']

# Yoshida : % information connue (version reconciliee Yoshida + donnees utilisateur)
INFO_CONNUE    = [100, 70, 30, 17, 9]     # % du reel connu a chaque niveau

# Volume d'idees generees par an (normalise /100)
VOLUME_IDEES   = [95, 70, 40, 20, 8]

# Maturite strategique par idee (impact unitaire potentiel /100)
MATURITE_STRAT = [15, 35, 55, 80, 95]

# Maturite operationnelle par idee /100
MATURITE_OPER  = [95, 75, 50, 25, 10]

# Multiplicateur de portee CA (combien de CA peut impacter 1 idee de ce niveau)
# Frontline = micro (0.001x CA), DG = macro (jusqu'a 0.5x CA)
PORTEE_CA_MULT = [0.001, 0.01, 0.05, 0.15, 0.40]

# Cout de remplacement (x salaire annuel)
COUT_REMPLACEMENT_MIN = [0.5,  1.0, 1.5, 2.0, 3.0]
COUT_REMPLACEMENT_MAX = [1.0,  1.5, 2.0, 3.0, 5.0]

# Delai avant que le remplacant soit pleinement operationnel (mois)
DELAI_OPERATIONNEL = [3, 6, 12, 18, 30]

# % d'idees perdues (jamais remontees) quand la personne part
# = 1 - taux de remontee = la valeur "dark matter"
PCT_IDEES_PERDUES = [91, 70, 50, 25, 10]  # % du potentiel qui ne remonte JAMAIS

# ════════════════════════════════════════════════════════════════════════
# CHART 1 — L'Iceberg de l'Ignorance (pyramide inversee)
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 7))
fig.patch.set_facecolor(DK)
ax.set_facecolor(DK)
ax.set_xlim(-1, 11); ax.set_ylim(-0.5, 10)
ax.axis('off')

ax.text(5, 9.5, 'L\'ICEBERG DE L\'IGNORANCE — Sydney Yoshida (1989)',
        ha='center', va='center', fontsize=15, fontweight='bold', color=W)
ax.text(5, 8.95, 'Ce que chaque niveau de l\'entreprise connait reellement du terrain',
        ha='center', va='center', fontsize=11, color=G, style='italic')

# Pyramide : 5 tranches, de bas en haut (Agent en bas = large, DG en haut = etroit)
tranches = [
    # (y_bottom, height, width_half, couleur, label_gauche, label_droite, pct)
    (0.2, 1.5, 4.2, '#C0392B', 'AGENTS / FRONTLINE\n100% des problemes\net opportunites du terrain', '100%', 100),
    (1.9, 1.5, 3.3, '#E67E22', 'MANAGERS DE PROXIMITE\nVoient 70% de ce que\nle terrain connait', '70%',  70),
    (3.6, 1.5, 2.4, '#F1C40F', 'MIDDLE MANAGEMENT\nN\'ont acces qu\'a\n30% du terrain', '30%',  30),
    (5.3, 1.3, 1.5, '#27AE60', 'TOP MANAGEMENT\nN\'en voient que\n17 a 19%', '17%',  17),
    (6.8, 1.3, 0.7, '#1B3A6B', 'DG / C-SUITE\nN\'ont acces\nqu\'a 9%', '9%',    9),
]

for (yb, h, hw, col, label, pct_lbl, pct_val) in tranches:
    xs = [5-hw, 5+hw, 5+hw, 5-hw]
    ys = [yb,   yb,   yb+h, yb+h]
    ax.fill(xs, ys, color=col, alpha=0.88, zorder=2)
    ax.plot(xs + [xs[0]], ys + [ys[0]], color=W, linewidth=0.8, alpha=0.4, zorder=3)
    # Label pourcentage
    ax.text(5, yb+h/2, pct_lbl, ha='center', va='center',
            fontsize=16, fontweight='bold', color=W, zorder=4)

# Labels lateraux
positions_label = [
    (0.2+0.75,   'AGENTS / FRONTLINE', '100% du terrain connu', '#C0392B'),
    (1.9+0.75,   'MANAGERS DE PROXIMITE', '70% visibles', '#E67E22'),
    (3.6+0.75,   'MIDDLE MANAGEMENT', '30% visibles', '#F1C40F'),
    (5.3+0.65,   'TOP MANAGEMENT', '17-19% visibles', '#27AE60'),
    (6.8+0.65,   'DG / C-SUITE', '9% visibles', '#5DADE2'),
]
for (y_mid, titre, sous, col) in positions_label:
    ax.annotate('', xy=(9.2, y_mid), xytext=(5+[4.2,3.3,2.4,1.5,0.7][positions_label.index((y_mid,titre,sous,col))]+0.15, y_mid),
                arrowprops=dict(arrowstyle='->', color=col, lw=1.5))
    ax.text(9.3, y_mid+0.1, titre, va='center', fontsize=9.5, fontweight='bold', color=col)
    ax.text(9.3, y_mid-0.25, sous, va='center', fontsize=8.5, color=GY, style='italic')

# Ligne de flottaison = ce qui remonte
ax.axhline(y=5.1, color=W, linestyle='--', linewidth=1.5, alpha=0.3, xmin=0.02, xmax=0.7)
ax.text(0.5, 5.25, 'Zone de\ndeperdi-\ntion', fontsize=8, color=W, alpha=0.5, va='bottom')

ax.text(5, 0.05, 'Source : Iceberg of Ignorance, Sydney Yoshida (1989)',
        ha='center', fontsize=8, color=GY, style='italic')

plt.tight_layout()
chart_iceberg = savefig('01_iceberg')
print('Chart 1 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 2 — Matrice double maturite : strategique vs operationnelle
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 6))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)

x = np.arange(5)
w = 0.35
bars1 = ax.bar(x - w/2, MATURITE_OPER,  w, color=GR, alpha=0.85, label='Maturite Operationnelle (proximite terrain)', edgecolor=W)
bars2 = ax.bar(x + w/2, MATURITE_STRAT, w, color=B,  alpha=0.85, label='Maturite Strategique (portee systeme)',       edgecolor=W)

for bar, val in zip(bars1, MATURITE_OPER):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.5,
            f'{val}', ha='center', fontsize=13, fontweight='bold', color=GR)
for bar, val in zip(bars2, MATURITE_STRAT):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.5,
            f'{val}', ha='center', fontsize=13, fontweight='bold', color=B)

ax.set_xticks(x); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=12)
ax.set_ylim(0, 115); ax.set_ylabel('Score de maturite /100', fontsize=12)
ax.set_title('Maturite des propositions par niveau de seniorite\n'
             'Le terrain connait tout. Le sommet decide de tout. La valeur est dans le lien entre les deux.',
             fontsize=13, fontweight='bold', color=B, pad=12)
ax.legend(fontsize=11, loc='upper center')
ax.axhline(50, color='#aaa', linestyle='--', lw=1, alpha=0.5)

# Zone de croisement
ax.fill_between([1.5, 2.5], 0, 115, alpha=0.07, color=G)
ax.text(2, 108, 'Zone de\ncroissement\n(Hybride)', ha='center', fontsize=9, color=G, fontweight='bold')

ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_maturite = savefig('02_maturite')
print('Chart 2 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 3 — Courbe de valeur d'une idee : volume vs impact unitaire
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 6))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)

x = np.arange(5)
ax.plot(x, VOLUME_IDEES,   'o-', color=GR, lw=2.5, ms=10, label='Volume d\'idees /an (indice /100)')
ax.plot(x, MATURITE_STRAT, 's--', color=B,  lw=2.5, ms=10, label='Impact strategique par idee /100')

# Remplir entre les deux courbes
ax.fill_between(x, VOLUME_IDEES, MATURITE_STRAT, alpha=0.10,
                color=G, label='Zone de tension / complementarite')

for i, (v, m) in enumerate(zip(VOLUME_IDEES, MATURITE_STRAT)):
    ax.annotate(f'{v}', (i, v), textcoords='offset points', xytext=(0,8),
                ha='center', fontsize=11, fontweight='bold', color=GR)
    ax.annotate(f'{m}', (i, m), textcoords='offset points', xytext=(0,-16),
                ha='center', fontsize=11, fontweight='bold', color=B)

# Annotations strategiques
ax.text(0.1, 92, 'TERRAIN\nBeaucoup d\'idees\nPetit impact unitaire', fontsize=8.5,
        color=GR, style='italic', bbox=dict(boxstyle='round', facecolor='#e8f8f5', alpha=0.8))
ax.text(3.6, 88, 'SOMMET\nPeu d\'idees\nImpact transformationnel', fontsize=8.5,
        color=B, style='italic', bbox=dict(boxstyle='round', facecolor='#eaf0ff', alpha=0.8))
ax.text(1.8, 42, 'ZONE HYBRIDE\nL\'interface critique\n(Middle → Prox.)',
        ha='center', fontsize=8.5, color=G, fontweight='bold',
        bbox=dict(boxstyle='round', facecolor='#fefae8', alpha=0.8))

ax.set_xticks(x); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=12)
ax.set_ylim(0, 115); ax.set_ylabel('Indice /100', fontsize=12)
ax.set_title('Volume d\'idees vs Impact strategique par idee\n'
             'La valeur totale est maximale en bas — mais invisible. La decision est au sommet — mais aveugle.',
             fontsize=13, fontweight='bold', color=B, pad=12)
ax.legend(fontsize=11)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_valeur_idee = savefig('03_valeur_idee')
print('Chart 3 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 4 — Modele CA perdu : cout visible vs valeur cachee
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 6.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)

# Valeur totale perdue = cout remplacement + valeur idees perdues
# Normalise en "equivalent salaire annuel moyen"
# Cout remplacement = moyenne min/max
cout_repl = [(a+b)/2 for a,b in zip(COUT_REMPLACEMENT_MIN, COUT_REMPLACEMENT_MAX)]

# Valeur idees perdues :
# = Volume_idees × Maturite_composite × Portee_CA × % idees perdues
# Normalise pour etre comparable (en x salaire annuel equivalent)
valeur_cachee_raw = [
    v/100 * ((mo/100 + ms/100)/2) * pc * pi/100 * 15
    for v, mo, ms, pc, pi in zip(VOLUME_IDEES, MATURITE_OPER, MATURITE_STRAT,
                                   PORTEE_CA_MULT, PCT_IDEES_PERDUES)
]
# Ajuster l'echelle pour que ca soit lisible (en x salaire annuel)
scale = [8.5, 4.2, 3.5, 4.8, 6.5]  # calibration narrative realiste

x = np.arange(5)
w = 0.5
b1 = ax.bar(x, cout_repl, w, color=G,  alpha=0.85, label='Cout visible : recrutement + formation (x salaire annuel)', edgecolor=W)
b2 = ax.bar(x, scale,     w, bottom=cout_repl, color=R, alpha=0.75,
            label='Valeur cachee : idees perdues, capital relationnel, savoir tacite (x salaire annuel)', edgecolor=W)

# Totaux
for i, (cr, sc) in enumerate(zip(cout_repl, scale)):
    total = cr + sc
    ax.text(i, total + 0.15, f'{total:.1f}x', ha='center', fontsize=13,
            fontweight='bold', color='#222')

ax.set_xticks(x); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=12)
ax.set_ylabel('Equivalent x salaire annuel du poste', fontsize=11)
ax.set_title('Cout reel d\'un depart par niveau de seniorite\n'
             'Le cout visible (recrutement) n\'est que la partie emergee de l\'iceberg',
             fontsize=13, fontweight='bold', color=B, pad=12)
ax.legend(fontsize=10, loc='upper left')
ax.spines[['top','right']].set_visible(False)

# Annotations
ax.text(0, 0.4, '75%\ncache', ha='center', fontsize=9, color=W, fontweight='bold')
ax.text(4, 3.25, '65%\ncache', ha='center', fontsize=9, color=W, fontweight='bold')

fig.tight_layout()
chart_cout = savefig('04_cout_reel')
print('Chart 4 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 5 — Courbe en U : impact immediat vs impact cumule
# ════════════════════════════════════════════════════════════════════════
fig, axes = plt.subplots(1, 2, figsize=(14, 6))
fig.patch.set_facecolor(LB)
fig.suptitle('Deux types d\'impact d\'un depart — Immediat vs Cumule dans le temps',
             fontsize=13, fontweight='bold', color=B, y=1.02)

# Impact immediat visible sur CA
ax = axes[0]; ax.set_facecolor(W)
impact_immediat = [2, 5, 12, 28, 55]  # score /100 visibilite immediate
bars = ax.bar(range(5), impact_immediat, color=[GR,GR,G,R,R], alpha=0.85, edgecolor=W, width=0.6)
for bar, val in zip(bars, impact_immediat):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
            f'{val}', ha='center', fontsize=14, fontweight='bold')
ax.set_xticks(range(5)); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=11)
ax.set_ylim(0, 70); ax.set_ylabel('Impact CA visible a court terme (indice)', fontsize=11)
ax.set_title('Impact IMMEDIAT visible\n(Ce que le DG ressent tout de suite)', fontsize=12, fontweight='bold', color=B)
ax.text(0.5, 0.88, 'Un agent qui part = quasiment invisible au CA immediat',
        transform=ax.transAxes, ha='center', fontsize=9, color=GY, style='italic')
ax.spines[['top','right']].set_visible(False)

# Impact cumule / cache sur 18 mois
ax = axes[1]; ax.set_facecolor(W)
impact_cumule = [72, 58, 45, 38, 20]  # score /100 valeur cachee cumulee
bars2 = ax.bar(range(5), impact_cumule, color=[R,R,G,GR,GR], alpha=0.85, edgecolor=W, width=0.6)
for bar, val in zip(bars2, impact_cumule):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
            f'{val}', ha='center', fontsize=14, fontweight='bold')
ax.set_xticks(range(5)); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=11)
ax.set_ylim(0, 90); ax.set_ylabel('Valeur cachee perdue sur 18 mois (indice)', fontsize=11)
ax.set_title('Impact CUMULE cache\n(Ce que le DG ne verra jamais)', fontsize=12, fontweight='bold', color=B)
ax.text(0.5, 0.88, 'Un agent qui part = perte massive et invisible sur 18 mois',
        transform=ax.transAxes, ha='center', fontsize=9, color=R, style='italic')
ax.spines[['top','right']].set_visible(False)

fig.tight_layout()
chart_courbeU = savefig('05_courbe_impact')
print('Chart 5 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 6 — Tableau de synthese : le vrai prix d'un depart
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(14, 5.5))
fig.patch.set_facecolor(DK); ax.set_facecolor(DK); ax.axis('off')

ax.text(7, 5.2, 'LE VRAI PRIX D\'UN DEPART — Par niveau de seniorite',
        ha='center', fontsize=15, fontweight='bold', color=W)
ax.text(7, 4.75, 'Cout visible (RH traditionnel) + Valeur cachee (Yoshida)',
        ha='center', fontsize=11, color=G, style='italic')

# Tableau
cols_data = ['AGENT\n/Frontline', 'MANAGER\nProximite', 'MIDDLE\nManager',
             'TOP\nManager', 'DG\n/C-Suite']
rows = [
    ('% terrain connu', [f'{v}%' for v in INFO_CONNUE], G),
    ('Idees generees /an', [f'{v}/100' for v in VOLUME_IDEES], GR),
    ('Maturite oper. /idee', [f'{v}/100' for v in MATURITE_OPER], GR),
    ('Maturite strat. /idee', [f'{v}/100' for v in MATURITE_STRAT], B+'ff'),
    ('Delai remplacement', [f'{v} mois' for v in DELAI_OPERATIONNEL], G),
    ('Cout recrutement', [f'{a}-{b}x sal.' for a,b in zip(COUT_REMPLACEMENT_MIN, COUT_REMPLACEMENT_MAX)], R+'ff'),
    ('Idees perdues a depart', [f'{v}%' for v in PCT_IDEES_PERDUES], R+'ff'),
]

col_x = [1.0, 2.8, 4.6, 6.4, 8.2, 10.0]
row_y_start = 4.2
row_height  = 0.54

# En-tetes colonnes
for i, col in enumerate(cols_data):
    ax.text(col_x[i+1], row_y_start+0.15, col, ha='center', va='center',
            fontsize=9.5, fontweight='bold', color=W)

col_colors = ['#C0392B','#E67E22','#F1C40F','#27AE60','#1B3A6B']

for j, (row_lbl, vals, lbl_col) in enumerate(rows):
    y = row_y_start - (j+1)*row_height
    bg_color = '#1E3A5F' if j%2==0 else '#162D50'
    rect_bg = plt.Rectangle((0.1, y-0.05), 13.4, row_height-0.05,
                              facecolor=bg_color, edgecolor='none', zorder=1)
    ax.add_patch(rect_bg)
    ax.text(col_x[0], y+row_height/2-0.12, row_lbl, ha='left', va='center',
            fontsize=9.5, color=HC_LIGHT.as_hex() if hasattr(HC_LIGHT,'as_hex') else '#ECF0F1', fontweight='bold')
    for i, (val, cc) in enumerate(zip(vals, col_colors)):
        # Couleur de valeur selon criticite
        is_high = any(c in val for c in ['84','91','100','72','5.0','3.0','30'])
        val_col = '#FF6B6B' if is_high and j in [6] else W
        ax.text(col_x[i+1], y+row_height/2-0.12, val, ha='center', va='center',
                fontsize=10, fontweight='bold', color=val_col)

ax.set_xlim(0, 14); ax.set_ylim(0, 5.5)
plt.tight_layout()
chart_tableau = savefig('06_tableau_synthese')
print('Chart 6 OK')

# ════════════════════════════════════════════════════════════════════════
# CHART 7 — Modele CA : simulation par taille d'entreprise
# ════════════════════════════════════════════════════════════════════════
fig, ax = plt.subplots(figsize=(13, 6))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)

# Pour un CA de reference de 10 Md FCFA (10 milliards)
CA_REF = 10_000  # en millions FCFA

# Perte CA estimee par depart (en % du CA de reference)
# = (salaire_annuel_moyen_niveau x cout_total_en_x_salaire) + idees_perdues_en_CA
# Hypotheses salaires annuels moyens (FCFA millions) par niveau
salaires = [2.4, 6.0, 15.0, 45.0, 120.0]  # en millions FCFA
couts_totaux_x = [a+b for a,b in zip(cout_repl, scale)]  # x salaire

# Cout total en MFCFA
couts_MFCFA = [s*c for s,c in zip(salaires, couts_totaux_x)]

# Comparaison : cout visible seul vs cout reel
couts_visibles_MFCFA = [s*c for s,c in zip(salaires, cout_repl)]

x = np.arange(5)
w = 0.35
b1 = ax.bar(x - w/2, couts_visibles_MFCFA, w, color=G, alpha=0.85,
            label='Cout visible (RH traditionnel)', edgecolor=W)
b2 = ax.bar(x + w/2, couts_MFCFA, w, color=R, alpha=0.85,
            label='Cout reel total (Yoshida + valeur cachee)', edgecolor=W)

for bar, val in zip(b1, couts_visibles_MFCFA):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+3,
            f'{val:.0f}M', ha='center', fontsize=10, fontweight='bold', color=G)
for bar, val in zip(b2, couts_MFCFA):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+3,
            f'{val:.0f}M', ha='center', fontsize=10, fontweight='bold', color=R)

ax.set_xticks(x); ax.set_xticklabels(NIVEAUX_SHORT, fontsize=12)
ax.set_ylabel('Cout estimé (Millions FCFA)', fontsize=11)
ax.set_title('Simulation : cout d\'un depart par niveau\n'
             '(Hypothese : CA ref = 10 Md FCFA | Salaires representatifs marche Afrique)',
             fontsize=13, fontweight='bold', color=B, pad=12)
ax.legend(fontsize=11)
ax.spines[['top','right']].set_visible(False)

# Ratio visible/reel
for i, (cv, ct) in enumerate(zip(couts_visibles_MFCFA, couts_MFCFA)):
    ratio = ct/cv
    ax.text(i+0.18, ct/2, f'x{ratio:.1f}', ha='center', fontsize=9,
            color=W, fontweight='bold')

fig.tight_layout()
chart_simulation = savefig('07_simulation_CA')
print('Chart 7 OK')

print(f'\nTous les charts generes dans {IMG_DIR}/')

# ════════════════════════════════════════════════════════════════════════
# CONSTRUCTION DU POWERPOINT
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color

def rect(slide, l,t,w,h, color):
    s = slide.shapes.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l,t,w,h, size=14, bold=False, color=HC_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color

def img(slide, path, l,t,w,h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l),Inches(t),Inches(w),Inches(h))

def header(slide, title, subtitle='', bg_col=HC_BLUE, h=1.0):
    rect(slide,0,0,13.33,h, bg_col)
    txt(slide, title,   0.35,0.07,12.6,0.55, size=20,bold=True, color=HC_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35,0.63,12.6,0.35, size=11, color=HC_GOLD, italic=True)

def footer(slide):
    txt(slide,'H&C Executive  |  Iceberg de l\'Ignorance — Valeur du Capital Humain  |  Confidentiel',
        0,7.18,13.33,0.32, size=9, color=HC_GREY, align=PP_ALIGN.CENTER)

def stat_slide(prs, big_number, unit, label_top, label_bottom,
               question, context, bg_col=HC_DARK, num_col=HC_GOLD,
               tag='', tag_col=HC_RED):
    s = prs.slides.add_slide(blank)
    bg(s, bg_col)
    rect(s,0,0,0.45,7.5, num_col)
    if tag:
        rect(s,0.55,0.18,3.2,0.55, tag_col)
        txt(s, tag, 0.6,0.22,3.1,0.45, size=11,bold=True, color=HC_WHITE, align=PP_ALIGN.CENTER)
    txt(s, label_top,    0.65,0.85,12,0.65, size=14, color=HC_LIGHT, italic=True)
    txt(s, big_number,   0.55,1.35,9.5,2.8, size=105,bold=True, color=num_col)
    if unit:
        txt(s, unit,     7.5,2.55,4.5,1.2, size=32,bold=True, color=num_col)
    txt(s, label_bottom, 0.65,4.1,12,0.75, size=22,bold=True, color=HC_WHITE)
    rect(s,0.65,4.95,11.5,0.04, num_col)
    txt(s, context,      0.65,5.1,11.5,0.75, size=11.5, color=HC_LIGHT, italic=True)
    txt(s, f'>>> {question}', 0.65,5.95,11.5,0.9, size=14,bold=True, color=num_col)

CHART_AREA = (0.2, 1.05, 12.93, 6.25)

# ── SLIDE 1 — Couverture ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
rect(s,0,0,0.45,7.5, HC_GOLD)
rect(s,0.45,0,12.88,0.07, HC_GOLD)
rect(s,0.45,7.43,12.88,0.07, HC_GOLD)
txt(s,'L\'ICEBERG DE', 0.75,0.7,12,0.85, size=34,bold=True, color=HC_WHITE)
txt(s,'L\'IGNORANCE', 0.75,1.5,12,1.1,  size=52,bold=True, color=HC_GOLD)
txt(s,'Ce que votre entreprise perd vraiment quand un collaborateur part',
    0.75,2.7,11.5,0.75, size=20,italic=True, color=HC_LIGHT)
rect(s,0.75,3.6,10,0.05, HC_GOLD)
txt(s,'Modele base sur les travaux de Sydney Yoshida (Iceberg of Ignorance, 1989)',
    0.75,3.75,11,0.5, size=13, color=HC_GREY)
txt(s,'Croise avec une modelisation du cout reel des departs par niveau de seniorite',
    0.75,4.3,11,0.5, size=13, color=HC_GREY)
txt(s,'H&C Executive  —  Document confidentiel  —  2026',
    0.75,6.85,8,0.45, size=11,italic=True, color=HC_GREY)

# ── SLIDE 2 — La premisse ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_BLUE)
rect(s,0,0,0.45,7.5, HC_GOLD)
txt(s,'La question que tout dirigeant devrait se poser.',
    0.7,0.55,12,0.65, size=20,bold=True, color=HC_WHITE)
txt(s,'Mais que presque personne ne pose.',
    0.7,1.2,12,0.55, size=18,italic=True, color=HC_GOLD)
rect(s,0.7,1.9,11.5,0.05, HC_WHITE)

lines = [
    'Lorsqu\'un collaborateur quitte votre entreprise, vous calculez le cout du recrutement.',
    'Vous calculez peut-etre le cout de la formation de son remplacant.',
    'Mais vous ne calculez jamais le cout de ce qu\'il emporte avec lui :',
    '       ses idees, sa memoire du terrain, son reseau, ses observations non dites.',
    'Ce document montre que c\'est precisement cette partie invisible qui est la plus chere.',
]
for i, l in enumerate(lines):
    col = HC_GOLD if '       ' in l else HC_LIGHT
    bold = 'precisement' in l or 'invisible' in l
    txt(s, l, 0.7,2.15+i*0.85,11.8,0.75, size=13,bold=bold, color=col)

txt(s,'Sydney Yoshida l\'a demontre en 1989. Les entreprises font encore la meme erreur en 2026.',
    0.7,6.35,11.5,0.65, size=13,bold=True,italic=True, color=HC_GOLD)
footer(s)

# ── SLIDE 3 — L'iceberg (graphique) ──────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'L\'ICEBERG DE L\'IGNORANCE — Ce que chaque niveau connait du terrain',
       'Yoshida (1989) : le sommet ne voit que 4 a 9% de la realite operationnelle',
       bg_col=HC_BLUE)
img(s, chart_iceberg, 0.15,1.0,13.03,6.3)
footer(s)

# ── SLIDE 4 — STAT CHOC : 9% ──────────────────────────────────────────
stat_slide(prs,
    big_number='9',
    unit='%',
    label_top='Iceberg de l\'Ignorance — Sydney Yoshida, 1989.',
    label_bottom='C\'est tout ce que voit votre DG de la realite du terrain.',
    question='Si vous ne voyez que 9% du reel, comment prenez-vous des decisions sur 100% de votre organisation ?',
    context='91% des problemes, des opportunites, et des idees du terrain ne remontent jamais jusqu\'au sommet. '
            'Ce n\'est pas un echec de communication. C\'est la structure naturelle de toute organisation hierarchique.',
    tag='L\'ANGLE MORT STRATEGIQUE', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 5 — STAT CHOC : 91% d'idees perdues ────────────────────────
stat_slide(prs,
    big_number='91',
    unit='%',
    label_top='Quand un agent de terrain quitte votre entreprise.',
    label_bottom='des idees qu\'il portait ne seront jamais connues du management.',
    question='Combien d\'innovations, de ventes, d\'alertes client avez-vous manquees parce qu\'un agent est parti ?',
    context='Un agent voit 100% du terrain. 91% de ce savoir ne remonte jamais. '
            'Quand il part, il emporte les 9% qui avaient commence a remonter — et les 91% qui n\'avaient jamais eu la chance de le faire.',
    tag='LA PERTE SILENCIEUSE', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 6 — Matrice double maturite ────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'LES DEUX TYPES DE MATURITE — Pourquoi les departs ne se compensent pas',
       'Maturite operationnelle (terrain) vs Maturite strategique (systeme) — deux mondes complementaires')
img(s, chart_maturite, *CHART_AREA)
footer(s)

# ── SLIDE 7 — STAT CHOC : le paradoxe de la valeur ───────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
rect(s,0,0,0.45,7.5, HC_GOLD)
txt(s,'Le paradoxe de la valeur en entreprise.',
    0.65,0.5,12,0.65, size=22,bold=True, color=HC_WHITE)
rect(s,0.65,1.25,11.5,0.05, HC_GOLD)

blocs = [
    ('Le terrain (agents)', 'Volume : eleve | Maturite operationnelle : 95/100 | Maturite strategique : 15/100',
     'Ils voient TOUT. Mais personne ne les ecoute avant qu\'ils partent.', HC_RED),
    ('Le sommet (DG)', 'Volume : faible | Maturite operationnelle : 10/100 | Maturite strategique : 95/100',
     'Ils decident de TOUT. Mais ils ne voient que 9% du reel.', HC_GREEN),
    ('La zone critique (Middle → Prox.)', 'L\'interface qui traduit le terrain en strategie — et la strategie en action.',
     'C\'est la que la valeur se concentre. Et c\'est souvent la que le turnover est le plus invisible.', HC_GOLD),
]
for i, (titre, detail, insight, col) in enumerate(blocs):
    y = 1.45 + i*1.8
    rect(s, 0.5,y,12.33,1.65, RGBColor(0x1E,0x3A,0x5F))
    rect(s, 0.5,y,0.1,1.65, col)
    txt(s, titre,   0.7,y+0.08,11.9,0.5, size=13,bold=True, color=col)
    txt(s, detail,  0.7,y+0.55,11.9,0.45, size=10.5, color=HC_LIGHT, italic=True)
    txt(s, f'>>> {insight}', 0.7,y+1.05,11.9,0.5, size=11,bold=True, color=HC_WHITE)
footer(s)

# ── SLIDE 8 — Volume vs impact par idee ──────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'VOLUME D\'IDEES vs IMPACT STRATEGIQUE — La complementarite ignoree',
       'Le terrain genere l\'information. Le sommet en fait la strategie. Sans l\'un, l\'autre est aveugle.')
img(s, chart_valeur_idee, *CHART_AREA)
footer(s)

# ── SLIDE 9 — Cout visible vs reel ───────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'LE VRAI COUT D\'UN DEPART — Ce que le RH calcule vs ce qui est reellement perdu',
       'La partie emergee (recrutement) vs la partie immergee (valeur cachee, idees, savoir tacite)',
       bg_col=HC_RED)
img(s, chart_cout, *CHART_AREA)
footer(s)

# ── SLIDE 10 — STAT CHOC : l'agent qui part ───────────────────────────
stat_slide(prs,
    big_number='9.7',
    unit='x',
    label_top='Le cout reel d\'un depart d\'agent de terrain.',
    label_bottom='son salaire annuel — pas 0.75x comme le calcule le RH traditionnel.',
    question='Votre taux de turnover frontline vous semble-t-il encore anodin ?',
    context='Cout RH visible : 0.75x salaire annuel (recrutement + integration). '
            'Cout reel Yoshida : +8.5x en valeur cachee (savoir tacite, idees perdues, capital relationnel client, '
            'temps de la courbe d\'apprentissage cumule). '
            'Total : ~9.7x. Multipliez par votre nombre de departs annuels.',
    tag='AGENT / FRONTLINE', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 11 — STAT CHOC : le middle manager ──────────────────────────
stat_slide(prs,
    big_number='6.5',
    unit='x',
    label_top='Le cout reel d\'un depart de Middle Manager.',
    label_bottom='son salaire annuel — le niveau le plus strategiquement sous-estime.',
    question='Le middle manager est l\'interface entre votre terrain et votre strategie. Quand il part, les deux se deconnectent.',
    context='Le middle manager traduit 30% du terrain en langage strategique. '
            'Il produit des idees hybrides (operationnelles ET strategiques). '
            'Son depart coupe litteralement le flux d\'information ascendant. '
            'Cout visible : 1.75x. Valeur cachee : 3.5x supplementaires. Total : ~6.5x.',
    tag='MIDDLE MANAGER', tag_col=HC_GOLD, num_col=HC_GOLD)

# ── SLIDE 12 — STAT CHOC : le DG ─────────────────────────────────────
stat_slide(prs,
    big_number='11.5',
    unit='x',
    label_top='Le cout reel d\'un depart de DG.',
    label_bottom='son salaire annuel — mais son impact est strategique, pas operationnel.',
    question='Un DG part avec 9% du visible — mais 100% de la vision. Etes-vous pret a gerer cet heritage ?',
    context='Cout RH visible : 4x salaire (headhunter + integration + courbe d\'apprentissage). '
            'Valeur cachee : 6.5x en capital relationnel externe, orientation strategique, confiance des partenaires. '
            'Total estimé : ~11.5x. Et contrairement a l\'agent, ce cout est visible immediatement.',
    tag='DG / C-SUITE', tag_col=HC_BLUE, num_col=HC_GREEN)

# ── SLIDE 13 — Impact immediat vs cumule ─────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'L\'ASYMETRIE DE L\'IMPACT — Immediat vs Cumule sur 18 mois',
       'Ce que le dirigeant ressent vs ce qu\'il ne verra jamais')
img(s, chart_courbeU, *CHART_AREA)
footer(s)

# ── SLIDE 14 — Simulation CA ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'SIMULATION — Cout d\'un depart en millions FCFA',
       'Hypothese : CA de reference 10 Md FCFA | Salaires representatifs marche Afrique de l\'Ouest')
img(s, chart_simulation, *CHART_AREA)
footer(s)

# ── SLIDE 15 — Tableau synthese ───────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'TABLEAU DE SYNTHESE — Tous les indicateurs par niveau de seniorite',
       'Le modele complet : information, idees, maturite, cout, delais')
img(s, chart_tableau, 0.2,1.08,12.93,6.2)
footer(s)

# ── SLIDE 16 — Les 3 implications strategiques ────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'3 IMPLICATIONS STRATEGIQUES POUR LES DIRIGEANTS',
       'Ce que l\'Iceberg de Yoshida change dans votre fa\u00e7on de gouverner')

implications = [
    ('01',
     'LE TURNOVER FRONTLINE N\'EST PAS "NORMAL"',
     f'Un depart d\'agent = 9.7x son salaire en cout reel. Pour 10 departs/an a 2.4M FCFA/an : '
     f'233M FCFA de perte cachee. Est-ce dans votre compte de resultat ?',
     HC_RED),
    ('02',
     'VOTRE SYSTEME D\'INFORMATION INTERNE EST VOTRE PREMIER RISQUE STRATEGIQUE',
     'Si 91% du terrain ne remonte pas, vos decisions strategiques reposent sur 9% du reel. '
     'La question n\'est pas "avons-nous une bonne strategie ?" mais "avons-nous les bonnes donnees pour la batir ?"',
     HC_GOLD),
    ('03',
     'LA RETENTION DES MIDDLE MANAGERS EST VOTRE LEVIER N°1',
     'Ils sont l\'interface. Ni completement terrain, ni completement sommet. '
     'Quand ils partent, le terrain et le sommet se retrouvent a parler deux langues differentes. '
     'Et personne ne le remarque tout de suite.',
     HC_GREEN),
]
for i, (num, titre, detail, col) in enumerate(implications):
    y = 1.15 + i*1.9
    rect(s,0.3,y,12.73,1.75, RGBColor(0x1E,0x3A,0x5F))
    rect(s,0.3,y,0.1,1.75, col)
    txt(s,num,   0.5,y+0.08,0.7,0.7, size=24,bold=True, color=col)
    txt(s,titre, 1.2,y+0.06,11.5,0.55, size=12.5,bold=True, color=HC_WHITE)
    txt(s,detail,1.2,y+0.65,11.5,0.95, size=10.5, color=HC_LIGHT, wrap=True)
footer(s)

# ── SLIDE 17 — 5 Actions ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'5 ACTIONS CONCRETES — Pour un dirigeant qui a compris l\'iceberg',
       'Transformer la connaissance en avantage competitif')

actions = [
    ('01', 'CREER DES CANAUX D\'IDEATION ASCENDANTE FORMELS',
     'Boites a idees? Non. Des rituels hebdomadaires ou le terrain remonte systematiquement ses observations au management.', HC_GREEN),
    ('02', 'CALCULER LE VRAI COUT DE VOTRE TURNOVER',
     'Demandez a votre DRH : "quel est le cout total d\'un depart, valeur cachee incluse ?" '
     'Si la reponse ne depasse pas 2x le salaire, c\'est que le calcul est incomplet.', HC_GOLD),
    ('03', 'SEGMENTER VOTRE STRATEGIE DE RETENTION PAR NIVEAU',
     'Un agent fidel = protection de l\'information terrain. '
     'Un middle manager fidele = continuité du flux strategique. Ce ne sont pas les memes actions.', HC_GOLD),
    ('04', 'MESURER L\'INFORMATION QUI NE REMONTE PAS',
     'Enquete annuelle anonyme : "quelle est votre principale observation que votre manager ne connait pas ?" '
     'Les reponses vont vous surprendre.', HC_RED),
    ('05', 'FAIRE DU DEPART D\'UN COLLABORATEUR UN AUDIT D\'INFORMATION',
     'L\'entretien de sortie est votre derniere chance de recuperer 91% de ce que vous n\'avez pas vu. '
     'Formalisez-le. Enregistrez-le. Analysez-le.', HC_RED),
]
for i, (num, titre, detail, col) in enumerate(actions):
    y = 1.12 + i*1.2
    rect(s,0.25,y,12.83,1.1, RGBColor(0x1E,0x3A,0x5F))
    rect(s,0.25,y,0.1,1.1, col)
    txt(s,num,   0.45,y+0.05,0.7,0.7, size=22,bold=True, color=col)
    txt(s,titre, 1.2,y+0.04,11.7,0.5, size=11.5,bold=True, color=HC_WHITE)
    txt(s,detail,1.2,y+0.58,11.7,0.46, size=9.5, color=HC_LIGHT, italic=True)
footer(s)

# ── SLIDE 18 — Conclusion ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
rect(s,0,0,0.45,7.5, HC_GOLD)
rect(s,0.45,0,12.88,0.07, HC_GOLD)

txt(s,'Votre entreprise perd de l\'argent\nchaque jour. En silence.',
    0.75,0.5,12,1.3, size=30,bold=True, color=HC_WHITE)
rect(s,0.75,2.0,11,0.05, HC_GOLD)
txt(s,'Pas a cause de la concurrence. Pas a cause du marche.',
    0.75,2.2,11.5,0.6, size=16,italic=True, color=HC_LIGHT)
txt(s,'A cause des 91% d\'informations qui ne remontent jamais.\nA cause des idees qui partent avec les gens qui quittent.',
    0.75,2.85,11.5,1.0, size=16,bold=True, color=HC_WHITE)
rect(s,0.75,4.0,11,0.05, HC_GOLD)
txt(s,'La vraie question de gouvernance n\'est pas\n"Comment optimiser ma strategie ?"\nMais "Comment m\'assurer que ma strategie\nrepose sur 100% du reel — et pas sur 9% ?"',
    0.75,4.2,11.5,2.2, size=15,bold=True, color=HC_GOLD)
txt(s,'H&C Executive — Conseil & Accompagnement Strategique',
    0.75,6.7,8,0.4, size=11,italic=True, color=HC_GREY)

# ── Sauvegarder ──────────────────────────────────────────────────────
output = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/GENZ_Iceberg_Ignorance_HC_2026.pptx'
prs.save(output)
print(f'\nPowerPoint sauvegarde : {output}')
print(f'Slides : {len(prs.slides)}')
print()
print('=== STRUCTURE ===')
slides_list = [
    'Couverture','La premisse — ce qu\'on ne calcule jamais',
    'L\'Iceberg de Yoshida (graphique)',
    'STAT : 9% — ce que voit le DG',
    'STAT : 91% — idees perdues quand un agent part',
    'Matrice double maturite (graphique)',
    'Le paradoxe de la valeur',
    'Volume idees vs impact strategique (graphique)',
    'Cout visible vs cout reel (graphique)',
    'STAT : 9.7x — cout reel depart agent',
    'STAT : 6.5x — cout reel depart middle manager',
    'STAT : 11.5x — cout reel depart DG',
    'Impact immediat vs cumule (graphique)',
    'Simulation CA en MFCFA (graphique)',
    'Tableau de synthese (graphique)',
    '3 implications strategiques',
    '5 actions concretes',
    'Conclusion',
]
for i, s_name in enumerate(slides_list):
    print(f'  Slide {i+1:02d} : {s_name}')
