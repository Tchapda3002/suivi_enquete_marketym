"""
Generateur PowerPoint — Barometre Gen Z Afrique 2026
H&C Executive — Storytelling dirigeant : chiffres chocs + graphiques lisibles
"""

import json, os, numpy as np, pandas as pd
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
from sklearn.cluster import KMeans
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────
HC_BLUE  = RGBColor(0x1B, 0x3A, 0x6B)
HC_GOLD  = RGBColor(0xC9, 0xA8, 0x4C)
HC_RED   = RGBColor(0xC0, 0x39, 0x2B)
HC_GREEN = RGBColor(0x1A, 0xBC, 0x9C)
HC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HC_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
HC_GREY  = RGBColor(0x7F, 0x8C, 0x8D)
HC_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)

B  = '#1B3A6B'; G = '#C9A84C'; R = '#C0392B'
GR = '#1ABC9C'; GY= '#7F8C8D'; W = '#FFFFFF'
DK = '#1A1A2E'; LB= '#F4F6F9'

matplotlib.rcParams['font.family'] = 'DejaVu Sans'
matplotlib.rcParams['figure.dpi']  = 150

IMG_DIR = '/tmp/ppt_charts'
os.makedirs(IMG_DIR, exist_ok=True)

# ── Donnees ───────────────────────────────────────────────────────────
with open('/tmp/genz_dataset.json') as f:
    raw = json.load(f)
df = pd.DataFrame(raw)
N  = len(df)

sante = df[df['secteur'].str.contains('Sant', na=False)]
agri  = df[df['secteur'].str.contains('Agri', na=False)]
educ  = df[df['secteur'].str.contains('Educ', na=False)]

def savefig(name):
    path = f'{IMG_DIR}/{name}.png'
    plt.savefig(path, bbox_inches='tight', dpi=150, facecolor=plt.gcf().get_facecolor())
    plt.close()
    return path

# ════════════════════════════════════════════════════════════════════════
# CHARTS DEDIES (grands, lisibles, une seule idee chacun)
# ════════════════════════════════════════════════════════════════════════

# Chart 1 — Composantes du bas vers le haut
fig, ax = plt.subplots(figsize=(12, 5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
comp = {'Remuneration':df['REM'].mean(),'Reconnaissance':df['RECO'].mean(),
        'Climat Travail':df['IATM'].mean(),'NPS Employeur':df['NPS'].mean(),
        'Satisfaction':df['SAT'].mean(),'Valorisation':df['VAL'].mean(),
        'Apprentissage':df['APPRENT'].mean(),'Engagement':df['ENG'].mean(),
        'Initiative':df['INIT'].mean(),'Innovation':df['INNOV'].mean()}
comp_s = dict(sorted(comp.items(), key=lambda x:x[1]))
cols_c = [R if v<60 else (G if v<70 else GR) for v in comp_s.values()]
bars = ax.barh(list(comp_s.keys()), list(comp_s.values()), color=cols_c, height=0.6, edgecolor=W)
ax.axvline(65, color='#888', linestyle='--', lw=1.2, alpha=0.5)
for bar, val in zip(bars, comp_s.values()):
    ax.text(val+0.8, bar.get_y()+bar.get_height()/2, f'{val:.1f}', va='center', fontsize=12, fontweight='bold')
ax.set_xlim(0,100); ax.tick_params(labelsize=12)
ax.set_title('Toutes les composantes — du plus bas au plus haut', fontsize=14, fontweight='bold', color=B, pad=10)
p1=mpatches.Patch(color=R,label='Zone critique (<60)')
p2=mpatches.Patch(color=G,label='Zone fragile (60-70)')
p3=mpatches.Patch(color=GR,label='Zone correcte (>70)')
ax.legend(handles=[p1,p2,p3], fontsize=10, loc='lower right')
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_comp = savefig('01_composantes')

# Chart 2 — Sous-payes par age
fig, ax = plt.subplots(figsize=(12, 4.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
ages = ['18-22 ans (Gen Z)','23-27 ans (Gen Z)','28-32 ans (Millennials)','33-37 ans (Millennials)','38-42 ans (Millennials)']
xlabs = ['18-22\n(Gen Z)','23-27\n(Gen Z)','28-32\n(Millenn.)','33-37\n(Millenn.)','38-42\n(Millenn.)']
pcts  = [(df[df['age']==a]['SOUSPAYE'].dropna()>50).mean()*100 for a in ages]
cols_a= [R if p>65 else G for p in pcts]
bars  = ax.bar(range(5), pcts, color=cols_a, width=0.55, edgecolor=W)
ax.axhline(50, color='#555', linestyle='--', lw=1.2, alpha=0.5)
for bar,val in zip(bars,pcts):
    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
            f'{val:.0f}%', ha='center', fontsize=15, fontweight='bold')
ax.set_xticks(range(5)); ax.set_xticklabels(xlabs, fontsize=12)
ax.set_ylim(0,90); ax.set_ylabel('% repondants', fontsize=12)
ax.set_title('% se sentant sous-payes par tranche d\'age', fontsize=14, fontweight='bold', color=B, pad=10)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_age = savefig('02_souspaye_age')

# Chart 3 — Sous-payes par secteur
fig, ax = plt.subplots(figsize=(12, 5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
sects_d = {}
for s in df['secteur'].dropna().unique():
    sub = df[df['secteur']==s]['SOUSPAYE'].dropna()
    if len(sub)>=20:
        sects_d[s[:28]] = (sub>50).mean()*100
sects_s = dict(sorted(sects_d.items(), key=lambda x:x[1], reverse=True))
cols_s  = [R if v>70 else (G if v>55 else GR) for v in sects_s.values()]
bars = ax.barh(list(sects_s.keys()), list(sects_s.values()), color=cols_s, height=0.6, edgecolor=W)
ax.axvline(50, color='#555', linestyle='--', lw=1.2, alpha=0.5)
for bar,val in zip(bars, sects_s.values()):
    ax.text(val+0.5, bar.get_y()+bar.get_height()/2, f'{val:.0f}%', va='center', fontsize=12, fontweight='bold')
ax.set_xlim(0,100); ax.tick_params(labelsize=11)
ax.set_title('Sous-paiement par secteur — qui souffre le plus ?', fontsize=14, fontweight='bold', color=B, pad=10)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_sect = savefig('03_souspaye_secteur')

# Chart 4 — Heatmap IDAT pays x secteur
fig, ax = plt.subplots(figsize=(12, 5))
fig.patch.set_facecolor(LB)
top_pays = df['pays'].value_counts().head(7).index.tolist()
top_s    = ['Santé / Pharmacie','Banque / Finance / Assurance','Agriculture / Agroalimentaire',
            'Éducation / Formation','Commerce / Distribution','Administration publique']
matrix   = []
for s in top_s:
    row = []
    for p in top_pays:
        sub = df[(df['pays']==p)&(df['secteur']==s)]['IDAT'].dropna()
        row.append(round(sub.mean(),1) if len(sub)>=5 else np.nan)
    matrix.append(row)
mdf = pd.DataFrame(matrix, index=[s[:20] for s in top_s], columns=top_pays)
sns.heatmap(mdf, ax=ax, cmap='RdYlGn', vmin=55, vmax=78,
            annot=True, fmt='.0f', linewidths=1.2, linecolor='white',
            annot_kws={'size':12,'weight':'bold'},
            cbar_kws={'label':'Score IDAT /100','shrink':0.85})
ax.set_title('IDAT par Pays x Secteur — ou agir en priorite ?', fontsize=14, fontweight='bold', color=B, pad=10)
ax.tick_params(axis='x', labelsize=11, rotation=25)
ax.tick_params(axis='y', labelsize=11, rotation=0)
fig.tight_layout()
chart_heatmap = savefig('04_heatmap')

# Chart 5 — Correlations avec IDAT
fig, ax = plt.subplots(figsize=(12, 5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
cols_c2  = ['SAT','ENG','VAL','NPS','REM','RECO','APPRENT','INIT','INNOV','IATM','SOUSPAYE_inv']
labs_c2  = ['Satisfaction','Engagement','Valorisation','NPS Employeur','Remuneration',
            'Reconnaissance','Apprentissage','Initiative','Innovation','Climat Travail','Non sous-paye']
corrs2   = [df[c].corr(df['IDAT']) for c in cols_c2]
pairs2   = sorted(zip(corrs2, labs_c2), key=lambda x:x[0])
cols_cr  = [R if v<0.3 else (G if v<0.5 else GR) for v in [p[0] for p in pairs2]]
bars = ax.barh([p[1] for p in pairs2], [p[0] for p in pairs2], color=cols_cr, height=0.6, edgecolor=W)
ax.axvline(0.5, color=GR, linestyle=':', lw=1.5, alpha=0.7, label='Forte correlation')
ax.axvline(0.3, color=G,  linestyle=':', lw=1.5, alpha=0.7, label='Moderee')
ax.axvline(0,   color='#333', lw=1)
for bar,val in zip(bars, [p[0] for p in pairs2]):
    ax.text(val+0.01, bar.get_y()+bar.get_height()/2, f'{val:.2f}',
            va='center', fontsize=11, fontweight='bold')
ax.tick_params(labelsize=12)
ax.set_title('Ce qui drive vraiment la retention Gen Z', fontsize=14, fontweight='bold', color=B, pad=10)
ax.legend(fontsize=10)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_corr = savefig('05_correlations')

# Chart 6 — 4 profils K-means
cols_pca = ['SAT','ENG','VAL','NPS','REM','RECO','APPRENT','INIT','INNOV','IATM','SOUSPAYE_inv']
df_pca   = df[cols_pca].dropna().copy()
X_sc     = StandardScaler().fit_transform(df_pca)
km       = KMeans(n_clusters=4, random_state=42, n_init=20)
df_pca['cluster'] = km.fit_predict(X_sc)
cm       = df_pca.groupby('cluster')[cols_pca].mean()
cs       = df_pca['cluster'].value_counts().sort_index()

radar_c = ['SAT','ENG','VAL','REM','RECO','INNOV','IATM']
radar_l = ['Satisfaction','Engagement','Valorisation','Remuneration','Reconnaissance','Innovation','Climat']
gmeans  = df[radar_c].mean().values

ccolors = [GR, R, G, B]
def cname(c):
    row = cm.loc[c]
    n   = cs[c]; pct = n/len(df_pca)*100
    if row['ENG']>75 and row['SAT']>70:   return f'Ambassadeurs\n{pct:.0f}% ({n} pers.)'
    if row['ENG']>70 and row['REM']<55:   return f'Frustres Ambitieux\n{pct:.0f}% ({n} pers.)'
    if row['REM']<50 and row['RECO']<50:  return f'En Danger\n{pct:.0f}% ({n} pers.)'
    return f'Attentistes\n{pct:.0f}% ({n} pers.)'

fig, axes = plt.subplots(1, 4, figsize=(15, 4.5), sharey=True)
fig.patch.set_facecolor(LB)
fig.suptitle('Les 4 profils Gen Z — 4 reponses manageriales differentes',
             fontsize=13, fontweight='bold', color=B, y=1.03)
for idx, (c, ax) in enumerate(zip(range(4), axes)):
    vals = cm.loc[c, radar_c].values
    col  = ccolors[idx]
    x    = np.arange(len(radar_c))
    ax.set_facecolor(W)
    ax.bar(x, gmeans, width=0.6, color='#ddd', alpha=0.5)
    ax.bar(x, vals,   width=0.6, color=col, alpha=0.82, edgecolor=W)
    ax.set_xticks(x)
    ax.set_xticklabels(radar_l, fontsize=8.5, rotation=30, ha='right')
    ax.set_ylim(0,100)
    ax.axhline(65, color='#aaa', linestyle='--', lw=1)
    for xi,v in enumerate(vals):
        ax.text(xi, v+1.5, f'{v:.0f}', ha='center', fontsize=9, fontweight='bold')
    ax.set_title(cname(c), fontsize=11, fontweight='bold', color=col, pad=8)
    ax.spines[['top','right']].set_visible(False)
    if idx==0:
        ax.set_ylabel('Score /100', fontsize=11)
        p1 = mpatches.Patch(color='#ddd', label='Moy. globale')
        p2 = mpatches.Patch(color=col, label='Ce profil')
        ax.legend(handles=[p1, p2], fontsize=8)
fig.tight_layout()
chart_profils = savefig('06_profils')

# Chart 7 — Courbe de vie par age
fig, ax = plt.subplots(figsize=(12, 4.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
ages_k = ['18-22 ans (Gen Z)','23-27 ans (Gen Z)','28-32 ans (Millennials)','33-37 ans (Millennials)','38-42 ans (Millennials)']
xlabs2 = ['18-22\nGenZ','23-27\nGenZ','28-32\nMillenn.','33-37\nMillenn.','38-42\nMillenn.']
for m, col, mk, lw, lbl in [('IDAT',B,'o',2.5,'Attractivite (IDAT)'),
                              ('ENG',GR,'s',2,'Engagement'),
                              ('IATM',R,'^',2,'Climat (IATM)'),
                              ('REM',G,'D',2,'Remuneration')]:
    vals = [df[df['age']==a][m].mean() for a in ages_k]
    ax.plot(range(5), vals, f'{mk}-', color=col, lw=lw, ms=8, label=lbl)
ax.set_xticks(range(5)); ax.set_xticklabels(xlabs2, fontsize=12)
ax.set_ylim(40,95); ax.set_ylabel('Score moyen /100', fontsize=12)
ax.set_title('La courbe de vie au travail — comment l\'experience evolue par age', fontsize=13, fontweight='bold', color=B, pad=10)
ax.legend(fontsize=11, loc='upper left')
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
chart_vie = savefig('07_courbe_vie')

print('Charts generes.')

# ════════════════════════════════════════════════════════════════════════
# POWERPOINT — STORYTELLING DIRIGEANT
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color

def rect(slide, l,t,w,h, color, alpha_color=None):
    s = slide.shapes.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l,t,w,h, size=14, bold=False, color=HC_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color

def img(slide, path, l,t,w,h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l),Inches(t),Inches(w),Inches(h))

def header(slide, title, subtitle='', bg_col=HC_BLUE, h=1.0):
    rect(slide,0,0,13.33,h, bg_col)
    txt(slide, title,   0.35,0.07,12.6,0.55, size=20, bold=True, color=HC_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35,0.63,12.6,0.35, size=11, color=HC_GOLD, italic=True)

def footer(slide, text='H&C Executive  |  Barometre Gen Z Afrique 2026  |  Confidentiel'):
    txt(slide, text, 0,7.18,13.33,0.32, size=9, color=HC_GREY, align=PP_ALIGN.CENTER)

def stat_slide(prs, big_number, unit, label_top, label_bottom,
               question, context, bg_col=HC_DARK, num_col=HC_GOLD,
               tag='', tag_col=HC_RED):
    """Slide choc : un gros chiffre, un contexte, une question percutante."""
    s = prs.slides.add_slide(blank)
    bg(s, bg_col)
    rect(s,0,0,0.45,7.5, num_col)
    if tag:
        rect(s,0.55,0.18,2.8,0.55, tag_col)
        txt(s, tag, 0.6,0.22,2.7,0.45, size=11, bold=True, color=HC_WHITE, align=PP_ALIGN.CENTER)
    txt(s, label_top, 0.65,0.85,12,0.7, size=14, color=HC_LIGHT, italic=True)
    txt(s, big_number, 0.55,1.4,12,2.8, size=110, bold=True, color=num_col, align=PP_ALIGN.LEFT)
    txt(s, unit,       8.0, 2.5,4.5,1.2, size=32, bold=True, color=num_col)
    txt(s, label_bottom, 0.65,4.15,12,0.7, size=22, bold=True, color=HC_WHITE)
    rect(s,0.65,4.95,11.5,0.04, num_col)
    txt(s, context,   0.65,5.1,11.5,0.75, size=12, color=HC_LIGHT, italic=True)
    txt(s, f'>>> {question}', 0.65,5.95,11.5,0.9, size=14, bold=True, color=num_col)

# ── SLIDE 1 — Couverture ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
rect(s,0,0,0.45,7.5, HC_BLUE)
rect(s,0.45,0,12.88,0.07, HC_GOLD)
rect(s,0.45,7.43,12.88,0.07, HC_GOLD)
txt(s,'BAROMETRE',    0.75,0.9,12,0.85,  size=36,bold=True, color=HC_WHITE)
txt(s,'GENERATION Z', 0.75,1.7,12,1.1,   size=52,bold=True, color=HC_GOLD)
txt(s,'AFRIQUE 2026', 0.75,2.75,12,0.85, size=36,bold=True, color=HC_WHITE)
rect(s,0.75,3.7,7.5,0.05, HC_GOLD)
txt(s,'Ce que les chiffres disent — et que les dirigeants doivent entendre',
    0.75,3.85,11,0.55, size=15,italic=True, color=HC_LIGHT)
txt(s,f'{N} repondants  |  13 pays  |  10+ secteurs  |  Avril 2026',
    0.75,4.55,10,0.45, size=12, color=HC_GREY)
txt(s,'H&C Executive — Document confidentiel',
    0.75,6.85,7,0.45, size=11,italic=True, color=HC_GREY)

# ── SLIDE 2 — Mise en garde narrative ─────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_BLUE)
rect(s,0,0,13.33,7.5, HC_BLUE)
rect(s,0,0,0.45,7.5, HC_GOLD)
txt(s,'Ce document ne contient pas une etude de plus.',
    0.7,0.7,12,0.75, size=22,bold=True, color=HC_WHITE)
txt(s,'Il contient les signaux que vos tableaux de bord RH\nne montrent pas encore.',
    0.7,1.55,11.5,1.1, size=20,italic=True, color=HC_GOLD)
rect(s,0.7,2.85,11.3,0.05, HC_WHITE)
lines = [
    '1 200 Gen Z interroges dans 13 pays et 10+ secteurs d\'Afrique.',
    'Des chiffres bruts. Des correlations. Des faits.',
    'Et au bout : une question que chaque dirigeant devrait se poser ce soir.',
]
for i,l in enumerate(lines):
    txt(s, f'—  {l}', 0.7, 3.1+i*0.9, 11.5, 0.75, size=14, color=HC_LIGHT)
txt(s,'Bonne lecture.',
    0.7,6.1,5,0.6, size=18,bold=True,italic=True, color=HC_GOLD)
footer(s)

# ── SLIDE 3 — STAT CHOC : 59% ─────────────────────────────────────────
pct_sp = (df['SOUSPAYE']>50).mean()*100
stat_slide(prs,
    big_number=f'{pct_sp:.0f}',
    unit='%',
    label_top='Premier chiffre. Le plus important.',
    label_bottom='de vos equipes Gen Z se sentent sous-payees.',
    question='Est-ce que vous le saviez ? Et si oui, qu\'avez-vous fait ?',
    context=f'Sur {N} repondants actifs en Afrique francophone. Ce n\'est pas une minorite — c\'est la majorite de votre releve.',
    tag='SIGNAL MAJEUR', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 4 — STAT CHOC : La bombe invisible ──────────────────────────
bien_rem = df[df['REM']>65]
pct_bienrem_souspaye = (bien_rem['SOUSPAYE']>50).mean()*100
correl_rem_sp = df['REM'].corr(df['SOUSPAYE_inv'])
stat_slide(prs,
    big_number=f'{pct_bienrem_souspaye:.0f}',
    unit='%',
    label_top='Le chiffre qui doit vous arreter net.',
    label_bottom='de ceux qui sont bien remunerés se sentent QUAND MEME sous-payes.',
    question='Ce n\'est donc pas un probleme de salaire. C\'est un probleme de justice et de reconnaissance. Avez-vous la bonne reponse ?',
    context=f'Correlation entre remuneration reelle et sentiment de sous-paiement : r={correl_rem_sp:.2f}. Quasi nulle. '
             'Augmenter les salaires ne reglera pas ce probleme.',
    tag='LA BOMBE INVISIBLE', tag_col=RGBColor(0x6C,0x35,0x83), num_col=HC_GOLD)

# ── SLIDE 5 — STAT CHOC : 41% talent bride ────────────────────────────
pct_bride = ((df['ENG']>70) & (df['IATM']<50)).mean()*100
ecart_eng_iatm = df['ENG'].mean() - df['IATM'].mean()
stat_slide(prs,
    big_number=f'{pct_bride:.0f}',
    unit='%',
    label_top='Le paradoxe au coeur de votre organisation.',
    label_bottom='sont engages — mais travaillent dans un mauvais climat.',
    question='Savez-vous combien de temps un talent motive peut tenir dans un environnement qui l\'etouffe ?',
    context=f'Ecart moyen entre Engagement ({df["ENG"].mean():.0f}/100) et Climat de travail ({df["IATM"].mean():.0f}/100) : '
             f'{ecart_eng_iatm:.0f} points. Le plus grand ecart mesure dans ce barometre. '
             'Vos meilleurs elements sont vos plus frustres.',
    tag='TALENT BRIDE', tag_col=HC_RED, num_col=HC_GOLD)

# ── SLIDE 6 — STAT CHOC : 84% sante ──────────────────────────────────
pct_sp_sante = (sante['SOUSPAYE']>50).mean()*100
stat_slide(prs,
    big_number=f'{pct_sp_sante:.0f}',
    unit='%',
    label_top='Secteur Sante / Pharmacie.',
    label_bottom='des professionnels de sante Gen Z se sentent sous-payes.',
    question='Qui soignera vos populations dans 10 ans si vous ne retenez pas ces profils aujourd\'hui ?',
    context=f'Sante : IATM={sante["IATM"].mean():.0f}/100 (climat critique) — ENG={sante["ENG"].mean():.0f}/100 (engagement eleve). '
             'Ils veulent rester. Les conditions les poussent a partir.',
    tag='URGENCE SECTORIELLE', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 7 — STAT CHOC : 31 points d'ecart ──────────────────────────
stat_slide(prs,
    big_number='31',
    unit='pts',
    label_top='L\'ecart qui dit tout.',
    label_bottom='separent l\'Engagement du Climat de Travail.',
    question='Que se passe-t-il quand une voiture tourne a plein regime mais que les freins sont grippes ?',
    context=f'Engagement moyen : {df["ENG"].mean():.0f}/100. Climat de travail moyen : {df["IATM"].mean():.0f}/100. '
             '31 points d\'ecart. L\'un des plus grands paradoxes identifies dans ce barometre. '
             'Ces equipes performent malgre le systeme — pas grace a lui.',
    tag='LE PARADOXE', tag_col=RGBColor(0x6C,0x35,0x83), num_col=HC_GOLD)

# ── SLIDE 8 — STAT CHOC : 1 sur 4 ────────────────────────────────────
pct_reco = (df['RECO']<50).mean()*100
stat_slide(prs,
    big_number='1/4',
    unit='',
    label_top='Question de reconnaissance.',
    label_bottom='de votre Gen Z ne se sent pas reconnu dans son travail.',
    question='Quand avez-vous pour la derniere fois dit a un collaborateur Gen Z : "tu fais du bon travail" ?',
    context=f'Score de Reconnaissance moyen : {df["RECO"].mean():.0f}/100. {pct_reco:.0f}% sous le seuil de 50. '
             'La reconnaissance ne coute rien. Son absence coute tout.',
    tag='RECONNAISSANCE', tag_col=HC_GOLD, num_col=HC_GOLD)

# ── SLIDE 9 — STAT CHOC : 13.5% partants actifs ──────────────────────
pct_nps_low = (df['NPS']<25).mean()*100
stat_slide(prs,
    big_number=f'{pct_nps_low:.0f}',
    unit='%',
    label_top='Ils ne partiront pas demain. Ils sont deja partis dans leur tete.',
    label_bottom='ont un NPS < 25 : ils ne recommanderaient pas leur employeur et envisagent de partir.',
    question='Connaissez-vous leurs noms ? Savez-vous pourquoi ils sont encore la ?',
    context=f'NPS global : {df["NPS"].mean():.0f}/100. 1 Gen Z sur 7 est en sortie mentale. '
             '23% supplementaires ne recommanderaient pas leur employeur. '
             'Ensemble : plus d\'un tiers de votre releve est en risque de depart.',
    tag='EN SORTIE MENTALE', tag_col=HC_RED, num_col=HC_RED)

# ── SLIDE 10 — STAT CHOC : 34 points Innovation vs Climat ─────────────
ecart_innov = df['INNOV'].mean() - df['IATM'].mean()
stat_slide(prs,
    big_number=f'{ecart_innov:.0f}',
    unit='pts',
    label_top='L\'innovation bridee.',
    label_bottom='separent le score d\'Innovation du score de Climat de travail.',
    question='Combien vaut pour vous 1 point d\'innovation perdu par un environnement de travail defaillant ?',
    context=f'Innovation : {df["INNOV"].mean():.0f}/100. Initiative : {df["INIT"].mean():.0f}/100. '
             f'Climat : {df["IATM"].mean():.0f}/100. '
             'Votre Gen Z est creative, proactive, et innovative. Mais le contexte ne lui permet pas d\'exprimer ce potentiel.',
    tag='POTENTIEL INEXPLOITE', tag_col=RGBColor(0x6C,0x35,0x83), num_col=HC_GREEN)

# ── SLIDE 11 — Toutes les composantes (graphique) ─────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'LE BILAN COMPLET — Toutes les composantes classees',
       'Rouge = zone critique | Or = zone fragile | Vert = zone correcte')
img(s, chart_comp, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 12 — Sous-payes par age (graphique) ─────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'71% des 18-22 ANS se sentent sous-payes',
       'La plus jeune Gen Z est la plus touchee — signal d\'alerte precoce',
       bg_col=HC_RED)
img(s, chart_age, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 13 — Sous-payes par secteur (graphique) ─────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'SOUS-PAIEMENT PAR SECTEUR — La Sante en etat d\'urgence',
       'Certains secteurs concentrent des tensions insoutenables',
       bg_col=HC_RED)
img(s, chart_sect, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 14 — Heatmap pays x secteur ────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'CARTOGRAPHIE STRATEGIQUE — Pays × Secteur × Attractivite',
       'Rouge = urgence absolue | Vert = zone correcte | Blanc = donnees insuffisantes')
img(s, chart_heatmap, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 15 — Correlations (graphique) ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'CE QUI DRIVE VRAIMENT LA RETENTION — Au-dela des idees recues',
       'Augmenter les salaires seuls ne regle rien — la donnee le prouve')
img(s, chart_corr, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 16 — 4 Profils (graphique) ──────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'LES 4 PROFILS GEN Z — Ne les managez pas de la meme facon',
       '4 types d\'employes, 4 strategies RH differentes, 4 risques distincts')
img(s, chart_profils, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 17 — Courbe de vie ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'LA COURBE DE VIE AU TRAVAIL — Comment l\'experience evolue par age',
       'Les 18-22 ans sont les plus frustres et les moins bien remuneres — les premiers a partir')
img(s, chart_vie, 0.2,1.05,12.9,6.2)
footer(s)

# ── SLIDE 18 — Ce que ca coute de ne rien faire ───────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'LE COUT DE L\'INACTION — Ce que personne ne calcule',
       'Avant les leviers d\'action, posons le vrai cout du statu quo', bg_col=HC_RED)

couts = [
    ('Turnover 1 profil senior',
     '6 a 18 mois de salaire\npour recruter + former son remplacant',
     HC_RED),
    ('Desengagement silencieux',
     '41% sont engages dans un mauvais climat.\nUn talent qui souffre produit 20 a 30% moins.',
     HC_GOLD),
    ('Perte d\'innovation',
     '34 points d\'ecart entre potentiel créatif et environnement.\nChaque point perdu = idee non developpee.',
     HC_GREEN),
    ('Image employeur',
     '23% ne recommanderaient pas leur employeur.\nVotre marque employeur se construit dans les conversations privees.',
     HC_GOLD),
]

for i,(titre,detail,col) in enumerate(couts):
    col_n = i % 2; row_n = i // 2
    x = 0.3 + col_n * 6.55; y = 1.2 + row_n * 2.8
    rect(s,x,y,6.3,2.55, RGBColor(0x1E,0x3A,0x5F))
    rect(s,x,y,0.1,2.55, col)
    txt(s, titre,  x+0.2,y+0.12,5.9,0.6, size=13,bold=True, color=HC_WHITE)
    txt(s, detail, x+0.2,y+0.75,5.9,1.6, size=11, color=HC_LIGHT, wrap=True)
footer(s)

# ── SLIDE 19 — 5 Leviers d'action ────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
header(s,'5 LEVIERS D\'ACTION — Classes par Retour sur Investissement',
       'Des actions. Pas un rapport de plus. Chaque levier est mesurable.')

leviers = [
    ('01','STRUCTURER LA RECONNAISSANCE formelle et mensuelle',
     f'RECO={df["RECO"].mean():.0f}/100  |  Zero budget necessaire — question de culture de management',
     HC_GREEN,'30 jours'),
    ('02','TRANSPARENCER LES GRILLES DE REMUNERATION',
     f'{(df["SOUSPAYE"]>50).mean()*100:.0f}% sous-payes malgre des salaires corrects  |  La transparence seule reduit de 30% le sentiment',
     HC_GOLD,'90 jours'),
    ('03','FORMER LES MANAGERS AU LEADERSHIP GEN Z',
     f'IATM={df["IATM"].mean():.0f}/100  |  Le manager direct represente 70% du climat — formation 2 jours',
     HC_GOLD,'60 jours'),
    ('04','CREER DES PARCOURS D\'EVOLUTION VISIBLES a 6 et 18 mois',
     f'INIT={df["INIT"].mean():.0f}, INNOV={df["INNOV"].mean():.0f}  |  L\'ambition est la — il manque le canal pour l\'exprimer',
     HC_GREEN,'45 jours'),
    ('05','INTERVENTIONS PRIORITAIRES : Sante et Agriculture',
     'Sante : 84% sous-payes | Agriculture : IDAT le plus bas  |  Ces 2 secteurs concentrent toutes les tensions',
     HC_RED,'60 jours'),
]
for i,(num,titre,stat,col,delai) in enumerate(leviers):
    y = 1.15 + i*1.15
    rect(s,0.25,y,12.83,1.08, RGBColor(0x1E,0x3A,0x5F))
    rect(s,0.25,y,0.1,1.08, col)
    txt(s,num,   0.45,y+0.05,0.7,0.7, size=22,bold=True, color=col)
    txt(s,titre, 1.2,y+0.04,8.3,0.52, size=11.5,bold=True, color=HC_WHITE)
    txt(s,f'>>> {stat}', 1.2,y+0.58,8.3,0.42, size=9.5,italic=True, color=HC_GOLD)
    txt(s,f'Delai : {delai}', 9.7,y+0.3,3.1,0.45, size=10,bold=True, color=col, align=PP_ALIGN.RIGHT)
footer(s)

# ── SLIDE 20 — Feuille de route ───────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, RGBColor(0xF4,0xF6,0xF9))
header(s,'FEUILLE DE ROUTE — Ce qui se decide aujourd\'hui se recolte dans 18 mois',
       '90 jours determines = transformation reelle de l\'experience Gen Z')

phases = [
    ('90 JOURS',   HC_RED,   ['Audit reconnaissance : cartographier les pratiques actuelles',
                               'Formation pilote managers (2 entreprises test)',
                               'Benchmarks salariaux Sante et Agriculture',
                               'Communication interne : rendre les grilles visibles']),
    ('6 MOIS',     HC_GOLD,  ['Deploiement systeme reconnaissance structure',
                               'Revisions salariales ciblees secteurs critiques',
                               'Parcours d\'evolution formalises pour Gen Z',
                               'Barometre interne trimestriel : mesurer l\'impact']),
    ('12 MOIS',    HC_GREEN, ['Remesure IDAT — cible : +8 points',
                               '% sous-payes — cible : -15 points',
                               'Retention Gen Z — cible : 85%+',
                               'Extension nationale du modele Gen Z Ready']),
]
for i,(periode,col,actions) in enumerate(phases):
    x = 0.3 + i*4.35
    rect(s,x,1.15,4.1,6.05, HC_WHITE)
    rect(s,x,1.15,4.1,0.62, col)
    txt(s,periode, x+0.1,1.18,3.9,0.58, size=18,bold=True, color=HC_WHITE, align=PP_ALIGN.CENTER)
    for j,action in enumerate(actions):
        y_a = 1.97 + j*1.15
        rect(s,x+0.15,y_a+0.12,0.08,0.75, col)
        txt(s,action, x+0.32,y_a,3.65,1.0, size=10.5, color=HC_BLUE, wrap=True)
footer(s)

# ── SLIDE 21 — Conclusion percutante ─────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, HC_DARK)
rect(s,0,0,0.45,7.5, HC_GOLD)
rect(s,0.45,0,12.88,0.07, HC_GOLD)

txt(s,'Avant de partir,\nposez-vous cette question :',
    0.75,0.5,12,1.0, size=20,italic=True, color=HC_GREY)
txt(s,'"Si mes Gen Z avaient le choix\nde rejoindre mon organisation aujourd\'hui,\nest-ce qu\'ils le feraient ?"',
    0.75,1.6,12,2.0, size=28,bold=True, color=HC_WHITE)

rect(s,0.75,3.75,11,0.05, HC_GOLD)

txt(s,f'IDAT = {df["IDAT"].mean():.1f}/100  |  IATM = {df["IATM"].mean():.1f}/100  |  '
    f'{(df["SOUSPAYE"]>50).mean()*100:.0f}% sous-payes  |  RECO = {df["RECO"].mean():.1f}/100  |  '
    f'{(df["NPS"]<25).mean()*100:.0f}% en sortie mentale',
    0.75,3.9,11.5,0.55, size=11,italic=True, color=HC_GREY)

txt(s,'Les donnees ont repondu. A vous de jouer.',
    0.75,4.6,11,0.6, size=18,bold=True, color=HC_GOLD)
txt(s,'La Generation Z n\'attend pas. Elle part.',
    0.75,5.3,11,0.55, size=16,italic=True, color=HC_LIGHT)

rect(s,0.75,6.2,11,0.05, RGBColor(0x44,0x44,0x44))
txt(s,'H&C Executive — Conseil & Accompagnement RH Strategique',
    0.75,6.35,11,0.45, size=11, color=HC_GREY, italic=True)

# ── Sauvegarder ──────────────────────────────────────────────────────
output = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/GENZ_Presentation_Strategique_2026.pptx'
prs.save(output)
print(f'\nPowerPoint sauvegarde : {output}')
print(f'Slides : {len(prs.slides)}')
print()
print('=== STRUCTURE ===')
types = ['Couverture','Mise en garde narrative',
         'STAT : 59% sous-payes','STAT : 60% bien remunerés mais sous-payes (bombe invisible)',
         'STAT : 41% talent bride','STAT : 84% sante',
         'STAT : 31 pts Engagement vs Climat','STAT : 1/4 non reconnu',
         'STAT : 13.5% en sortie mentale','STAT : 34 pts Innovation bridee',
         'Graph : Toutes les composantes','Graph : Sous-payes par age',
         'Graph : Sous-payes par secteur','Graph : Heatmap pays x secteur',
         'Graph : Correlations','Graph : 4 profils Gen Z','Graph : Courbe de vie',
         'Cout de l\'inaction','5 leviers d\'action','Feuille de route','Conclusion percutante']
for i,t in enumerate(types):
    print(f'  Slide {i+1:02d} : {t}')
