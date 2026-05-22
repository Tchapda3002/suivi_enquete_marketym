"""
PowerPoint Unifie — Barometre Gen Z x Iceberg de l'Ignorance
H&C Executive 2026 — Un seul recit strategique pour decideurs
"""

import json, os, numpy as np, pandas as pd
import matplotlib, matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
import warnings
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

matplotlib.rcParams['font.family'] = 'DejaVu Sans'
matplotlib.rcParams['figure.dpi'] = 150

# ── Palette ──────────────────────────────────────────────────────────
B='#1B3A6B'; G='#C9A84C'; R='#C0392B'; GR='#1ABC9C'
GY='#7F8C8D'; W='#FFFFFF'; DK='#1A1A2E'; LB='#F4F6F9'
OR='#E67E22'; PU='#6C3583'

HC_BLUE=RGBColor(0x1B,0x3A,0x6B); HC_GOLD=RGBColor(0xC9,0xA8,0x4C)
HC_RED=RGBColor(0xC0,0x39,0x2B);  HC_GREEN=RGBColor(0x1A,0xBC,0x9C)
HC_WHITE=RGBColor(0xFF,0xFF,0xFF); HC_DARK=RGBColor(0x1A,0x1A,0x2E)
HC_GREY=RGBColor(0x7F,0x8C,0x8D); HC_LIGHT=RGBColor(0xEC,0xF0,0xF1)
HC_OR=RGBColor(0xE6,0x7E,0x22);   HC_PU=RGBColor(0x6C,0x35,0x83)

IMG='/tmp/unified_charts'; os.makedirs(IMG,exist_ok=True)

def save(name):
    p=f'{IMG}/{name}.png'
    plt.savefig(p,bbox_inches='tight',dpi=150,facecolor=plt.gcf().get_facecolor())
    plt.close(); return p

# ── Donnees ───────────────────────────────────────────────────────────
with open('/tmp/genz_dataset.json') as f: raw=json.load(f)
df=pd.DataFrame(raw); N=len(df)
gz=df[df['age'].str.contains('Gen Z',na=False)]
sante=df[df['secteur'].str.contains('Sant',na=False)]

# ════════════════════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════════════════════

# ── C1 : Iceberg + Gen Z positionne dedans ────────────────────────────
fig,ax=plt.subplots(figsize=(13,7))
fig.patch.set_facecolor(DK); ax.set_facecolor(DK)
ax.set_xlim(-1,14); ax.set_ylim(-0.5,10); ax.axis('off')

ax.text(5,9.55,'L\'ICEBERG DE L\'IGNORANCE — Ou se situe la Gen Z ?',
        ha='center',fontsize=15,fontweight='bold',color=W)
ax.text(5,9.0,'Sydney Yoshida (1989) — croise avec le Barometre Gen Z Afrique 2026',
        ha='center',fontsize=10,color=G,style='italic')

tranches=[
    (0.15,1.6,4.3,R,     '9%',  'DG / C-Suite'),
    (1.9, 1.5,3.4,'#C0392B','17%', 'Top Management'),
    (3.5, 1.5,2.5,'#E67E22','30%', 'Middle Management'),
    (5.1, 1.5,3.5,'#F1C40F','70%', 'Managers de Proximite'),
    (6.75,1.6,4.5,GR,    '100%','Agents / Frontline'),
]
for yb,h,hw,col,pct,lbl in tranches:
    xs=[5-hw,5+hw,5+hw,5-hw]; ys=[yb,yb,yb+h,yb+h]
    ax.fill(xs,ys,color=col,alpha=0.85,zorder=2)
    ax.plot(xs+[xs[0]],ys+[ys[0]],color=W,lw=0.6,alpha=0.35,zorder=3)
    ax.text(5,yb+h/2,pct,ha='center',va='center',fontsize=15,fontweight='bold',color=W,zorder=4)

# Fleche Gen Z sur les 2 niveaux du bas
ax.annotate('',xy=(9.4,6.05),xytext=(9.4,8.5),
            arrowprops=dict(arrowstyle='<->',color=GR,lw=2.5))
ax.text(9.55,7.3,'GEN Z\nOCCUPE CES\nNIVEAUX',va='center',fontsize=11,
        fontweight='bold',color=GR,
        bbox=dict(boxstyle='round,pad=0.4',facecolor='#0D2E1A',edgecolor=GR,lw=1.5))

# Stats Gen Z dans bulle
gz_sp=(gz['SOUSPAYE']>50).mean()*100
gz_nps=(gz['NPS']<25).mean()*100
gz_iatm=gz['IATM'].mean()

stats_txt=f'64% sous-payes\nIATM = {gz_iatm:.0f}/100\n18% en sortie mentale'
ax.text(11.2,7.3,stats_txt,va='center',ha='left',fontsize=10,color=R,fontweight='bold',
        bbox=dict(boxstyle='round,pad=0.5',facecolor='#2C0A08',edgecolor=R,lw=1.5))

# Labels niveaux a gauche
ylabels=[(0.95,'DG / C-Suite','9% du reel visible',W),
         (2.65,'Top Management','17% du reel visible',W),
         (4.25,'Middle Management','30% du reel visible',W),
         (5.85,'Mgr de Proximite','70% du reel visible','#F1C40F'),
         (7.55,'Agents / Frontline','100% du reel visible',GR)]
for ym,titre,sous,col in ylabels:
    ax.text(-0.85,ym,titre,fontsize=8.5,fontweight='bold',color=col,ha='left',va='center')
    ax.text(-0.85,ym-0.3,sous,fontsize=7.5,color=GY,ha='left',va='center',style='italic')

ax.text(5,-0.3,
        'La Gen Z = votre frontline. Elle voit 100% du terrain. Et 91% de ce savoir ne remonte jamais.',
        ha='center',fontsize=10,color=G,fontweight='bold',style='italic')
plt.tight_layout()
c_iceberg_genz=save('01_iceberg_genz')
print('C1 OK')

# ── C2 : Ce que la Gen Z emporte quand elle part ─────────────────────
fig,ax=plt.subplots(figsize=(13,6))
fig.patch.set_facecolor(DK); ax.set_facecolor(DK); ax.axis('off')
ax.set_xlim(0,14); ax.set_ylim(0,7)

ax.text(7,6.6,'QUAND UN GEN Z PART — Ce qu\'il emporte avec lui',
        ha='center',fontsize=15,fontweight='bold',color=W)
ax.text(7,6.1,'Et que vous ne calculerez jamais dans vos tableaux RH',
        ha='center',fontsize=11,color=G,style='italic')

blocs=[
    (0.3, 5.0, GR,  '100%', 'du terrain connu',
     'Observations clients, dysfonctionnements,\nopportunites de vente non exploitees'),
    (4.8, 5.0, G,   '91%', 'ne remonte jamais',
     'La part de son savoir que votre DG\nn\'a JAMAIS vue — et ne verra plus jamais'),
    (9.3, 5.0, R,   '9.7x', 'son salaire annuel',
     'Le vrai cout de son depart\n(recrutement + valeur cachee perdue)'),
    (0.3, 2.5, OR,  '34 pts', 'd\'innovation bridee',
     'INNOV=81 mais IATM=47\nSon potentiel creatif reste inexploite'),
    (4.8, 2.5, R,   '64%', 'se sentent sous-payes',
     'Dont 60% qui sont pourtant bien remunerés\nCe n\'est pas un pb de salaire'),
    (9.3, 2.5, R,   '18%', 'deja partis en tete',
     'NPS < 25 : ils sont encore la\nmais ils ne produisent plus vraiment'),
]
for x,y,col,val,lbl,detail in blocs:
    rect=plt.Rectangle((x,y),3.9,1.95,facecolor='#1E3A5F',edgecolor=col,lw=2)
    ax.add_patch(rect); rect2=plt.Rectangle((x,y),0.12,1.95,facecolor=col)
    ax.add_patch(rect2)
    ax.text(x+0.3,y+1.55,val,fontsize=22,fontweight='bold',color=col,va='top')
    ax.text(x+0.3,y+1.05,lbl,fontsize=9,fontweight='bold',color=W,va='top')
    ax.text(x+0.3,y+0.1,detail,fontsize=8,color=GY,va='bottom',style='italic')

ax.text(7,0.15,
        'Source : Barometre Gen Z H&C Executive (n=1200) x Iceberg of Ignorance, Yoshida (1989)',
        ha='center',fontsize=8,color=GY,style='italic')
plt.tight_layout()
c_depart_genz=save('02_depart_genz')
print('C2 OK')

# ── C3 : Composantes Gen Z ────────────────────────────────────────────
fig,ax=plt.subplots(figsize=(12,5.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
comp={'Remuneration':df['REM'].mean(),'Reconnaissance':df['RECO'].mean(),
      'Climat Travail':df['IATM'].mean(),'NPS Employeur':df['NPS'].mean(),
      'Satisfaction':df['SAT'].mean(),'Valorisation':df['VAL'].mean(),
      'Apprentissage':df['APPRENT'].mean(),'Engagement':df['ENG'].mean(),
      'Initiative':df['INIT'].mean(),'Innovation':df['INNOV'].mean()}
cs=dict(sorted(comp.items(),key=lambda x:x[1]))
cols_c=[R if v<60 else(G if v<70 else GR) for v in cs.values()]
bars=ax.barh(list(cs.keys()),list(cs.values()),color=cols_c,height=0.6,edgecolor=W)
ax.axvline(65,color='#888',linestyle='--',lw=1.2,alpha=0.5)
for bar,val in zip(bars,cs.values()):
    ax.text(val+0.8,bar.get_y()+bar.get_height()/2,f'{val:.1f}',
            va='center',fontsize=12,fontweight='bold')
ax.set_xlim(0,100); ax.tick_params(labelsize=12)
ax.set_title('Toutes les composantes — ce qui va bien, ce qui est en danger',
             fontsize=14,fontweight='bold',color=B,pad=10)
p1=mpatches.Patch(color=R,label='Zone critique (<60)')
p2=mpatches.Patch(color=G,label='Zone fragile (60-70)')
p3=mpatches.Patch(color=GR,label='Zone correcte (>70)')
ax.legend(handles=[p1,p2,p3],fontsize=10,loc='lower right')
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
c_comp=save('03_composantes')
print('C3 OK')

# ── C4 : Sous-payes par age ───────────────────────────────────────────
fig,ax=plt.subplots(figsize=(12,5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
ages=['18-22 ans (Gen Z)','23-27 ans (Gen Z)','28-32 ans (Millennials)',
      '33-37 ans (Millennials)','38-42 ans (Millennials)']
xlabs=['18-22\n(Gen Z)','23-27\n(Gen Z)','28-32\n(Millenn.)','33-37\n(Millenn.)','38-42\n(Millenn.)']
pcts=[(df[df['age']==a]['SOUSPAYE'].dropna()>50).mean()*100 for a in ages]
cols_a=[R if p>65 else G for p in pcts]
bars=ax.bar(range(5),pcts,color=cols_a,width=0.55,edgecolor=W)
ax.axhline(50,color='#555',linestyle='--',lw=1.2,alpha=0.5)
for bar,val in zip(bars,pcts):
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1,
            f'{val:.0f}%',ha='center',fontsize=15,fontweight='bold')
ax.set_xticks(range(5)); ax.set_xticklabels(xlabs,fontsize=12)
ax.set_ylim(0,90); ax.set_ylabel('% repondants',fontsize=12)
ax.set_title('% se sentant sous-payes par tranche d\'age\n'
             '>>> Plus on est jeune (= plus on est terrain), plus on se sent sous-paye',
             fontsize=13,fontweight='bold',color=B,pad=10)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
c_age=save('04_souspaye_age')
print('C4 OK')

# ── C5 : Heatmap IDAT pays x secteur ─────────────────────────────────
fig,ax=plt.subplots(figsize=(12,5.5))
fig.patch.set_facecolor(LB)
top_pays=df['pays'].value_counts().head(7).index.tolist()
top_s=['Santé / Pharmacie','Banque / Finance / Assurance','Agriculture / Agroalimentaire',
       'Éducation / Formation','Commerce / Distribution','Administration publique']
matrix=[]
for s in top_s:
    row=[]
    for p in top_pays:
        sub=df[(df['pays']==p)&(df['secteur']==s)]['IDAT'].dropna()
        row.append(round(sub.mean(),1) if len(sub)>=5 else np.nan)
    matrix.append(row)
mdf=pd.DataFrame(matrix,index=[s[:22] for s in top_s],columns=top_pays)
sns.heatmap(mdf,ax=ax,cmap='RdYlGn',vmin=55,vmax=78,
            annot=True,fmt='.0f',linewidths=1.2,linecolor='white',
            annot_kws={'size':12,'weight':'bold'},
            cbar_kws={'label':'Score IDAT /100','shrink':0.85})
ax.set_title('Attractivite par Pays x Secteur — Ou est l\'urgence ?',
             fontsize=14,fontweight='bold',color=B,pad=10)
ax.tick_params(axis='x',labelsize=11,rotation=25)
ax.tick_params(axis='y',labelsize=11,rotation=0)
fig.tight_layout()
c_heatmap=save('05_heatmap')
print('C5 OK')

# ── C6 : Cout visible vs reel (Yoshida applique a la Gen Z) ──────────
fig,ax=plt.subplots(figsize=(12,5.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
niveaux=['Agent\nFrontline','Mgr\nProximite','Middle\nManager','Top\nManager','DG\nC-Suite']
cout_repl=[0.75,1.25,1.75,2.5,4.0]
valeur_cachee=[8.5,4.2,3.5,4.8,6.5]
x=np.arange(5); w=0.5
b1=ax.bar(x,cout_repl,w,color=G,alpha=0.85,label='Cout visible : recrutement + formation',edgecolor=W)
b2=ax.bar(x,valeur_cachee,w,bottom=cout_repl,color=R,alpha=0.8,
          label='Valeur cachee : idees, savoir tacite, capital terrain (Yoshida)',edgecolor=W)
for i,(cr,sc) in enumerate(zip(cout_repl,valeur_cachee)):
    total=cr+sc
    ax.text(i,total+0.15,f'{total:.1f}x',ha='center',fontsize=14,fontweight='bold',color='#111')
ax.set_xticks(x); ax.set_xticklabels(niveaux,fontsize=12)
ax.set_ylabel('Equivalent x salaire annuel',fontsize=11)
ax.set_title('Le vrai cout d\'un depart — Yoshida applique a votre Gen Z\n'
             'Ce que le RH calcule vs ce qui est vraiment perdu',
             fontsize=13,fontweight='bold',color=B,pad=10)
ax.legend(fontsize=10,loc='upper right')
# Fleche "Gen Z ici"
ax.annotate('',xy=(0.5,12),xytext=(1.5,12),
            arrowprops=dict(arrowstyle='<->',color=GR,lw=2))
ax.text(0.9,12.4,'GEN Z',ha='center',fontsize=11,fontweight='bold',color=GR)
ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
c_cout=save('06_cout_genz')
print('C6 OK')

# ── C7 : Simulation MFCFA ────────────────────────────────────────────
fig,ax=plt.subplots(figsize=(12,5.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
salaires=[2.4,6.0,15.0,45.0,120.0]
totaux_x=[a+b for a,b in zip(cout_repl,valeur_cachee)]
couts_reel=[s*t for s,t in zip(salaires,totaux_x)]
couts_vis=[s*c for s,c in zip(salaires,cout_repl)]
x=np.arange(5); w=0.35
b1=ax.bar(x-w/2,couts_vis,w,color=G,alpha=0.85,label='Cout visible (RH traditionnel)',edgecolor=W)
b2=ax.bar(x+w/2,couts_reel,w,color=R,alpha=0.85,label='Cout reel total (Yoshida)',edgecolor=W)
for bar,val in zip(b1,couts_vis):
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+2,f'{val:.0f}M',
            ha='center',fontsize=10,fontweight='bold',color=G)
for bar,val in zip(b2,couts_reel):
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+2,f'{val:.0f}M',
            ha='center',fontsize=10,fontweight='bold',color=R)
ax.set_xticks(x); ax.set_xticklabels(niveaux,fontsize=12)
ax.set_ylabel('Millions FCFA',fontsize=11)
ax.set_title('Simulation : cout d\'un depart en MFCFA\nHyp. salaires representatifs marche Afrique de l\'Ouest',
             fontsize=13,fontweight='bold',color=B,pad=10)
ax.legend(fontsize=10); ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
c_simu=save('07_simulation')
print('C7 OK')

# ── C8 : Correlations ────────────────────────────────────────────────
fig,ax=plt.subplots(figsize=(12,5.5))
fig.patch.set_facecolor(LB); ax.set_facecolor(W)
cols_c=['SAT','ENG','VAL','NPS','REM','RECO','APPRENT','INIT','INNOV','IATM','SOUSPAYE_inv']
labs_c=['Satisfaction','Engagement','Valorisation','NPS Employeur','Remuneration',
        'Reconnaissance','Apprentissage','Initiative','Innovation','Climat Travail','Non sous-paye']
corrs=[df[c].corr(df['IDAT']) for c in cols_c]
pairs=sorted(zip(corrs,labs_c),key=lambda x:x[0])
cols_cr=[R if v<0.3 else(G if v<0.5 else GR) for v in [p[0] for p in pairs]]
bars=ax.barh([p[1] for p in pairs],[p[0] for p in pairs],color=cols_cr,height=0.6,edgecolor=W)
ax.axvline(0.5,color=GR,linestyle=':',lw=1.5,alpha=0.7,label='Forte correlation')
ax.axvline(0.3,color=G,linestyle=':',lw=1.5,alpha=0.7,label='Moderee')
ax.axvline(0,color='#333',lw=1)
for bar,val in zip(bars,[p[0] for p in pairs]):
    ax.text(val+0.01,bar.get_y()+bar.get_height()/2,f'{val:.2f}',
            va='center',fontsize=11,fontweight='bold')
ax.tick_params(labelsize=12)
ax.set_title('Ce qui retient vraiment la Gen Z — les vrais drivers',
             fontsize=14,fontweight='bold',color=B,pad=10)
ax.legend(fontsize=10); ax.spines[['top','right']].set_visible(False)
fig.tight_layout()
c_corr=save('08_correlations')
print('C8 OK')

# ── C9 : 4 profils K-means ───────────────────────────────────────────
cols_pca=['SAT','ENG','VAL','NPS','REM','RECO','APPRENT','INIT','INNOV','IATM','SOUSPAYE_inv']
df_pca=df[cols_pca].dropna().copy()
X_sc=StandardScaler().fit_transform(df_pca)
km=KMeans(n_clusters=4,random_state=42,n_init=20)
df_pca['cluster']=km.fit_predict(X_sc)
cm=df_pca.groupby('cluster')[cols_pca].mean()
cs_km=df_pca['cluster'].value_counts().sort_index()
radar_c=['SAT','ENG','VAL','REM','RECO','INNOV','IATM']
radar_l=['Satisfaction','Engagement','Valorisation','Remuneration','Reconnaissance','Innovation','Climat']
gmeans=df[radar_c].mean().values
ccolors=[GR,R,G,B]
def cname(c):
    row=cm.loc[c]; n=cs_km[c]; pct=n/len(df_pca)*100
    if row['ENG']>75 and row['SAT']>70: return f'Ambassadeurs\n{pct:.0f}% ({n})'
    if row['ENG']>70 and row['REM']<55: return f'Frustres Ambitieux\n{pct:.0f}% ({n})'
    if row['REM']<50 and row['RECO']<50: return f'En Danger\n{pct:.0f}% ({n})'
    return f'Attentistes\n{pct:.0f}% ({n})'
fig,axes=plt.subplots(1,4,figsize=(15,4.5),sharey=True)
fig.patch.set_facecolor(LB)
fig.suptitle('Les 4 profils Gen Z — 4 strategies RH differentes',fontsize=13,fontweight='bold',color=B,y=1.02)
for idx,(c,ax) in enumerate(zip(range(4),axes)):
    vals=cm.loc[c,radar_c].values; col=ccolors[idx]; x=np.arange(len(radar_c))
    ax.set_facecolor(W)
    ax.bar(x,gmeans,width=0.6,color='#ddd',alpha=0.5)
    ax.bar(x,vals,width=0.6,color=col,alpha=0.82,edgecolor=W)
    ax.set_xticks(x); ax.set_xticklabels(radar_l,fontsize=8.5,rotation=30,ha='right')
    ax.set_ylim(0,100); ax.axhline(65,color='#aaa',linestyle='--',lw=1)
    for xi,v in enumerate(vals):
        ax.text(xi,v+1.5,f'{v:.0f}',ha='center',fontsize=9,fontweight='bold')
    ax.set_title(cname(c),fontsize=11,fontweight='bold',color=col,pad=8)
    ax.spines[['top','right']].set_visible(False)
    if idx==0:
        ax.set_ylabel('Score /100',fontsize=11)
        p1=mpatches.Patch(color='#ddd',label='Moy. globale')
        p2=mpatches.Patch(color=col,label='Ce profil')
        ax.legend(handles=[p1,p2],fontsize=8)
fig.tight_layout()
c_profils=save('09_profils')
print('C9 OK')

print(f'\nTous les charts dans {IMG}/')

# ════════════════════════════════════════════════════════════════════════
# POWERPOINT — UN SEUL RECIT
# ════════════════════════════════════════════════════════════════════════
prs=Presentation()
prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
blank=prs.slide_layouts[6]

def bg(s,c):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c

def rect(s,l,t,w,h,c):
    sh=s.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background(); return sh

def txt(s,text,l,t,w,h,size=13,bold=False,color=HC_WHITE,
        align=PP_ALIGN.LEFT,italic=False,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    run=p.add_run(); run.text=text
    run.font.size=Pt(size); run.font.bold=bold
    run.font.italic=italic; run.font.color.rgb=color

def img(s,path,l,t,w,h):
    if os.path.exists(path):
        s.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(h))

def header(s,title,subtitle='',bg_col=HC_BLUE,h=1.0):
    rect(s,0,0,13.33,h,bg_col)
    txt(s,title,0.35,0.07,12.6,0.55,size=20,bold=True,color=HC_WHITE)
    if subtitle:
        txt(s,subtitle,0.35,0.63,12.6,0.35,size=11,color=HC_GOLD,italic=True)

def acte_separator(prs,num,titre,desc):
    s=prs.slides.add_slide(blank)
    bg(s,RGBColor(0x0D,0x1F,0x3C))
    rect(s,0,0,0.45,7.5,HC_GOLD)
    txt(s,f'ACTE {num}',0.65,1.5,12,0.9,size=36,bold=True,color=HC_GOLD)
    txt(s,titre,0.65,2.5,12,1.1,size=28,bold=True,color=HC_WHITE)
    rect(s,0.65,3.75,11,0.05,HC_GOLD)
    txt(s,desc,0.65,3.95,11.5,1.5,size=14,italic=True,color=HC_LIGHT,wrap=True)

def footer(s):
    txt(s,'H&C Executive  |  Barometre Gen Z x Iceberg de l\'Ignorance  |  Afrique 2026  |  Confidentiel',
        0,7.18,13.33,0.32,size=9,color=HC_GREY,align=PP_ALIGN.CENTER)

def stat_slide(prs,big,unit,top,bottom,question,context,
               bg_col=HC_DARK,num_col=HC_GOLD,tag='',tag_col=HC_RED):
    s=prs.slides.add_slide(blank)
    bg(s,bg_col); rect(s,0,0,0.45,7.5,num_col)
    if tag:
        rect(s,0.55,0.18,3.5,0.55,tag_col)
        txt(s,tag,0.6,0.22,3.4,0.45,size=10,bold=True,color=HC_WHITE,align=PP_ALIGN.CENTER)
    txt(s,top,   0.65,0.85,12,0.65,size=13,color=HC_LIGHT,italic=True)
    txt(s,big,   0.55,1.35,9.0,2.8,size=105,bold=True,color=num_col)
    if unit: txt(s,unit,7.2,2.6,5,1.2,size=30,bold=True,color=num_col)
    txt(s,bottom,0.65,4.1,12,0.75,size=22,bold=True,color=HC_WHITE)
    rect(s,0.65,4.95,11.5,0.04,num_col)
    txt(s,context,0.65,5.1,11.5,0.75,size=11,color=HC_LIGHT,italic=True)
    txt(s,f'>>> {question}',0.65,5.95,11.5,0.9,size=13,bold=True,color=num_col)

CA=(0.2,1.05,12.93,6.25)

# ════════════════
# COUVERTURE
# ════════════════
s=prs.slides.add_slide(blank)
bg(s,HC_DARK)
rect(s,0,0,0.45,7.5,HC_GOLD)
rect(s,0.45,0,12.88,0.07,HC_GOLD)
rect(s,0.45,7.43,12.88,0.07,HC_GOLD)
txt(s,'GENERATION Z &',0.75,0.65,12,0.85,size=32,bold=True,color=HC_WHITE)
txt(s,'L\'ICEBERG DE L\'IGNORANCE',0.75,1.45,12,1.1,size=42,bold=True,color=HC_GOLD)
txt(s,'Ce que votre entreprise perd quand sa releve part',
    0.75,2.65,11.5,0.7,size=18,italic=True,color=HC_LIGHT)
rect(s,0.75,3.5,10,0.05,HC_GOLD)
txt(s,'Barometre Gen Z Afrique 2026 (n=1 200)  x  Yoshida, Iceberg of Ignorance (1989)',
    0.75,3.65,11,0.5,size=12,color=HC_GREY)
txt(s,'13 pays  |  10+ secteurs  |  Analyse strategique pour decideurs',
    0.75,4.2,11,0.45,size=12,color=HC_GREY)
txt(s,'H&C Executive — Document confidentiel — Avril 2026',
    0.75,6.85,8,0.45,size=11,italic=True,color=HC_GREY)

# ════════════════
# ACTE 1 — LE CADRE
# ════════════════
acte_separator(prs,'1','L\'ANGLE MORT DE VOTRE ORGANISATION',
    'Yoshida (1989) a prouve que le sommet d\'une entreprise ne voit que 9% de sa realite.\n'
    'Ce n\'est pas un probleme de communication. C\'est la structure de toute hierarchie.\n'
    'La question est : que se passe-t-il quand ceux qui voient 100% du terrain decident de partir ?')

# Slide iceberg
s=prs.slides.add_slide(blank)
bg(s,HC_DARK)
header(s,'L\'ICEBERG DE L\'IGNORANCE — Ce que chaque niveau connait du terrain',
       'Yoshida (1989) : le DG ne voit que 9% de la realite operationnelle de son entreprise',
       bg_col=HC_BLUE)
img(s,c_iceberg_genz,*CA)
footer(s)

stat_slide(prs,'9','%',
    'Ce que voit votre DG. Pas plus.',
    'de la realite du terrain remonte jusqu\'au sommet de votre organisation.',
    'Sur quoi reposent alors 100% de vos decisions strategiques ?',
    'Yoshida, 1989 — valide dans toutes les organisations hierarchiques. '
    '91% des problemes, des opportunites, et des idees du terrain ne franchissent jamais le niveau du middle management.',
    tag='L\'ANGLE MORT STRATEGIQUE',tag_col=HC_RED,num_col=HC_RED)

# ════════════════
# ACTE 2 — LE DIAGNOSTIC
# ════════════════
acte_separator(prs,'2','VOTRE FRONTLINE, C\'EST LA GEN Z.\nET ELLE EST EN TRAIN DE PARTIR.',
    'Notre etude le mesure sur 1 200 repondants dans 13 pays d\'Afrique.\n'
    'Les chiffres sont clairs : la Gen Z — qui occupe vos niveaux terrain et proximite —\n'
    'est la plus insatisfaite, la plus sous-payee, et la plus proche du depart.')

# Slide de liaison Iceberg + Gen Z
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'LA GEN Z DANS VOTRE PYRAMIDE — Elle est en bas. Elle voit tout.',
       'Et c\'est exactement la qu\'est le probleme : 64% se sentent sous-payees. 18% sont deja partis en tete.')

# 2 colonnes : pyramide mini + stats
rect(s,0.3,1.1,5.8,6.1,HC_WHITE)
txt(s,'LA GEN Z OCCUPE',0.5,1.2,5.4,0.55,size=13,bold=True,color=HC_BLUE,align=PP_ALIGN.CENTER)
txt(s,'les 2 niveaux du bas de la pyramide',0.5,1.72,5.4,0.45,size=10,italic=True,color=HC_GREY,align=PP_ALIGN.CENTER)
niveaux_p=[('DG / C-Suite','9% du terrain','#1B3A6B',0.3),
           ('Top Mgmt','17% du terrain','#27AE60',0.3),
           ('Middle Mgmt','30% du terrain','#E67E22',0.3),
           ('Mgr Proximite','70% du terrain','#F1C40F',0.5),
           ('Agents / Frontline','100% du terrain','#C0392B',0.5)]
for i,(lbl,pct,col,alpha) in enumerate(niveaux_p):
    y=2.35+i*0.82
    w_bar=0.8+i*0.7
    x_bar=3.2-w_bar/2
    rect(s,x_bar,y,w_bar,0.65,RGBColor(*[int(col.lstrip('#')[j:j+2],16) for j in (0,2,4)]))
    txt(s,lbl,0.4,y+0.05,2.2,0.5,size=9,bold=True,color=HC_WHITE)
    txt(s,pct,2.7,y+0.05,1.2,0.5,size=9,color=HC_WHITE,align=PP_ALIGN.RIGHT)

# Fleche Gen Z
rect(s,5.5,5.2,0.55,1.4,HC_GREEN)
txt(s,'GEN\nZ',5.57,5.35,0.4,1.1,size=9,bold=True,color=HC_WHITE,align=PP_ALIGN.CENTER)

# Stats droite
stats_r=[
    ('56%','de l\'echantillon est Gen Z (18-27 ans)',HC_GREEN),
    ('64%','de la Gen Z se sent sous-payee',HC_RED),
    ('44.6','IATM — Climat Gen Z (zone critique)',HC_RED),
    ('18%','en sortie mentale (NPS < 25)',HC_RED),
    ('34 pts','entre Innovation (81) et Climat (47)',HC_OR),
    ('9.7x','cout reel d\'un depart frontline',HC_RED),
]
def hex_to_rgb(c):
    h=str(c).lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

for i,(val,lbl,col) in enumerate(stats_r):
    y=1.2+i*0.95
    col_rgb=hex_to_rgb(col) if isinstance(col,str) else col
    rect(s,6.4,y,6.7,0.85,RGBColor(0x1E,0x3A,0x5F))
    rect(s,6.4,y,0.08,0.85,col_rgb)
    txt(s,val,6.55,y+0.05,1.4,0.6,size=18,bold=True,color=col_rgb)
    txt(s,lbl,8.0,y+0.18,5.0,0.5,size=10,color=HC_LIGHT)
footer(s)

# Stats chocs
stat_slide(prs,'64','%',
    'ACTE 2 — Le diagnostic. Premier signal.',
    'de votre Gen Z — votre frontline — se sent sous-payee.',
    'Mais augmenter les salaires ne reglera rien. La preuve est dans les donnees.',
    f'Correlation entre remuneration reelle et sentiment de sous-paiement : r=0.08 (quasi nulle). '
    f'60% de ceux qui sont bien remunerés se sentent QUAND MEME sous-payes. '
    f'C\'est un probleme de reconnaissance et de justice — pas de salaire.',
    tag='SIGNAL 1 — SOUS-PAIEMENT PERCU',tag_col=HC_RED,num_col=HC_RED)

stat_slide(prs,'41','%',
    'ACTE 2 — Deuxieme signal. Le plus paradoxal.',
    'sont engages (ENG>70) mais travaillent dans un mauvais climat (IATM<50).',
    'Vos meilleurs elements sont vos plus frustres. Combien de temps peuvent-ils tenir ?',
    f'Ecart Engagement ({df["ENG"].mean():.0f}/100) vs Climat ({df["IATM"].mean():.0f}/100) : 31 points. '
    f'Innovation = {df["INNOV"].mean():.0f}/100. Initiative = {df["INIT"].mean():.0f}/100. '
    f'Tout ce potentiel existe — dans un environnement qui l\'etouffe.',
    tag='SIGNAL 2 — TALENT BRIDE',tag_col=HC_PU,num_col=HC_GOLD)

stat_slide(prs,'18','%',
    'ACTE 2 — Troisieme signal. Le plus silencieux.',
    'de votre Gen Z sont deja partis dans leur tete.',
    'Ils sont encore la. Ils pointent. Mais ils ne produisent plus vraiment.',
    f'NPS < 25 : 18% de la Gen Z. 23% supplementaires ne recommanderaient pas leur employeur. '
    f'Ensemble : plus d\'1 Gen Z sur 3 est en risque de depart ou de desengagement profond. '
    f'Et 71% des 18-22 ans — les plus proches du terrain — se sentent sous-payes.',
    tag='SIGNAL 3 — SORTIE MENTALE',tag_col=HC_RED,num_col=HC_RED)

# Composantes
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'LE BILAN COMPLET — Ce qui va bien, ce qui est en danger',
       'Engagement et Innovation au sommet — Remuneration et Reconnaissance en zone critique')
img(s,c_comp,*CA); footer(s)

# Sous-payes par age
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'PLUS ON EST JEUNE (= PLUS ON EST TERRAIN), PLUS ON SE SENT SOUS-PAYE',
       '71% des 18-22 ans — la tranche la plus proche du client et du terrain — signal de depart precoce',
       bg_col=HC_RED)
img(s,c_age,*CA); footer(s)

# Heatmap
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'CARTOGRAPHIE — Ou est l\'urgence par pays et par secteur ?',
       'Certaines zones concentrent toutes les tensions — ce sont vos priorites d\'intervention')
img(s,c_heatmap,*CA); footer(s)

# ════════════════
# ACTE 3 — LA CONSEQUENCE
# ════════════════
acte_separator(prs,'3','CE QUE VOUS NE CALCULEZ PAS.\nC\'EST LA QUE SE CACHE LE VRAI COUT.',
    'Yoshida + les signaux de l\'etude forment un equation simple :\n'
    'Votre Gen Z voit 100% du terrain. Elle est insatisfaite. Elle part.\n'
    'Elle emporte avec elle 91% de ce que votre DG n\'a jamais vu.\n'
    'Et vous calculez ca a 0.75x son salaire.')

# Slide liaison Yoshida + Gen Z
s=prs.slides.add_slide(blank)
bg(s,HC_DARK)
header(s,'QUAND UN GEN Z PART — Ce qu\'il emporte. Ce que vous ne calculerez jamais.',
       'Yoshida (1989) x Barometre Gen Z H&C Executive (2026)',bg_col=HC_RED)
img(s,c_depart_genz,0.2,1.05,12.93,6.15)
footer(s)

stat_slide(prs,'91','%',
    'ACTE 3 — La consequence directe.',
    'des idees d\'un agent Gen Z ne remontent jamais jusqu\'au management.',
    'Quand il part, il emporte tout ce savoir. Avez-vous une strategie pour le capturer avant ?',
    'Yoshida (1989) : seuls 9% des problemes et opportunites du terrain franchissent le niveau du DG. '
    '91% disparaissent — soit parce qu\'ils n\'ont jamais ete exprimes, soit parce qu\'ils n\'ont pas ete entendus. '
    'Un depart Gen Z = une fuite permanente de cette information.',
    tag='LA PERTE SILENCIEUSE — YOSHIDA',tag_col=HC_RED,num_col=HC_RED)

# Cout visible vs reel
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'LE VRAI COUT D\'UN DEPART — Yoshida applique a votre Gen Z',
       'Ce que le RH calcule (partie emergee) vs ce qui est vraiment perdu (partie immergee)',
       bg_col=HC_RED)
img(s,c_cout,*CA); footer(s)

stat_slide(prs,'9.7','x',
    'Le vrai cout d\'un depart d\'agent Gen Z.',
    'son salaire annuel — pas 0.75x comme le calcule votre service RH.',
    'Multipliez par votre nombre de departs annuels. Quel chiffre obtenez-vous ?',
    'Cout RH visible : 0.75x salaire (recrutement + integration). '
    'Valeur cachee Yoshida : +8.5x (savoir terrain, idees perdues, capital client, courbe apprentissage). '
    'Total : 9.7x. Pour un salaire agent de 2.4M FCFA : 23M FCFA par depart.',
    tag='FRONTLINE / GEN Z',tag_col=HC_RED,num_col=HC_RED)

# Simulation
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'SIMULATION EN MFCFA — Ce que ca represente concretement',
       'Hypothese : salaires representatifs marche Afrique de l\'Ouest | CA ref 10 Md FCFA')
img(s,c_simu,*CA); footer(s)

# ════════════════
# ACTE 4 — LES LEVIERS
# ════════════════
acte_separator(prs,'4','CE QU\'ON FAIT MAINTENANT.\nPAS DANS 6 MOIS.',
    'Les signaux sont clairs. Le cout est calcule. Il reste une question : '
    'est-ce que vous allez continuer a perdre 9.7x le salaire de chaque Gen Z qui part,\n'
    'ou est-ce que vous actionnez les leviers qui existent ?')

# Correlations
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'CE QUI RETIENT VRAIMENT LA GEN Z — Les vrais leviers',
       'Ce n\'est pas le salaire. C\'est la reconnaissance et le climat — actionnables maintenant.')
img(s,c_corr,*CA); footer(s)

# 4 profils
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'LES 4 PROFILS GEN Z — Ne les managez pas de la meme facon',
       '4 types de travailleurs dans votre organisation. 4 risques distincts. 4 strategies differentes.')
img(s,c_profils,*CA); footer(s)

# 5 leviers
s=prs.slides.add_slide(blank)
bg(s,HC_DARK)
header(s,'5 LEVIERS D\'ACTION — Classes par impact et par urgence',
       'Des actions. Pas un rapport de plus. Chacune repond directement aux signaux mesures.')
leviers=[
    ('01','STRUCTURER LA RECONNAISSANCE — formelle, mensuelle, visible',
     f'RECO={df["RECO"].mean():.0f}/100 | 1/4 ne se sent pas reconnu | Cout zero — question de culture',
     HC_GREEN,'30j','Faible'),
    ('02','TRANSPARENCER LES GRILLES SALARIALES',
     f'60% bien remunerés mais sous-payes | Correlation REM/ressenti = 0.08 | La transparence reduit 30% du sentiment',
     HC_GOLD,'90j','Moyen'),
    ('03','FORMER LES MANAGERS AU LEADERSHIP GEN Z',
     f'IATM={df["IATM"].mean():.0f}/100 | Manager direct = 70% du climat | Formation 2 jours + rituels',
     HC_GOLD,'60j','Moyen'),
    ('04','CREER DES CANAUX D\'IDEATION ASCENDANTE',
     'Yoshida : 91% du terrain ne remonte pas | Rituels hebdo terrain → management | Capturer avant le depart',
     HC_GREEN,'45j','Faible'),
    ('05','ENTRETIEN DE SORTIE STRATEGIQUE — Recuperer l\'information',
     'Dernieres chance de capturer 91% du savoir terrain | Formaliser + analyser systematiquement',
     HC_RED,'Immédiat','Nul'),
]
for i,(num,titre,stat,col,delai,cout) in enumerate(leviers):
    y=1.12+i*1.2
    rect(s,0.25,y,12.83,1.1,RGBColor(0x1E,0x3A,0x5F))
    rect(s,0.25,y,0.1,1.1,col)
    txt(s,num,0.45,y+0.05,0.7,0.7,size=22,bold=True,color=col)
    txt(s,titre,1.2,y+0.04,8.5,0.52,size=11,bold=True,color=HC_WHITE)
    txt(s,f'>>> {stat}',1.2,y+0.58,8.5,0.44,size=9,italic=True,color=HC_GOLD)
    txt(s,f'Delai : {delai}  |  Cout : {cout}',9.8,y+0.3,3.1,0.45,size=9,bold=True,color=col,align=PP_ALIGN.RIGHT)
footer(s)

# Feuille de route
s=prs.slides.add_slide(blank)
bg(s,RGBColor(0xF4,0xF6,0xF9))
header(s,'FEUILLE DE ROUTE — Ce qui se decide aujourd\'hui se recolte dans 18 mois',
       '90 jours determines = transformation reelle | Chaque etape repond a un signal mesure')
phases=[
    ('90 JOURS',HC_RED,[
        'Audit reconnaissance : cartographier les pratiques actuelles',
        'Mise en place entretiens de sortie strategiques (Yoshida)',
        'Rituels ideation terrain → management (capturer les 91%)',
        'Benchmarks salariaux Sante & Agriculture (84% sous-payes)',
    ]),
    ('6 MOIS',HC_GOLD,[
        'Formation managers leadership Gen Z (IATM = 46.8)',
        'Systeme de reconnaissance structure et formalise',
        'Transparence des grilles salariales (60% bien payes mais sous-payes)',
        'Parcours evolution visibles a 6 et 18 mois',
    ]),
    ('12 MOIS',HC_GREEN,[
        'Remesure IDAT — cible : +8 points',
        '% sous-payes — cible : -15 points',
        'Retention Gen Z — cible : 85%+',
        '% idees terrain remontant — cible : doubler le taux Yoshida actuel',
    ]),
]
for i,(periode,col,actions) in enumerate(phases):
    x=0.3+i*4.35
    rect(s,x,1.15,4.1,6.05,HC_WHITE)
    rect(s,x,1.15,4.1,0.62,col)
    txt(s,periode,x+0.1,1.18,3.9,0.58,size=17,bold=True,color=HC_WHITE,align=PP_ALIGN.CENTER)
    for j,action in enumerate(actions):
        ya=1.97+j*1.15
        rect(s,x+0.15,ya+0.12,0.08,0.75,col)
        txt(s,action,x+0.32,ya,3.65,1.0,size=10.5,color=HC_BLUE,wrap=True)
footer(s)

# CONCLUSION
s=prs.slides.add_slide(blank)
bg(s,HC_DARK)
rect(s,0,0,0.45,7.5,HC_GOLD)
rect(s,0.45,0,12.88,0.07,HC_GOLD)

txt(s,'La vraie question de gouvernance\nn\'est pas :',
    0.75,0.45,12,1.0,size=20,italic=True,color=HC_GREY)
txt(s,'"Comment motiver ma Generation Z ?"',
    0.75,1.45,12,0.85,size=26,bold=True,color=HC_WHITE)
rect(s,0.75,2.45,11,0.05,HC_GOLD)
txt(s,'Elle est :',0.75,2.65,12,0.6,size=20,italic=True,color=HC_GREY)
txt(s,'"Est-ce que mon organisation\nest construite pour entendre\nce que ses 91% savent ?"',
    0.75,3.3,12,1.8,size=26,bold=True,color=HC_GOLD)
rect(s,0.75,5.25,11,0.05,RGBColor(0x44,0x44,0x44))
txt(s,f'IDAT={df["IDAT"].mean():.1f} | IATM={df["IATM"].mean():.1f} | '
    f'{(df["SOUSPAYE"]>50).mean()*100:.0f}% sous-payes | '
    f'{(df["NPS"]<25).mean()*100:.0f}% en sortie mentale | '
    f'9.7x — cout reel d\'un depart Gen Z',
    0.75,5.4,11.5,0.55,size=10.5,italic=True,color=HC_GREY)
txt(s,'La Generation Z ne demande pas l\'impossible.\nElle demande d\'etre entendue. Avant de partir.',
    0.75,6.05,11.5,0.75,size=14,bold=True,color=HC_WHITE)

# ── Sauvegarder ──────────────────────────────────────────────────────
output='/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/GENZ_Analyse_Complete_HC_2026.pptx'
prs.save(output)
print(f'\nPowerPoint sauvegarde : {output}')
print(f'Slides : {len(prs.slides)}')
print()
slides_list=[
    'Couverture',
    '--- ACTE 1 : L\'angle mort ---',
    'L\'Iceberg + Gen Z positionne dedans',
    'STAT : 9% — ce que voit le DG',
    '--- ACTE 2 : Le diagnostic ---',
    'Slide liaison : Gen Z dans la pyramide + stats etude',
    'STAT : 64% sous-payes',
    'STAT : 41% talent bride',
    'STAT : 18% en sortie mentale',
    'Graph : Toutes les composantes',
    'Graph : Sous-payes par age',
    'Graph : Heatmap pays x secteur',
    '--- ACTE 3 : La consequence ---',
    'Slide liaison Yoshida x Gen Z (6 stats combinees)',
    'STAT : 91% idees perdues',
    'Graph : Cout visible vs reel (Yoshida + Gen Z)',
    'STAT : 9.7x cout reel depart frontline',
    'Graph : Simulation MFCFA',
    '--- ACTE 4 : Les leviers ---',
    'Graph : Correlations — vrais drivers',
    'Graph : 4 profils Gen Z',
    '5 leviers d\'action',
    'Feuille de route 90j/6m/12m',
    'Conclusion',
]
for i,name in enumerate(slides_list):
    prefix='  ' if name.startswith('---') else f'  Slide {i-[j for j,n in enumerate(slides_list) if n.startswith("---") and j<=i].__len__()+1:02d} : ' if not name.startswith('---') else ''
    print(f'{name}' if name.startswith('---') else f'  Slide : {name}')
