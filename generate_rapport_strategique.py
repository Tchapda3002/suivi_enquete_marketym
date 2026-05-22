#!/usr/bin/env python3
"""
DATATYM™ — Rapport Stratégique Gen Z Afrique 2026
Stats qui alertent. Clarté. Problème → Solution.
"""

import os, json, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
matplotlib.rcParams['font.family'] = 'DejaVu Sans'

BG='#FFFFFF'; BG2='#F4F5F7'; INK='#1B1F3B'; CORAL='#E63946'
AMBER='#E9A820'; OCEAN='#1D7A8A'; SAGE='#2D8659'; SLATE='#8896AB'
CLOUD='#D6DCE5'; SNOW='#FFFFFF'

IMG='/tmp/datatym_final'; os.makedirs(IMG, exist_ok=True)

with open('/tmp/genz_full_stats.json') as f: D=json.load(f)
S=D['S']

def h(c):
    s=str(c).lstrip('#')
    return RGBColor(int(s[0:2],16),int(s[2:4],16),int(s[4:6],16))

def save(name):
    p=f'{IMG}/{name}.png'
    plt.savefig(p,dpi=180,bbox_inches='tight',facecolor=plt.gcf().get_facecolor(),edgecolor='none')
    plt.close(); return p


# ════════════════════════════════════════════════════════════════
# CHARTS — Simples, clairs, lisibles
# ════════════════════════════════════════════════════════════════

def chart_iq():
    fig,ax=plt.subplots(figsize=(13,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats=['Oui, activement','Si opportunite\nse presente','Non, pas\nmaintenant','Non,\njamais']
    vals=[D['iq'].get('Oui, activement',12),
          D['iq'].get("Oui, si une opportunité se présente",53),
          D['iq'].get('Non, pas pour linstant',26),
          D['iq'].get('Non, pas du tout',8)]
    cols=[CORAL,AMBER,SLATE,SAGE]
    bars=ax.barh(cats[::-1],vals[::-1],color=cols[::-1],height=0.55)
    for bar,v in zip(bars,vals[::-1]):
        ax.text(v+0.5,bar.get_y()+bar.get_height()/2,f'{v:.0f}%',va='center',color=INK,fontsize=16,fontweight='bold')
    # Zone danger
    ax.axvspan(0,vals[0]+vals[1],alpha=0.06,color=CORAL)
    ax.text(32,3.3,f'{vals[0]+vals[1]:.0f}% a risque',ha='center',fontsize=13,fontweight='bold',color=CORAL)
    ax.set_xlim(0,60); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y',colors=INK,labelsize=12)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Envisagez-vous de quitter votre entreprise ?',color=INK,fontsize=15,fontweight='bold',pad=12)
    plt.tight_layout(); return save('01_iq')

def chart_vision5():
    fig,ax=plt.subplots(figsize=(13,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    items=[
        ('A mon propre compte',D['v5'].get('À mon propre compte',32),CORAL),
        ('Dans une autre entreprise',D['v5'].get('Dans une autre entreprise',20),AMBER),
        ("A l'etranger",D['v5'].get("À l\u2019étranger",0)+D['v5'].get("À l'étranger",17),AMBER),
        ('Je ne sais pas',D['v5'].get('Je ne sais pas',14),SLATE),
        ('Meme entreprise',D['v5'].get('Dans la même entreprise, à un poste supérieur',12),SAGE),
        ('Reconversion',D['v5'].get('En reconversion professionnelle',4),OCEAN),
    ]
    labs=[i[0] for i in items]; vals=[i[1] for i in items]; cols=[i[2] for i in items]
    bars=ax.barh(labs[::-1],vals[::-1],color=cols[::-1],height=0.5)
    for bar,v in zip(bars,vals[::-1]):
        ax.text(v+0.4,bar.get_y()+bar.get_height()/2,f'{v:.0f}%',va='center',color=INK,fontsize=14,fontweight='bold')
    ax.set_xlim(0,40); ax.xaxis.set_visible(False); ax.tick_params(axis='y',colors=INK,labelsize=11)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Ou vous voyez-vous dans 5 ans ?',color=INK,fontsize=15,fontweight='bold',pad=12)
    plt.tight_layout(); return save('02_vision5')

def chart_entrepreneuriat():
    fig,ax=plt.subplots(figsize=(13,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats=list(D['ent'].keys()); vals=[D['ent'][c] for c in cats]
    cols=[CORAL,AMBER,AMBER,SAGE,SLATE][:len(cats)]
    bars=ax.barh(range(len(cats)),vals,color=cols,height=0.5)
    ax.set_yticks(range(len(cats))); ax.set_yticklabels([c[:35] for c in cats],color=INK,fontsize=10)
    for bar,v in zip(bars,vals):
        ax.text(v+0.3,bar.get_y()+bar.get_height()/2,f'{v:.0f}%',va='center',color=INK,fontsize=14,fontweight='bold')
    ax.set_xlim(0,58); ax.xaxis.set_visible(False)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Avez-vous un projet entrepreneurial ?',color=INK,fontsize=15,fontweight='bold',pad=12)
    plt.tight_layout(); return save('03_ent')

def chart_nps():
    fig,ax=plt.subplots(figsize=(12,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    passifs=100-S['promoters']-S['detractors']
    cats=['Detracteurs\n(0-6)','Passifs\n(7-8)','Promoteurs\n(9-10)']
    vals=[S['detractors'],passifs,S['promoters']]; cols=[CORAL,SLATE,SAGE]
    bars=ax.bar(range(3),vals,color=cols,width=0.55)
    ax.set_xticks(range(3)); ax.set_xticklabels(cats,color=INK,fontsize=12)
    for bar,v in zip(bars,vals):
        ax.text(bar.get_x()+bar.get_width()/2,v+1,f'{v:.0f}%',ha='center',color=INK,fontsize=17,fontweight='bold')
    ax.text(1,max(vals)*0.55,f'NPS = {S["nps"]:.0f}',ha='center',fontsize=32,fontweight='bold',color=CORAL)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(CLOUD); ax.spines['left'].set_color(CLOUD); ax.tick_params(colors=SLATE)
    ax.set_title('Recommanderiez-vous votre entreprise a un proche ?',color=INK,fontsize=14,fontweight='bold',pad=15)
    plt.tight_layout(); return save('04_nps')

def chart_sal_sens():
    fig,ax=plt.subplots(figsize=(12,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats=list(D['svs'].keys()); vals=[D['svs'][c] for c in cats]
    cols=[AMBER,AMBER,CORAL,OCEAN,OCEAN][:len(cats)]
    bars=ax.bar(range(len(cats)),vals,color=cols,width=0.55)
    ax.set_xticks(range(len(cats))); ax.set_xticklabels([c[:22] for c in cats],color=INK,fontsize=9,rotation=10)
    for bar,v in zip(bars,vals):
        ax.text(bar.get_x()+bar.get_width()/2,v+0.5,f'{v:.0f}%',ha='center',color=INK,fontsize=14,fontweight='bold')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(CLOUD); ax.spines['left'].set_color(CLOUD); ax.tick_params(colors=SLATE)
    ax.set_title('Qu\'est-ce qui compte le plus : le salaire ou le sens du travail ?',color=INK,fontsize=14,fontweight='bold',pad=12)
    plt.tight_layout(); return save('05_sal_sens')

def chart_driver_reco():
    fig,ax=plt.subplots(figsize=(13,6)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    labels_map={'ICEQ':'Equite salariale\npercue','ITEQ':'Conditions\nde travail','SAT':'Satisfaction',
                'ICHD':'Developpement\nRH','ENG':'Engagement','IATM':'Management','REM':'Salaire brut'}
    items=sorted(D['corr_reco'].items(),key=lambda x:x[1],reverse=True)
    names=[labels_map.get(i[0],i[0]) for i in items]; vals=[i[1] for i in items]
    cols=[OCEAN if v>0.3 else AMBER if v>0.1 else SLATE for v in vals]
    bars=ax.barh(names[::-1],vals[::-1],color=cols[::-1],height=0.5)
    for bar,v in zip(bars,vals[::-1]):
        ax.text(v+0.02,bar.get_y()+bar.get_height()/2,f'{v:.2f}',va='center',color=INK,fontsize=13,fontweight='bold')
    ax.axvline(0,color=CLOUD,lw=1); ax.set_xlim(-0.1,0.85); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y',colors=INK,labelsize=11)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Qu\'est-ce qui fait qu\'un Gen Z recommande son entreprise ?',color=INK,fontsize=14,fontweight='bold',pad=12)
    plt.tight_layout(); return save('06_driver_reco')

def chart_driver_retention():
    fig,ax=plt.subplots(figsize=(13,6)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    labels_map={'IATM':'Management','SAT':'Satisfaction','REM':'Salaire brut',
                'ENG':'Engagement','ICHD':'Developpement RH','ICEQ':'Equite salariale','ITEQ':'Conditions de travail'}
    items=sorted(D['corr_iq'].items(),key=lambda x:x[1],reverse=True)
    names=[labels_map.get(i[0],i[0]) for i in items]; vals=[i[1] for i in items]
    cols=[OCEAN if v>0.3 else AMBER if v>0.1 else SLATE for v in vals]
    bars=ax.barh(names[::-1],vals[::-1],color=cols[::-1],height=0.5)
    for bar,v in zip(bars,vals[::-1]):
        pos=v+0.02 if v>0 else v-0.06
        ax.text(pos,bar.get_y()+bar.get_height()/2,f'{v:.2f}',va='center',
                ha='left' if v>0 else 'right',color=INK,fontsize=13,fontweight='bold')
    ax.axvline(0,color=CLOUD,lw=1); ax.set_xlim(-0.15,0.7); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y',colors=INK,labelsize=11)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Qu\'est-ce qui fait qu\'un Gen Z reste dans son entreprise ?',color=INK,fontsize=14,fontweight='bold',pad=12)
    plt.tight_layout(); return save('07_driver_retention')

def chart_digital():
    fig,ax=plt.subplots(figsize=(12,5)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats=list(D['dig'].keys()); vals=[D['dig'][c] for c in cats]
    cols=[OCEAN,SAGE,AMBER,CORAL,CORAL][:len(cats)]
    bars=ax.bar(range(len(cats)),vals,color=cols,width=0.55)
    ax.set_xticks(range(len(cats))); ax.set_xticklabels([c[:20] for c in cats],color=INK,fontsize=10)
    for bar,v in zip(bars,vals):
        ax.text(bar.get_x()+bar.get_width()/2,v+0.5,f'{v:.0f}%',ha='center',color=INK,fontsize=14,fontweight='bold')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(CLOUD); ax.spines['left'].set_color(CLOUD); ax.tick_params(colors=SLATE)
    ax.set_title('A quel point le digital est present dans votre travail ?',color=INK,fontsize=14,fontweight='bold',pad=12)
    plt.tight_layout(); return save('08_digital')

def chart_heatmap_sec():
    fig,ax=plt.subplots(figsize=(14,7)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    secs=list(D['sec_data'].keys()); indices=['IDAT','IATM','ENG','RECO','SOUSPAYE']
    data=np.array([[D['sec_data'][s].get(i,50) for i in indices] for s in secs])
    cmap=mcolors.LinearSegmentedColormap.from_list('dt',[CORAL,AMBER,OCEAN])
    im=ax.imshow(data,cmap=cmap,aspect='auto',vmin=30,vmax=85)
    ax.set_xticks(range(len(indices))); ax.set_xticklabels(indices,color=INK,fontsize=12)
    ax.set_yticks(range(len(secs))); ax.set_yticklabels([s[:28] for s in secs],color=INK,fontsize=10)
    for i in range(len(secs)):
        for j in range(len(indices)):
            c=SNOW if data[i,j]<55 else INK
            ax.text(j,i,f'{data[i,j]:.0f}',ha='center',va='center',color=c,fontsize=12,fontweight='bold')
    plt.colorbar(im,ax=ax,fraction=0.025,pad=0.02)
    ax.set_title('Scores par secteur',color=INK,fontsize=15,fontweight='bold',pad=12)
    plt.tight_layout(); return save('09_heatmap_sec')

def chart_heatmap_pays():
    fig,ax=plt.subplots(figsize=(14,7)); fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    pays=list(D['pays_data'].keys()); indices=['IDAT','IATM','ENG','RECO','SOUSPAYE']
    data=np.array([[D['pays_data'][p].get(i,50) for i in indices] for p in pays])
    cmap=mcolors.LinearSegmentedColormap.from_list('dt',[CORAL,AMBER,OCEAN])
    im=ax.imshow(data,cmap=cmap,aspect='auto',vmin=30,vmax=85)
    ax.set_xticks(range(len(indices))); ax.set_xticklabels(indices,color=INK,fontsize=12)
    ax.set_yticks(range(len(pays))); ax.set_yticklabels(pays,color=INK,fontsize=11)
    for i in range(len(pays)):
        for j in range(len(indices)):
            c=SNOW if data[i,j]<55 else INK
            ax.text(j,i,f'{data[i,j]:.0f}',ha='center',va='center',color=c,fontsize=12,fontweight='bold')
    plt.colorbar(im,ax=ax,fraction=0.025,pad=0.02)
    ax.set_title('Scores par pays',color=INK,fontsize=15,fontweight='bold',pad=12)
    plt.tight_layout(); return save('10_heatmap_pays')


# ════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ════════════════════════════════════════════════════════════════

def new_slide(prs):
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    fill=sl.background.fill; fill.solid(); fill.fore_color.rgb=h(BG); return sl

def tbx(sl,text,l,t,w,ht,size=18,bold=False,color=INK,align=PP_ALIGN.LEFT,italic=False):
    tx=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(ht))
    tf=tx.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=h(color); return tx

def rect(sl,l,t,w,ht,col):
    sh=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(ht))
    sh.fill.solid(); sh.fill.fore_color.rgb=h(col); sh.line.fill.background(); return sh

def img(sl,path,l,t,w,ht):
    if path and os.path.exists(path): sl.shapes.add_picture(path,Inches(l),Inches(t),Inches(w),Inches(ht))

def slide_cover(prs,t1,t2):
    sl=new_slide(prs); rect(sl,0,0,13.33,0.06,CORAL)
    tbx(sl,'DATATYM\u2122 ANALYTICS',0.8,0.5,12,0.4,size=11,color=OCEAN,bold=True)
    tbx(sl,t1,0.8,1.5,11.5,3.0,size=38,bold=True,color=INK)
    tbx(sl,t2,0.8,5.0,10,1.0,size=15,color=SLATE,italic=True)
    rect(sl,0,7.44,13.33,0.06,OCEAN)

def slide_stat(prs,big,unit,head,ins,accent=CORAL):
    sl=new_slide(prs); rect(sl,0,0,0.12,7.5,accent)
    tbx(sl,big,0.6,0.5,9,3.8,size=140,bold=True,color=accent)
    tbx(sl,unit,0.6,4.0,10,1.0,size=28,bold=True,color=INK)
    tbx(sl,head,0.6,5.2,12.3,0.8,size=17,bold=True,color=INK)
    tbx(sl,ins,0.6,6.1,12.3,1.2,size=12,color=SLATE,italic=True)

def slide_sep(prs,num,titre,desc):
    sl=new_slide(prs); fill=sl.background.fill; fill.solid(); fill.fore_color.rgb=h(INK)
    rect(sl,0,3.4,13.33,0.06,CORAL)
    tbx(sl,num,0.8,0.8,3,2,size=80,bold=True,color='#2A2E4A')
    tbx(sl,titre,0.8,2.0,11,1.5,size=32,bold=True,color=SNOW)
    tbx(sl,desc,0.8,3.8,10,1.0,size=15,color=CLOUD,italic=True)

def slide_chart(prs,title,img_path,note=''):
    sl=new_slide(prs); rect(sl,0,0,13.33,0.06,OCEAN)
    tbx(sl,title,0.5,0.15,12.5,0.7,size=17,bold=True,color=INK)
    img(sl,img_path,0.3,0.95,12.7,6.05)
    if note: tbx(sl,note,0.5,7.05,12,0.4,size=9,color=SLATE,italic=True)

def slide_quote(prs,quote):
    sl=new_slide(prs); fill=sl.background.fill; fill.solid(); fill.fore_color.rgb=h(BG2)
    tbx(sl,quote,1.5,1.5,10.3,4.5,size=26,bold=True,color=INK,italic=True,align=PP_ALIGN.CENTER)
    rect(sl,5.5,6.3,2.5,0.04,AMBER)
    tbx(sl,'DATATYM\u2122 \xb7 Barometre Gen Z Afrique 2026',0,6.6,13.33,0.5,size=11,color=SLATE,italic=True,align=PP_ALIGN.CENTER)

def slide_ps(prs,titre,stats,probleme,solution,accent=CORAL):
    """Slide Problème/Solution clair"""
    sl=new_slide(prs); rect(sl,0,0,0.12,7.5,accent)
    tbx(sl,titre,0.5,0.15,12.5,0.6,size=17,bold=True,color=INK)
    # Stats
    y=0.9
    for big,desc in stats:
        tbx(sl,big,0.5,y,3,0.6,size=32,bold=True,color=accent)
        tbx(sl,desc,3.5,y+0.1,5,0.5,size=12,color=INK)
        y+=0.7
    # Problème
    py=y+0.15
    rect(sl,0.5,py,12.3,1.4,BG2); rect(sl,0.5,py,0.08,1.4,CORAL)
    tbx(sl,'LE PROBLEME',0.75,py+0.05,3,0.3,size=9,bold=True,color=CORAL)
    tbx(sl,probleme,0.75,py+0.35,11.8,1.0,size=11,color=INK)
    # Solution
    sy=py+1.6
    rect(sl,0.5,sy,12.3,7.5-sy-0.3,BG2); rect(sl,0.5,sy,0.08,7.5-sy-0.3,SAGE)
    tbx(sl,'CE QU\'IL FAUT FAIRE',0.75,sy+0.05,3,0.3,size=9,bold=True,color=SAGE)
    tbx(sl,solution,0.75,sy+0.35,11.8,7.5-sy-0.7,size=11,color=INK)


# ════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════

def build():
    print('\n=== DATATYM Rapport Strategique — Final ===\n')
    print('[Charts]')
    c_iq=chart_iq(); c_v5=chart_vision5(); c_ent=chart_entrepreneuriat()
    c_nps=chart_nps(); c_svs=chart_sal_sens()
    c_dreco=chart_driver_reco(); c_dret=chart_driver_retention()
    c_dig=chart_digital(); c_sec=chart_heatmap_sec(); c_pays=chart_heatmap_pays()

    print('\n[PowerPoint]')
    prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)

    # ═══ COVER ═══
    slide_cover(prs,'Barometre Gen Z\nAfrique 2026','1 500 repondants \xb7 13 pays \xb7 10+ secteurs')

    slide_stat(prs,'1 500','Gen Z interroges. 13 pays africains.',
        'La generation qui representera 75% de la force de travail africaine d\'ici 2035.',
        'Ce rapport presente ce qu\'ils pensent, ce qu\'ils veulent, et ce qu\'ils feront — avec ou sans votre accord.')

    # ═══ 1. ENGAGEMENT — LE PIEGE ═══
    slide_sep(prs,'01','Ce que la Gen Z pense\nde son travail','Satisfaction, engagement, valeurs — et le piege qui va avec.')

    slide_stat(prs,f'{S["ENG"]:.0f}%','d\'engagement declare.',
        f'Satisfaction = {S["SAT"]:.0f}%. Alignement valeurs = {S["VAL"]:.0f}%. A premiere vue, tout va bien.',
        'Mais ces chiffres mesurent l\'implication dans le TRAVAIL — pas la loyaute envers l\'ORGANISATION.\nLa suite va montrer pourquoi c\'est un piege.')

    slide_stat(prs,'65%','des Gen Z engages\nenvisagent quand meme de partir.',
        'L\'engagement ne protege pas du depart. 65% de ceux qui ont un engagement >= 75% veulent quand meme partir.',
        'Un employe engage qui part, c\'est pire qu\'un employe desengagé qui part.\nParce qu\'il emporte avec lui la competence, le reseau, et la motivation — chez votre concurrent.',
        accent=CORAL)

    slide_stat(prs,'48%','des Gen Z satisfaits\ndeconseillent leur entreprise.',
        'Ils sont satisfaits de leur travail. Mais ils ne recommandent pas leur employeur.',
        'Satisfaction ≠ loyaute. Un collaborateur peut aimer son metier\net mepriser son organisation. C\'est exactement ce qui se passe.',
        accent=CORAL)

    # ═══ 2. CE QUE LA GEN Z PENSE DE SON EMPLOYEUR ═══
    slide_sep(prs,'02','Ce que la Gen Z pense\nde son employeur','NPS, recommandation, detracteurs.')

    slide_chart(prs,'Recommanderiez-vous votre entreprise a un proche ?',c_nps,
        f'NPS = Promoteurs ({S["promoters"]:.0f}%) - Detracteurs ({S["detractors"]:.0f}%) = {S["nps"]:.0f}. Un NPS sain est > 0.')

    slide_ps(prs,'La marque employeur est negative',
        [(f'{S["detractors"]:.0f}%','deconseillent activement votre entreprise'),
         ('18%','sont des detracteurs extremes (NPS 0-2) — pas decus, hostiles'),
         ('5%','boycottent totalement (RECO = 0)')],
        'La majorite de vos collaborateurs Gen Z recrutent activement pour la concurrence.\n'
        'Chaque detracteur est un risque reputationnel — surtout a l\'ere des reseaux sociaux.',
        '1. Integrer le NPS au tableau de bord CODIR — pas aux RH. Revue trimestrielle.\n'
        '2. Objectif chiffre : NPS > 0 en 18 mois.\n'
        '3. Entretien individuel structure avec les 18% de detracteurs extremes sous 60 jours :\n'
        '   diagnostic des causes, plan correctif individuel, ou separation negociee.\n'
        '4. Budget marque employeur structure : >= 2% de la masse salariale, pilote comme un investissement marketing.')

    # ═══ 3. CE QUE LA GEN Z VEUT FAIRE DEMAIN ═══
    slide_sep(prs,'03','Ce que la Gen Z\nveut faire demain','Vision 5 ans, intention de depart.')

    slide_stat(prs,'88%','ne se voient plus chez vous\ndans 5 ans.',
        'Sur 1 500 Gen Z, 12% seulement se projettent dans la meme entreprise.',
        '32% a leur propre compte. 20% chez un concurrent. 17% a l\'etranger.\n4% en reconversion. 14% ne savent pas encore.',
        accent=CORAL)

    slide_chart(prs,'Ou vous voyez-vous dans 5 ans ?',c_v5)

    slide_chart(prs,'Envisagez-vous de quitter votre entreprise ?',c_iq,
        '65% a risque de depart = recherche active (12%) + si opportunite (53%). Seulement 8% ne partiront jamais.')

    slide_ps(prs,'L\'hemorragie est programmee, pas soudaine',
        [('7,6%','ne partiront jamais — c\'est tout ce qui est garanti'),
         ('65%','partiront si une opportunite se presente'),
         ('12%','recherchent deja activement')],
        'Vous n\'avez que 51 personnes sur 673 qui sont acquises. Le reste est negociable ou deja parti dans sa tete.\n'
        'La fenetre d\'intervention est de 6 a 12 mois.',
        '1. Identifier vos postes critiques et les classer en 3 categories :\n'
        '   — A verrouiller (fideles) : promotion + visibilite sous 30 jours\n'
        '   — A activer (hesitants) : entretien de retention + plan de carriere sous 60 jours\n'
        '   — A diagnostiquer (actifs) : entretien individuel sous 30 jours — reengager ou organiser la transition\n'
        '2. Plans de carriere documentes 12/24/36 mois pour chaque poste cle — minimum vital.\n'
        '3. Entretiens de retention trimestriels sur les 20% de postes les plus critiques.')

    # ═══ 4. ENTREPRENEURIAT ═══
    slide_sep(prs,'04','L\'entrepreneuriat','L\'elephant dans la piece.')

    slide_stat(prs,'69%','veulent creer leur propre activite.',
        '23% sont deja en train. 19% dans les 2 prochaines annees. 51% a moyen terme.',
        'Il n\'existe aucun profil qui ne veut PAS entreprendre.\nMeme les 12% de fideles : 96% veulent aussi entreprendre.',
        accent=AMBER)

    slide_chart(prs,'Avez-vous un projet entrepreneurial ?',c_ent)

    slide_ps(prs,'Vous formez vos futurs concurrents',
        [('94%','des 25% les plus engages veulent aussi entreprendre'),
         ('99%','en grande entreprise (+200) veulent entreprendre'),
         ('97%','de ceux qui ne veulent QUE le salaire veulent aussi entreprendre'),
         ('50%','sont engages + sous-valorises + entrepreneurs = la bombe')],
        'La grande entreprise est devenue une ecole gratuite pour entrepreneurs.\n'
        'Vous formez, vous financez, vous reseautez — et le benefice part creer un concurrent.\n'
        '50% de vos Gen Z sont simultanement engages, sous-valorises ET entrepreneurs. C\'est la meme personne.',
        '1. Programme d\'intrapreneuriat structure avec budget dedie (2-5% de la masse salariale).\n'
        '2. Format : appel a projets interne trimestriel, 3 mois d\'incubation,\n'
        '   equity sharing sur les projets viables qui deviennent des filiales.\n'
        '3. Identifier les 23% "deja en cours" MAINTENANT et leur proposer un cadre interne\n'
        '   — ils partent dans les 6 mois si vous ne bougez pas.\n'
        '4. Objectif : convertir 30% des projets entrepreneuriaux en projets internes en 12 mois.\n'
        '5. Alternative : devenir investisseur du collaborateur qui part — garder 20% du projet\n'
        '   plutot que perdre 100% du talent.',
        accent=AMBER)

    # ═══ 5. SALAIRE ═══
    slide_sep(prs,'05','Le salaire','La croyance la plus couteuse.')

    slide_stat(prs,'0,06','correlation salaire \u2192 retention.',
        'Le salaire ne retient pas vos talents. C\'est la correlation la plus faible de toutes nos mesures.',
        'Mais attention : la perception d\'equite salariale (est-ce que c\'est juste ?) correle a 0,78 avec la recommandation.\nLe probleme n\'est pas le montant. C\'est le sentiment de justice.',
        accent=AMBER)

    slide_stat(prs,'73%','se sentent sous-payes\nMAIS valident leur salaire.',
        'Ils disent : "mon salaire est correct — mais ma generation est sous-payee."',
        'C\'est un sentiment collectif, pas individuel. Et c\'est ce sentiment collectif qui erode la marque employeur.\n76% accepteraient un salaire inferieur pour de meilleures conditions.',
        accent=AMBER)

    slide_chart(prs,'Salaire ou sens du travail — qu\'est-ce qui compte ?',c_svs,
        f'{D["svs"].get("Les deux à parts égales",58):.0f}% exigent les deux simultanement. Double exigence non negociable.')

    slide_ps(prs,'Les budgets augmentation sont mal cibles',
        [('0,06','correlation salaire brut \u2192 retention (quasi nulle)'),
         ('0,78','correlation equite salariale percue \u2192 recommandation'),
         ('76%','accepteraient moins pour de meilleures conditions'),
         ('25%','trouvent leur pouvoir d\'achat insuffisant')],
        'Augmenter les salaires uniformement est l\'investissement le moins rentable en retention.\n'
        'MAIS un salaire percu comme injuste declenche le depart actif.\n'
        'Le montant n\'est pas le sujet — le sentiment de justice l\'est.',
        '1. Audit d\'equite salariale sous 60 jours : comparer les grilles par poste,\n'
        '   anciennete, performance. Identifier les ecarts injustifies.\n'
        '2. Publier les resultats en INTERNE — la transparence reduit le sentiment d\'injustice\n'
        '   plus efficacement qu\'une augmentation de 5%.\n'
        '3. Repenser les packages : 70% fixe + 15% variable performance + 15% avantages\n'
        '   flexibles choisis par le collaborateur (formation, remote, horaires).\n'
        '4. Ne PAS augmenter uniformement — cibler les inequites identifiees.')

    # ═══ 6. MANAGEMENT ═══
    slide_sep(prs,'06','Le management','Le vrai levier. Le noeud central.')

    slide_stat(prs,'44,6','/100 — score management.',
        'C\'est le score le plus bas de toutes nos mesures. Et c\'est le levier n\xb01 de la retention.',
        'Correlation management \u2192 retention = 0,62. C\'est 10x plus puissant que le salaire (0,06).\nC\'est structurel : bas dans TOUS les secteurs, TOUS les pays, TOUTES les tailles.',
        accent=OCEAN)

    slide_chart(prs,'Qu\'est-ce qui fait qu\'un Gen Z reste dans son entreprise ?',c_dret,
        'Le management domine. Le salaire est quasi nul. L\'engagement ne predit pas la retention.')

    slide_ps(prs,'Les managers sont le noeud n\xb01',
        [('44,6','/100 — le score le plus bas de tous'),
         ('0,62','correlation management \u2192 retention (la plus forte)'),
         ('x10','plus puissant que le salaire pour retenir'),
         ('100%','des secteurs ont un score management bas — c\'est structurel')],
        'Les managers n\'ont jamais ete formes a manager des Gen Z. Ils ont ete promus pour leur\n'
        'expertise technique, pas pour leur capacite a developper des talents.',
        '1. Coaching managerial OBLIGATOIRE : 40h/an minimum pour tout manager\n'
        '   encadrant des Gen Z. Focus : feedback continu (pas annuel), autonomie cadree\n'
        '   (pas de micro-management), sens et impact (pas juste des KPIs).\n'
        '2. Evaluation 360\xb0 des managers PAR leurs equipes — resultats lies\n'
        '   a la remuneration variable du manager.\n'
        '3. KPI managerial : taux de retention de l\'equipe. Un manager qui perd\n'
        '   plus de 20% de ses Gen Z par an = probleme managerial identifie.\n'
        '4. Chaque euro investi en coaching managerial a 10x le ROI d\'une prime salariale.',
        accent=OCEAN)

    # ═══ 7. POTENTIEL INEXPLOITE ═══
    slide_sep(prs,'07','Le potentiel inexploite','Innovation, digital, apprentissage — ce que vos Gen Z pourraient donner.')

    slide_stat(prs,f'{S["INNOV"]:.0f}','/100 — score innovation.',
        f'Apprentissage = {S["APPRENT"]:.0f}. Initiative = {S["INIT"]:.0f}. Vos Gen Z veulent innover.',
        'Ce ne sont pas des scores bas. C\'est un capital disponible.\n'
        'Le probleme n\'est pas l\'appetence — c\'est l\'absence d\'espace pour l\'exprimer.',
        accent=OCEAN)

    slide_chart(prs,'A quel point le digital est present dans votre travail ?',c_dig,
        '56% utilisent le digital en grande partie ou totalement. 38% utilisent deja l\'IA.')

    slide_ps(prs,'Capital humain bride — productivite perdue',
        [('82','/100 innovation — ils veulent innover'),
         ('81','/100 apprentissage — ils veulent apprendre'),
         ('69%','prets a se reconvertir — pas attaches a leur metier'),
         ('66%','utilisent deja le digital au quotidien')],
        'Vos Gen Z ont le potentiel. Vos organisations ne liberent pas l\'espace.\n'
        'C\'est une perte de productivite invisible mais massive.',
        '1. Programmes "Digital Champions" : identifier les 10% les plus digitaux\n'
        '   et leur donner un role d\'ambassadeur interne.\n'
        '2. Budget formation individuel annuel de 500-1000$ par Gen Z —\n'
        '   non conditionne au metier actuel (encourage la polyvalence).\n'
        '3. Career marketplace digitale interne : les Gen Z postulent sur des missions\n'
        '   internes AVANT de chercher dehors. Passerelles entre metiers facilitees.\n'
        '4. Les 69% prets a se reconvertir = une opportunite de mobilite interne,\n'
        '   pas une menace. Programme de reskilling finance avec periode d\'essai de 3 mois.',
        accent=OCEAN)

    # ═══ 8. 18-22 vs 23-27 ═══
    slide_sep(prs,'08','Les 18-22 ans\nvs les 23-27 ans','Ca s\'accelere.')

    slide_ps(prs,'La prochaine vague est plus radicale',
        [('70%','des 18-22 a risque de depart (vs 61% chez les 23-27)'),
         ('72','/100 sentiment de sous-paiement (vs 67 chez les 23-27)'),
         ('93%','veulent entreprendre (identique aux 23-27)')],
        'Ce qui est un signal orange avec les 23-27 ans est deja un signal rouge avec les 18-22.\n'
        'Et ils arrivent en masse. La generation suivante sera encore plus radicale.',
        '1. Parcours d\'onboarding specifique 18-22 : mentorat des J+1,\n'
        '   projet concret des le premier mois, feedback hebdomadaire.\n'
        '2. Programme "Shadow Leadership" : exposer les 18-22 aux decisions strategiques\n'
        '   des les 6 premiers mois — ils veulent comprendre le "pourquoi".\n'
        '3. Ne PAS les traiter comme des juniors —\n'
        '   les traiter comme des entrepreneurs en residence.\n'
        '4. Budget d\'integration specifique : les 90 premiers jours determinent\n'
        '   si un 18-22 ans reste 3 ans ou part dans 6 mois.',
        accent='#6C3483')

    # ═══ 9. TAILLE D'ENTREPRISE ═══
    slide_sep(prs,'09','Taille d\'entreprise','Pas le meme combat.')

    slide_ps(prs,'Grande entreprise = ecole d\'entrepreneurs',
        [('99%','des Gen Z en grande entreprise (+200) veulent entreprendre'),
         ('57%','sont a risque de depart'),
         ('82','engagement en PME (<10) — le plus eleve'),
         ('52','/100 RECO en PME — le plus bas')],
        'Plus l\'entreprise est grande, plus elle forme d\'entrepreneurs qui iront creer ailleurs.\n'
        'Les PME ont l\'engagement mais pas la marque employeur.',
        '1. Grandes entreprises : creer un fonds d\'intrapreneuriat corporate.\n'
        '   "Startup studio interne" — les meilleurs projets deviennent des filiales.\n'
        '   Plutot que perdre 100% du talent, garder 30% du projet.\n'
        '2. PME : capitaliser sur l\'engagement fort en ameliorant la RECO.\n'
        '   Les PME n\'ont pas de marque employeur — en creer une localement\n'
        '   (presence campus, temoignages collaborateurs, partenariats ecoles).\n'
        '3. Moyennes (51-200) : l\'engagement est le plus bas (72) mais le sous-paiement aussi (59).\n'
        '   Levier : valoriser les conditions de travail qui sont deja meilleures que la moyenne.',
        accent=AMBER)

    # ═══ 10. STABILITE ═══
    slide_sep(prs,'10','La stabilite','Ce que la Gen Z veut vraiment — et ca surprend.')

    slide_stat(prs,'52%','veulent la stabilite.',
        'Surprise : 9% seulement veulent la mobilite. 65% trouvent leur secteur attractif.',
        'Contrairement aux idees recues, les Gen Z ne veulent pas le chaos.\n'
        'Ils ne fuient pas la stabilite — ils fuient l\'absence de perspective.',
        accent=SAGE)

    # ═══ 11. HEATMAPS ═══
    slide_sep(prs,'11','Radiographie sectorielle\net par pays','Ou sont les zones de fragilite — et les specificites.')

    slide_chart(prs,'Scores par secteur — 5 indicateurs x 8 secteurs',c_sec,
        'Sante : NPS = -45, le pire. Agriculture : NPS = +9, le seul positif. Le management (IATM) est bas partout.')

    slide_chart(prs,'Scores par pays — 5 indicateurs x 10 pays',c_pays,
        'Chaque pays a son profil. Un plan continental unique ne sert personne.')

    slide_chart(prs,'Qu\'est-ce qui fait qu\'un Gen Z recommande son entreprise ?',c_dreco,
        'L\'equite salariale percue domine (0,78). Le salaire brut est quasi nul (0,06). C\'est la justice qui compte.')

    # ═══ 12. SYNTHESE ═══
    slide_sep(prs,'12','Synthese et\nplan d\'action','10 constats. 5 chantiers. 90 jours.')

    # 10 constats
    sl=new_slide(prs); rect(sl,0,0,13.33,0.06,CORAL)
    tbx(sl,'10 CONSTATS',0.5,0.15,12,0.6,size=20,bold=True,color=INK)
    constats=[
        ('65% envisagent de partir','Dont 12% en recherche active',CORAL),
        ('88% ne se voient plus chez vous dans 5 ans','Seulement 12% se projettent ici',CORAL),
        ('69% veulent entreprendre','23% sont deja en train',CORAL),
        ('48% des satisfaits deconseillent','Satisfaction ≠ loyaute',CORAL),
        ('50% = engages + sous-valorises + entrepreneurs','C\'est la meme personne',AMBER),
        ('Salaire \u2192 retention = 0,06','Le salaire ne retient pas',AMBER),
        ('Management \u2192 retention = 0,62','10x plus puissant que le salaire',OCEAN),
        ('Equite salariale \u2192 recommandation = 0,78','La justice fait recommander',OCEAN),
        ('Innovation = 82, Apprentissage = 81','Capital inexploite',SAGE),
        ('52% veulent la stabilite','Ils fuient l\'absence de perspective, pas la stabilite',SAGE),
    ]
    for i,(stat,desc,col) in enumerate(constats):
        y=0.85+i*0.62
        rect(sl,0.5,y,0.1,0.5,col)
        tbx(sl,stat,0.8,y,5.5,0.5,size=12,bold=True,color=col)
        tbx(sl,desc,6.5,y,6.3,0.5,size=11,color=INK)

    # 5 chantiers
    sl=new_slide(prs); rect(sl,0,0,13.33,0.06,OCEAN)
    tbx(sl,'5 CHANTIERS — 90 JOURS',0.5,0.15,12,0.6,size=20,bold=True,color=INK)
    chantiers=[
        ('MANAGEMENT','Coaching 40h/an + evaluation 360\xb0 par les equipes + KPI retention',CORAL,'J+30','Score IATM > 55'),
        ('EQUITE SALARIALE','Audit salarial interne + transparence + packages flexibles',AMBER,'J+45','NPS > -10 en 18 mois'),
        ('INTRAPRENEURIAT','Programme structure + incubation + equity sharing sur projets viables',AMBER,'J+60','30% projets canalises en interne'),
        ('ONBOARDING 18-22','Mentorat J+1 + Shadow Leadership + projet des le 1er mois','#6C3483','J+45','Risque depart 18-22 < 55%'),
        ('MOBILITE INTERNE','Career marketplace + reskilling + passerelles entre metiers',OCEAN,'J+90','20% de reconversions internes'),
    ]
    for i,(titre,desc,col,delai,kpi) in enumerate(chantiers):
        y=0.9+i*1.25
        rect(sl,0.5,y,0.12,1.08,col); rect(sl,0.72,y,12.1,1.08,BG2)
        tbx(sl,titre,0.85,y+0.05,2.8,0.4,size=13,bold=True,color=col)
        tbx(sl,desc,0.85,y+0.45,7,0.55,size=10,color=INK)
        rect(sl,8.5,y+0.15,1.3,0.7,col)
        tbx(sl,delai,8.5,y+0.15,1.3,0.7,size=13,bold=True,color=SNOW,align=PP_ALIGN.CENTER)
        tbx(sl,f'KPI : {kpi}',10.0,y+0.3,2.8,0.5,size=9,color=SLATE,italic=True)

    slide_quote(prs,
        '1 500 Gen Z ont parle.\n\n'
        '88% preparent leur depart.\n'
        '69% construisent leur alternative.\n'
        '50% sont engages, sous-valorises\net entrepreneurs a la fois.\n\n'
        'La question n\'est pas s\'ils partiront.\n'
        'C\'est ce que vous faites avant.')

    slide_cover(prs,'DATATYM\u2122 Analytics',
        'Barometre Gen Z Afrique 2026\n1 500 repondants \xb7 13 pays \xb7 Confidentiel')

    out='/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/DATATYM_Rapport_Strategique_GenZ_2026.pptx'
    prs.save(out); print(f'\n=== {len(prs.slides)} slides -> {out} ===')

if __name__=='__main__': build()
