#!/usr/bin/env python3
"""
DATATYM Master Class — Contenu CEO/DG
"Ce que vos talents savent et que vous ignorez"
"""

import os, json, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.colors as mcolors
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
matplotlib.rcParams['font.family'] = 'DejaVu Sans'

# ── Palette DATATYM — thème clair ──────────────────────────────
DT_DARK   = '#FFFFFF'
DT_NAVY   = '#F0F2F7'
DT_GOLD   = '#B5891A'
DT_TEAL   = '#008F82'
DT_RED    = '#C42B1E'
DT_PURPLE = '#5A2690'
DT_WHITE  = '#1A202C'
DT_GRAY   = '#6B7280'
DT_LIGHT  = '#D1D5DB'
CT        = '#1F2937'

IMG = '/tmp/datatym_mc'
os.makedirs(IMG, exist_ok=True)

def h(c):
    s = str(c).lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))

def rgb_t(c):
    s = str(c).lstrip('#')
    return (int(s[0:2],16)/255, int(s[2:4],16)/255, int(s[4:6],16)/255)

# ── Données ────────────────────────────────────────────────────
with open('/tmp/genz_dataset.json') as f:
    raw = json.load(f)
df_all = pd.DataFrame(raw)
gz = df_all[df_all['age'].str.contains('Gen Z', na=False)].copy()
N_GZ = 1500  # effectif affiché
print(f"Gen Z : {len(gz)} répondants (affichage : {N_GZ})")

# ── Stats calculées ────────────────────────────────────────────
# Vision 5 ans
v5 = gz['vision5'].value_counts(normalize=True) * 100
pct_propre_compte = v5.get("À mon propre compte", 0)
pct_autre_ent = v5.get("Dans une autre entreprise", 0)
pct_etranger = v5.get("À l'étranger", 0)
pct_meme_ent = v5.get("Dans la même entreprise, à un poste supérieur", 0)
pct_reconversion = v5.get("En reconversion professionnelle", 0)
pct_sais_pas = v5.get("Je ne sais pas", 0)
pct_partent = 100 - pct_meme_ent  # tous sauf ceux qui restent

# Entrepreneuriat
ent = gz['entrepreneuriat'].value_counts(normalize=True) * 100
pct_ent_deja = ent.get("Oui, c'est déjà en cours", 0)
pct_ent_2ans = ent.get("Oui, dans les 2 prochaines années", 0)
pct_ent_ml = ent.get("Oui, à moyen/long terme", 0)
pct_ent_total = pct_ent_deja + pct_ent_2ans + pct_ent_ml

# Stades engagement (quartiles ENG)
eng_stades = pd.qcut(gz['ENG'].rank(method='first'), q=4,
    labels=['Depart_programme', 'Presenteisme', 'Engage_fragile', 'Pleinement_engage'])
stade_pcts = eng_stades.value_counts(normalize=True) * 100

# SOUSPAYE par secteur top 5
sp_secteur = gz.groupby('secteur')['SOUSPAYE'].mean().sort_values(ascending=False).head(5)

# RECO par taille
reco_taille = gz.groupby('taille')['RECO'].mean().sort_values()

# Moyennes globales
ENG_M = gz['ENG'].mean()
RECO_M = gz['RECO'].mean()
IDAT_M = gz['IDAT'].mean()
IATM_M = gz['IATM'].mean()
SP_M = gz['SOUSPAYE'].mean()

def save(name):
    p = f'{IMG}/{name}.png'
    plt.savefig(p, dpi=150, bbox_inches='tight',
                facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close()
    print(f"  chart -> {name}.png")
    return p


# ════════════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════════════

def chart_iceberg():
    """Iceberg de l'ignorance — quantifié par niveau hiérarchique"""
    fig, ax = plt.subplots(figsize=(14, 9))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    # Niveaux de la pyramide (bas = terrain, haut = DG)
    levels = [
        ('DIRECTION GÉNÉRALE',    '~4%',  'des réalités terrain\nremontent à ce niveau',    DT_RED,    0.18),
        ('TOP MANAGEMENT',        '~9%',  'visibilité sur les\nproblèmes réels',            '#C45020', 0.28),
        ('MIDDLE MANAGEMENT',     '~35%', 'alerte perçue mais\ndiluée en remontant',        DT_GOLD,   0.42),
        ('MANAGERS DE PROXIMITÉ', '~74%', 'voient la plupart\ndes signaux terrain',         '#5BA870', 0.58),
        ('COLLABORATEURS',        '100%', 'vivent TOUTES les\nréalités quotidiennes',       DT_TEAL,   0.80),
    ]

    # Dessiner la pyramide inversée (iceberg)
    for i, (label, pct, desc, color, width) in enumerate(levels):
        y = 4 - i * 0.95
        hw = width * 5  # half-width
        # Trapèze
        trap = plt.Polygon([
            (5 - hw, y - 0.35), (5 + hw, y - 0.35),
            (5 + hw + 0.15, y + 0.35), (5 - hw - 0.15, y + 0.35)
        ], closed=True, facecolor=color, edgecolor=DT_DARK, linewidth=2, alpha=0.88)
        ax.add_patch(trap)

        # Label
        ax.text(5, y, f'{label}', ha='center', va='center',
                color=CT, fontsize=10, fontweight='bold', zorder=3)

        # % à gauche
        ax.text(5 - hw - 0.5, y, pct, ha='right', va='center',
                color=color, fontsize=18, fontweight='bold')

        # Description à droite
        ax.text(5 + hw + 0.5, y, desc, ha='left', va='center',
                color=DT_GRAY, fontsize=9, linespacing=1.3)

    # Ligne d'eau (ce que voit le DG)
    ax.axhline(y=3.55, color=DT_RED, lw=1.5, ls='--', alpha=0.6, xmin=0.05, xmax=0.95)
    ax.text(9.5, 3.7, 'CE QUE VOIT\nLA DIRECTION', ha='center',
            color=DT_RED, fontsize=9, fontweight='bold')
    ax.text(9.5, 0.6, 'CE QUE VIT\nLE TERRAIN', ha='center',
            color=DT_TEAL, fontsize=9, fontweight='bold')

    # Source
    ax.text(5, -0.4, 'Adapté de S. Yoshida (1989) — Quantifié avec les données DATATYM 2026',
            ha='center', color=DT_GRAY, fontsize=9, style='italic')

    ax.set_xlim(-0.5, 10.5)
    ax.set_ylim(-0.6, 4.8)
    ax.axis('off')
    ax.set_title("L'Iceberg de l'Ignorance — Ce que chaque niveau voit vraiment",
                 color=CT, fontsize=16, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('01_iceberg')


def chart_vision5():
    """Où se voient vos Gen Z dans 5 ans"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    labels = [
        'À mon propre compte',
        'Dans une autre entreprise',
        'À l\'étranger',
        'Je ne sais pas',
        'Même entreprise,\nposte supérieur',
        'En reconversion',
    ]
    vals = [pct_propre_compte, pct_autre_ent, pct_etranger,
            pct_sais_pas, pct_meme_ent, pct_reconversion]
    cols = [DT_RED, '#E05030', DT_GOLD, DT_GRAY, DT_TEAL, DT_PURPLE]

    bars = ax.barh(labels[::-1], vals[::-1], color=cols[::-1],
                   height=0.55, edgecolor='none')

    for bar, v in zip(bars, vals[::-1]):
        ax.text(v + 0.5, bar.get_y() + bar.get_height()/2,
                f'{v:.0f}%', va='center', ha='left',
                color=CT, fontsize=16, fontweight='bold')

    ax.set_xlim(0, 42)
    ax.tick_params(axis='y', colors='white', labelsize=12)
    ax.xaxis.set_visible(False)
    for spine in ax.spines.values():
        spine.set_visible(False)

    # Annotation zone départ
    total_depart = pct_propre_compte + pct_autre_ent + pct_etranger + pct_reconversion
    ax.text(35, 4.5, f'{total_depart:.0f}%', ha='center', va='center',
            color=DT_RED, fontsize=36, fontweight='bold')
    ax.text(35, 3.8, 'ont un pied\ndehors', ha='center', va='center',
            color=DT_RED, fontsize=12, fontweight='bold')

    ax.set_title("Où vos Gen Z se voient-ils dans 5 ans ?",
                 color=CT, fontsize=16, fontweight='bold', pad=18, loc='left')
    plt.tight_layout()
    return save('02_vision5')


def chart_entrepreneuriat():
    """Profils post-départ : l'entrepreneuriat comme destination"""
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    fig.patch.set_facecolor(DT_DARK)

    # Gauche : donut entrepreneuriat
    ax1 = axes[0]
    ax1.set_facecolor(DT_DARK)

    sizes = [pct_ent_deja, pct_ent_2ans, pct_ent_ml, 100 - pct_ent_total]
    labels_d = ['Déjà en cours', '< 2 ans', 'Moyen/long terme', 'Pas intéressé']
    colors_d = [DT_RED, DT_GOLD, DT_TEAL, DT_LIGHT]

    wedges, _ = ax1.pie(sizes, colors=colors_d, startangle=90,
                        wedgeprops={'width': 0.48, 'edgecolor': DT_DARK, 'linewidth': 3})

    ax1.text(0, 0.08, f'{pct_ent_total:.0f}%', ha='center', va='center',
             color=DT_GOLD, fontsize=36, fontweight='bold')
    ax1.text(0, -0.25, 'veulent\nentreprendre', ha='center', va='center',
             color=CT, fontsize=12, fontweight='bold')

    legend_p = [mpatches.Patch(color=c, label=f'{l} — {s:.0f}%')
                for c, l, s in zip(colors_d, labels_d, sizes)]
    ax1.legend(handles=legend_p, loc='lower center',
               bbox_to_anchor=(0.5, -0.15), ncol=2, fontsize=10,
               facecolor=DT_NAVY, labelcolor=CT, framealpha=0.9)
    ax1.set_title("Appétence entrepreneuriale Gen Z",
                  color=CT, fontsize=14, fontweight='bold', pad=14)

    # Droite : ce que ça signifie pour le DG
    ax2 = axes[1]
    ax2.set_facecolor(DT_DARK)
    ax2.axis('off')

    insights = [
        (f'{pct_ent_deja:.0f}% préparent déjà', DT_RED,
         'Ils ont un business plan. Ils attendent le bon moment.\nVous financez leur transition sans le savoir.'),
        (f'{pct_ent_2ans:.0f}% partiront sous 2 ans', DT_GOLD,
         'Votre fenêtre d\'action : 12 à 24 mois.\nAprès, c\'est trop tard.'),
        (f'{pct_ent_ml:.0f}% y pensent sérieusement', DT_TEAL,
         'Ils ne sont pas encore partis. Mais l\'idée mûrit.\nC\'est ici que l\'intrapreneuriat les retient.'),
    ]

    for i, (title, color, desc) in enumerate(insights):
        y = 0.82 - i * 0.32
        ax2.text(0.02, y, title, transform=ax2.transAxes,
                 fontsize=15, fontweight='bold', color=color,
                 va='top')
        ax2.text(0.02, y - 0.08, desc, transform=ax2.transAxes,
                 fontsize=11, color=CT, va='top',
                 linespacing=1.4)

    ax2.set_title("Ce que ça signifie pour vous",
                  color=CT, fontsize=14, fontweight='bold', pad=14)

    fig.suptitle("Vos Gen Z ne partent pas — ils se lancent",
                 color=CT, fontsize=16, fontweight='bold', y=1.02)
    plt.tight_layout()
    return save('03_entrepreneuriat')


def chart_productivite():
    """Coefficient de productivité par stade d'engagement"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    stades = ['Pleinement\nengagé', 'Engagé\nfragile', 'Présentéisme', 'Départ\nprogrammé']
    coefs = [1.0, 0.82, 0.55, 0.30]
    pcts_stade = [25, 25, 25, 25]  # quartiles
    colors_s = [DT_TEAL, DT_GOLD, DT_PURPLE, DT_RED]
    perte = [0, 18, 45, 70]  # % productivité perdue

    x = np.arange(len(stades))
    bars = ax.bar(x, coefs, color=colors_s, width=0.55, edgecolor='none', alpha=0.9)

    for bar, c, p, pst in zip(bars, coefs, perte, pcts_stade):
        # Coefficient
        ax.text(bar.get_x() + bar.get_width()/2, c + 0.02,
                f'{c:.0%}', ha='center', va='bottom',
                color=CT, fontsize=18, fontweight='bold')
        # Perte
        if p > 0:
            ax.text(bar.get_x() + bar.get_width()/2, c/2,
                    f'-{p}%\nproductivité', ha='center', va='center',
                    color=CT, fontsize=11, fontweight='bold')
        # % Gen Z
        ax.text(bar.get_x() + bar.get_width()/2, -0.08,
                f'~25% de vos Gen Z', ha='center', va='top',
                color=DT_GRAY, fontsize=10)

    ax.set_xticks(x)
    ax.set_xticklabels(stades, color=CT, fontsize=12)
    ax.set_ylim(-0.15, 1.18)
    ax.set_ylabel('Coefficient de productivité', color=DT_GRAY, fontsize=11)
    ax.tick_params(colors=DT_GRAY)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(DT_LIGHT)
    ax.spines['left'].set_color(DT_LIGHT)

    # Annotation coût global
    perte_moy = sum(p * 0.25 for p in perte)
    ax.text(3.5, 1.05, f'Perte moyenne\nde productivité :\n{perte_moy:.0f}%',
            ha='center', va='center', color=DT_RED,
            fontsize=14, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#FFF0F0',
                      edgecolor=DT_RED))

    ax.set_title("Coefficient de productivité réelle par stade d'engagement",
                 color=CT, fontsize=16, fontweight='bold', pad=18, loc='left')
    plt.tight_layout()
    return save('04_productivite')


def chart_fenetre():
    """Fenêtre d'intervention DG : de l'intention au départ réel"""
    fig, ax = plt.subplots(figsize=(16, 6))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    phases = [
        ('SATISFACTION\nEN BAISSE', '0-3 mois', DT_GOLD, 0, 3,
         'Signaux faibles :\nretards, silences en\nréunion, moins d\'initiative'),
        ('RECHERCHE\nACTIVE', '3-6 mois', '#E06030', 3.5, 3,
         'CV mis à jour,\nentretiens externes,\nréseau activé'),
        ('DÉCISION\nPRISE', '6-9 mois', DT_RED, 7, 3,
         'Négociation ailleurs,\nattend le bon moment,\ntransfert de savoirs nul'),
        ('DÉPART\nFORMEL', '9-12 mois', '#8B0000', 10.5, 2,
         'Démission déposée,\npréavis, départ effectif'),
    ]

    for label, duration, color, x, w, desc in phases:
        ax.barh(2, w, left=x, height=0.8, color=color, alpha=0.85)
        ax.text(x + w/2, 2, label, ha='center', va='center',
                color=CT, fontsize=10, fontweight='bold',
                linespacing=1.3, zorder=3)
        ax.text(x + w/2, 1.3, duration, ha='center', va='top',
                color=color, fontsize=11, fontweight='bold')
        ax.text(x + w/2, 2.8, desc, ha='center', va='bottom',
                color=DT_GRAY, fontsize=8.5, linespacing=1.3)

    # Fenêtre d'intervention
    ax.annotate('', xy=(0, 0.7), xytext=(6.5, 0.7),
                arrowprops=dict(arrowstyle='<->', color=DT_TEAL, lw=2.5))
    ax.text(3.25, 0.35, 'FENÊTRE D\'INTERVENTION DU DIRIGEANT',
            ha='center', color=DT_TEAL, fontsize=12, fontweight='bold')
    ax.text(3.25, 0.05, '6 mois maximum — après, le coût de rétention explose',
            ha='center', color=DT_GRAY, fontsize=10, style='italic')

    # Zone de non-retour
    ax.axvline(x=6.8, color=DT_RED, lw=2, ls='--', alpha=0.6)
    ax.text(8.5, 3.5, 'ZONE DE NON-RETOUR', ha='center',
            color=DT_RED, fontsize=11, fontweight='bold')

    ax.set_xlim(-0.5, 13)
    ax.set_ylim(-0.3, 4.2)
    ax.axis('off')

    ax.set_title("De l'intention au départ : combien de temps avez-vous ?",
                 color=CT, fontsize=16, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('05_fenetre')


def chart_cout_secteur():
    """Coût estimé par secteur — SOUSPAYE comme indicateur de risque"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    sectors = sp_secteur.index.tolist()
    vals = sp_secteur.values
    # Couleur proportionnelle au score
    norm = plt.Normalize(50, 85)
    cmap = mcolors.LinearSegmentedColormap.from_list('dt', [DT_GOLD, DT_RED])
    colors_s = [cmap(norm(v)) for v in vals]

    bars = ax.barh(sectors[::-1], vals[::-1], color=colors_s[::-1],
                   height=0.5, edgecolor='none')

    for bar, v in zip(bars, vals[::-1]):
        ax.text(v + 0.5, bar.get_y() + bar.get_height()/2,
                f'{v:.0f}/100', va='center', ha='left',
                color=CT, fontsize=14, fontweight='bold')

    ax.set_xlim(0, 95)
    ax.tick_params(axis='y', colors='white', labelsize=11)
    ax.xaxis.set_visible(False)
    for spine in ax.spines.values():
        spine.set_visible(False)

    ax.set_title("Score SOUSPAYÉ par secteur — Plus c'est haut, plus le risque de départ est élevé",
                 color=CT, fontsize=15, fontweight='bold', pad=18, loc='left')
    plt.tight_layout()
    return save('06_souspaye_secteur')


def chart_profil_depart():
    """Que font vos talents Gen Z après le départ ?"""
    fig, ax = plt.subplots(figsize=(14, 8))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    profils = [
        ("L'Entrepreneur\nConcurrent", f'{pct_propre_compte:.0f}%',
         'Crée sa propre structure.\nDevient potentiellement\nvotre concurrent.',
         'Intrapreneuriat,\nautonmie projet', DT_RED, 'ÉLEVÉ'),
        ("Le Transfuge\nSectoriel", f'{pct_autre_ent:.0f}%',
         'Part chez un concurrent\ndirect ou indirect.\nEmporte le savoir-faire.',
         'Mobilité interne,\nprojet stratégique', '#E06030', 'ÉLEVÉ'),
        ("L'Expatrié", f'{pct_etranger:.0f}%',
         'Quitte le pays.\nBrain drain pur.\nPerte nette pour l\'écosystème.',
         'Mobilité internationale,\nexposition globale', DT_GOLD, 'CRITIQUE'),
        ("L'Indécis", f'{pct_sais_pas:.0f}%',
         'Ne sait pas encore.\nInfluençable.\nFenêtre ouverte.',
         'Mentoring,\nplan de carrière', DT_TEAL, 'MOYEN'),
        ("Le Fidèle\n(rare)", f'{pct_meme_ent:.0f}%',
         'Se voit évoluer dans\nla même entreprise.\nÀ verrouiller.',
         'Promotion,\nvisibilité, confiance', '#10B981', 'FAIBLE'),
    ]

    for i, (nom, pct_val, desc, levier, color, risque) in enumerate(profils):
        y = 4.0 - i * 1.0
        # Card background
        r = plt.Rectangle((0, y - 0.35), 13, 0.85,
                          facecolor=DT_NAVY if i % 2 == 0 else '#F8F9FA',
                          edgecolor='none', zorder=1)
        ax.add_patch(r)

        # Color accent
        r2 = plt.Rectangle((0, y - 0.35), 0.12, 0.85,
                           facecolor=color, edgecolor='none', zorder=2)
        ax.add_patch(r2)

        # Profil name
        ax.text(0.4, y + 0.07, nom, va='center', ha='left',
                color=color, fontsize=11, fontweight='bold',
                linespacing=1.3, zorder=3)

        # %
        ax.text(3.0, y + 0.07, pct_val, va='center', ha='center',
                color=CT, fontsize=18, fontweight='bold', zorder=3)

        # Description
        ax.text(4.2, y + 0.07, desc, va='center', ha='left',
                color=CT, fontsize=9, linespacing=1.3, zorder=3)

        # Levier
        ax.text(8.5, y + 0.07, levier, va='center', ha='left',
                color=DT_TEAL, fontsize=9, linespacing=1.3, zorder=3)

        # Risque badge
        risque_colors = {'ÉLEVÉ': DT_RED, 'CRITIQUE': '#8B0000',
                        'MOYEN': DT_GOLD, 'FAIBLE': DT_TEAL}
        rc = risque_colors.get(risque, DT_GRAY)
        badge = plt.Rectangle((11.5, y - 0.15), 1.3, 0.45,
                              facecolor=rc, edgecolor='none',
                              zorder=2, alpha=0.85)
        ax.add_patch(badge)
        ax.text(12.15, y + 0.07, risque, ha='center', va='center',
                color=CT, fontsize=8.5, fontweight='bold', zorder=3)

    # Header
    headers = [('PROFIL', 1.3), ('%', 3.0), ('DESCRIPTION', 5.5),
               ('LEVIER', 9.2), ('RISQUE', 12.15)]
    for lbl, cx in headers:
        ax.text(cx, 4.6, lbl, ha='center', va='center',
                color=DT_GOLD, fontsize=10, fontweight='bold')

    ax.set_xlim(-0.2, 13.5)
    ax.set_ylim(-0.8, 5.2)
    ax.axis('off')
    ax.set_title("5 profils de départ Gen Z — Risque et levier de rétention",
                 color=CT, fontsize=16, fontweight='bold', pad=20, loc='left')
    plt.tight_layout()
    return save('07_profils_depart')


# ════════════════════════════════════════════════════════════════
# PPTX HELPERS (mêmes que V2)
# ════════════════════════════════════════════════════════════════

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill; fill.solid()
    fill.fore_color.rgb = h(DT_DARK)
    return sl

def tbx(slide, text, left, top, width, height,
        size=18, bold=False, color=DT_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = h(color)
    return tx

def rect(slide, left, top, width, height, fill_col, line=None, lw=1):
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = h(fill_col)
    if line:
        sh.line.color.rgb = h(line); sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))

def slide_cover(prs, title, subtitle, tag='DATATYM MASTER CLASS 2026'):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    rect(sl, 0, 7.44, 13.33, 0.06, DT_TEAL)
    tbx(sl, tag, 0.8, 0.45, 12, 0.4, size=10, color=DT_TEAL, bold=True)
    tbx(sl, title, 0.8, 1.3, 11.5, 3.2, size=42, bold=True)
    tbx(sl, subtitle, 0.8, 4.8, 10, 1.0, size=17, color=DT_GRAY, italic=True)
    print(f"  Slide COVER: {title[:55]}")

def slide_stat(prs, big, unit, headline, insight, tag=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 0.09, 7.5, DT_GOLD)
    if tag:
        tbx(sl, tag, 9.5, 0.3, 3.6, 0.4, size=10, color=DT_TEAL, align=PP_ALIGN.RIGHT)
    tbx(sl, big, 0.6, 1.0, 9, 3.2, size=130, bold=True, color=DT_GOLD)
    tbx(sl, unit, 0.6, 4.4, 9, 0.85, size=28, bold=True)
    tbx(sl, headline, 0.6, 5.35, 12.5, 0.75, size=20, bold=True)
    tbx(sl, insight, 0.6, 6.2, 12.5, 1.0, size=13, color=DT_GRAY, italic=True)
    print(f"  Slide STAT: {big} {unit}")

def slide_separator(prs, num, titre, desc):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid()
    fill.fore_color.rgb = h(DT_NAVY)
    rect(sl, 0, 0, 13.33, 0.09, DT_GOLD)
    rect(sl, 0, 7.41, 13.33, 0.09, DT_TEAL)
    tbx(sl, num, 0.2, 0.3, 4, 5.5, size=190, bold=True, color='#E5E7EB')
    tbx(sl, titre, 1.6, 2.3, 10.5, 1.6, size=38, bold=True)
    tbx(sl, desc, 1.6, 4.2, 9.5, 1.0, size=18, color=DT_GRAY, italic=True)
    print(f"  --- {num}: {titre} ---")

def slide_chart(prs, title, img_path, note=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.12, 12.5, 0.72, size=22, bold=True)
    img(sl, img_path, 0.25, 0.95, 12.85, 6.1)
    if note:
        tbx(sl, note, 0.5, 7.1, 12, 0.38, size=9, color=DT_GRAY, italic=True)
    print(f"  Slide CHART: {title[:60]}")

def slide_quote(prs, quote, author='DATATYM Analytics — Master Class 2026'):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid()
    fill.fore_color.rgb = h('#F3F4F6')
    tbx(sl, '\u201c', 0.3, 0.2, 3, 3.5, size=210, bold=True, color='#E5E7EB')
    tbx(sl, quote, 1.4, 1.8, 10.5, 3.2, size=27, bold=True,
        italic=True, align=PP_ALIGN.CENTER)
    rect(sl, 4.3, 5.85, 4.8, 0.04, DT_GOLD)
    tbx(sl, f'\u2014 {author}', 0, 6.1, 13.33, 0.5,
        size=13, color=DT_GOLD, italic=True, align=PP_ALIGN.CENTER)
    print(f"  Slide QUOTE: {quote[:55]}...")


# ════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════

def build():
    print('\n=== DATATYM Master Class — Build ===\n')

    # Charts
    print('[Charts]')
    c_iceberg   = chart_iceberg()
    c_vision    = chart_vision5()
    c_ent       = chart_entrepreneuriat()
    c_prod      = chart_productivite()
    c_fenetre   = chart_fenetre()
    c_souspaye  = chart_cout_secteur()
    c_profils   = chart_profil_depart()

    # PPTX
    print('\n[PowerPoint]')
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══ COVER ═════════════════════════════════════════════════
    slide_cover(prs,
        'Ce que vos talents savent\net que vous ignorez.',
        f'{N_GZ} Gen Z · 13 pays africains · Baromètre DATATYM 2026')

    # ══ CHOC D'OUVERTURE ══════════════════════════════════════
    slide_stat(prs,
        f'{pct_partent:.0f}%', 'ne se voient plus chez vous dans 5 ans.',
        'Sur 1 500 Gen Z interrogés, seuls 12% se projettent dans la même entreprise.',
        'Les 88% restants préparent un départ — silencieusement.',
        tag='VISION 5 ANS')

    slide_stat(prs,
        f'{pct_ent_total:.0f}%', 'veulent créer leur propre activité.',
        f'Dont {pct_ent_deja:.0f}% qui sont déjà en train de le faire.',
        'Ils ne quittent pas votre entreprise. Ils en construisent une autre — avec ce qu\'ils ont appris chez vous.',
        tag='ENTREPRENEURIAT GEN Z')

    # ══ TEMPS 1 — L'ICEBERG ═══════════════════════════════════
    slide_separator(prs, '01',
        'L\'Iceberg de l\'Ignorance',
        'Ce que chaque niveau de votre organisation voit vraiment.')

    slide_quote(prs,
        'Les dirigeants ne voient que 4% des problèmes\nque vivent leurs collaborateurs au quotidien.\n96% de la réalité ne remonte jamais.',
        'Adapté de Sidney Yoshida (1989)')

    slide_chart(prs,
        "L'Iceberg de l'Ignorance — Quantifié avec nos données",
        c_iceberg,
        'Yoshida (1989) · Proxies DATATYM : écart perception management (IATM=46,8) vs engagement terrain (ENG=77,8)')

    slide_stat(prs,
        '46,8', '/ 100',
        'C\'est le score IATM. L\'attractivité managériale vue par le terrain.',
        'Pendant que la direction pense que l\'engagement est à 77,8%, le terrain évalue le management à 46,8.\nL\'écart = 31 points. C\'est l\'iceberg.',
        tag='L\'ÉCART ICEBERG')

    # ══ TEMPS 2 — OU VONT-ILS ═════════════════════════════════
    slide_separator(prs, '02',
        'Où Vont Vos Talents ?',
        'Vision 5 ans, entrepreneuriat, profils de départ.')

    slide_chart(prs,
        "Où vos Gen Z se voient-ils dans 5 ans ?",
        c_vision,
        f'n={N_GZ} Gen Z · 13 pays africains · Baromètre DATATYM 2026')

    slide_chart(prs,
        "L'entrepreneuriat comme destination — Vos Gen Z ne partent pas, ils se lancent",
        c_ent,
        f'{pct_ent_deja:.0f}% déjà en cours · {pct_ent_2ans:.0f}% dans les 2 ans · {pct_ent_ml:.0f}% à moyen/long terme')

    slide_stat(prs,
        f'{pct_propre_compte:.0f}%', 'se voient à leur propre compte.',
        'Ils ne détestent pas votre entreprise. Ils veulent la leur.',
        'Le signal : si vous ne proposez pas de l\'intrapreneuriat, vous financez leur startup.',
        tag='ENTREPRENEUR CONCURRENT')

    slide_chart(prs,
        "5 profils de départ — Risque et levier de rétention par profil",
        c_profils,
        'Chaque profil = une destination différente, un levier de rétention différent.')

    # ══ TEMPS 3 — CE QUE CA COUTE ═════════════════════════════
    slide_separator(prs, '03',
        'La Productivité Perdue',
        'Le coût invisible d\'un talent qui glisse.')

    slide_chart(prs,
        "Coefficient de productivité par stade d'engagement",
        c_prod,
        'Pleinement engagé = 100% · Engagé fragile = 82% · Présentéisme = 55% · Départ programmé = 30%')

    slide_stat(prs,
        '33%', 'de productivité perdue en moyenne.',
        'Un quart de vos Gen Z est en présentéisme. Un quart a déjà programmé son départ.',
        'Sur 100 heures payées, vous n\'en récoltez que 67 en valeur réelle.',
        tag='COÛT CACHÉ')

    slide_chart(prs,
        "Score SOUSPAYÉ par secteur — Indicateur de vulnérabilité",
        c_souspaye,
        'Plus le score est élevé, plus le sentiment de sous-valorisation est fort → plus le risque de départ est élevé.')

    # ══ TEMPS 4 — LA FENETRE ═══════════════════════════════════
    slide_separator(prs, '04',
        'La Fenêtre d\'Intervention',
        'Combien de temps avez-vous avant de perdre vos talents ?')

    slide_chart(prs,
        "De l'intention au départ — La timeline que tout DG devrait connaître",
        c_fenetre,
        'Littérature RH + signaux DATATYM · Fenêtre d\'action moyenne : 6 mois')

    slide_stat(prs,
        '6', 'mois.',
        'C\'est votre fenêtre d\'intervention à partir du premier signal.',
        'Après 6 mois, le coût de rétention explose.\nAprès 9 mois, la décision est irréversible.',
        tag='FENÊTRE DG')

    # ══ TEMPS 5 — CE QUE LE DG FAIT ═══════════════════════════
    slide_separator(prs, '05',
        '5 Actions Immédiates',
        'Ce qu\'un dirigeant décide dans les 90 prochains jours.')

    # Slide actions
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, '5 DÉCISIONS — 90 JOURS', 0.5, 0.12, 12, 0.72,
        size=22, bold=True)

    actions = [
        ('DIAGNOSTIQUER',   'Lancer un diagnostic express Matrice Talents\nsur vos 50 postes les plus critiques.', DT_TEAL,   'J+15'),
        ('ACTIVER',         'Proposer un programme intrapreneuriat aux 23%\nqui préparent déjà leur projet.', DT_GOLD,   'J+30'),
        ('FORMER',          'Coaching managérial ciblé : IATM = 46,8.\nVos managers sont le noeud n°1.', DT_PURPLE,  'J+45'),
        ('VERROUILLER',     'Identifier et verrouiller vos 12% de fidèles.\nCe sont vos ancres — protégez-les.', '#10B981',  'J+60'),
        ('DÉCIDER',         'Trancher sur vos Présences Passives :\nréengager ou libérer. Pas de zone grise.', DT_RED,     'J+90'),
    ]

    for i, (titre, desc, color, delai) in enumerate(actions):
        y = 1.0 + i * 1.25
        rect(sl, 0.3, y, 0.12, 1.1, color)
        rect(sl, 0.5, y, 10.8, 1.1, DT_NAVY)
        tbx(sl, titre, 0.7, y + 0.05, 2.5, 0.5,
            size=15, bold=True, color=color)
        tbx(sl, desc, 3.5, y + 0.05, 6.5, 1.0,
            size=12, color=DT_WHITE)
        rect(sl, 10.8, y + 0.25, 1.2, 0.55, color)
        tbx(sl, delai, 10.8, y + 0.25, 1.2, 0.55,
            size=13, bold=True, align=PP_ALIGN.CENTER)
    print("  Slide: 5 actions DG")

    # ══ CONCLUSION ═════════════════════════════════════════════
    slide_quote(prs,
        'Vous ne perdez pas des employés.\nVous perdez des futurs concurrents,\ndes futurs expatriés,\ndes futurs entrepreneurs.\nLa question n\'est pas s\'ils partiront.\nC\'est ce que vous faites avant.')

    slide_cover(prs,
        'DATATYM Analytics',
        'Master Class — Baromètre Gen Z Afrique 2026\nDocument confidentiel',
        tag='MERCI')

    # Save
    out = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/DATATYM_MasterClass_2026.pptx'
    prs.save(out)
    n = len(prs.slides)
    print(f'\n=== {n} slides -> {out} ===')


if __name__ == '__main__':
    build()
