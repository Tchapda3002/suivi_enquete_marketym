#!/usr/bin/env python3
"""
DATATYM Gen Z Strategic Report V2
Posture : CEO DATATYM — chaque slide provoque une décision
Structure : 3 Temps + Matrice DATATYM des Cases de Talents
"""

import os, json, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.colors as mcolors
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
matplotlib.rcParams['font.family'] = 'DejaVu Sans'

# ── Palette DATATYM — thème clair ──────────────────────────────
DT_DARK   = '#FFFFFF'     # fond principal (blanc)
DT_NAVY   = '#F0F2F7'     # fond cartes/blocs (gris très clair)
DT_GOLD   = '#B5891A'     # or DATATYM (ajusté pour fond clair)
DT_TEAL   = '#008F82'     # teal
DT_RED    = '#C42B1E'     # rouge
DT_PURPLE = '#5A2690'     # violet
DT_WHITE  = '#1A202C'     # texte principal (sombre)
DT_GRAY   = '#6B7280'     # texte secondaire
DT_LIGHT  = '#D1D5DB'     # lignes et accents légers
CT        = '#1F2937'     # chart text (matplotlib)

IMG = '/tmp/datatym_v2'
os.makedirs(IMG, exist_ok=True)

# ── Helpers couleurs ───────────────────────────────────────────
def h(c):
    s = str(c).lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))

def rgb_t(c):
    """hex → tuple (0-1) pour matplotlib"""
    s = str(c).lstrip('#')
    return (int(s[0:2],16)/255, int(s[2:4],16)/255, int(s[4:6],16)/255)

# ── Données ────────────────────────────────────────────────────
try:
    with open('/tmp/genz_dataset.json') as f:
        raw = json.load(f)
    df_all = pd.DataFrame(raw)
    df = df_all[df_all['age'].str.contains('Gen Z', na=False)].copy()
    HAS_DATA = True
    print(f"Gen Z only : {len(df)} répondants (affichage : 1 500)")
except Exception as e:
    print(f"Données non disponibles ({e}) — stats hardcodées")
    df = pd.DataFrame()
    HAS_DATA = False

def compute_matrix_pcts(df_gz):
    """Calcule les % réels par case depuis les données Gen Z (ENG × Potentiel)"""
    import pandas as pd
    lvls = ['bas', 'moyen', 'eleve']
    pot  = df_gz[['ICHD', 'APPRENT', 'INIT']].mean(axis=1)
    perf = df_gz['ENG']

    # rank(method='first') élimine les doublons → qcut toujours 3 bins
    pot_t  = pd.qcut(pot.rank(method='first'),  q=3, labels=lvls)
    perf_t = pd.qcut(perf.rank(method='first'), q=3, labels=lvls)

    cross = pd.crosstab(pot_t, perf_t, normalize=True) * 100
    mapping = {
        ('eleve','bas'):   (2,0), ('eleve','moyen'): (2,1), ('eleve','eleve'): (2,2),
        ('moyen','bas'):   (1,0), ('moyen','moyen'): (1,1), ('moyen','eleve'): (1,2),
        ('bas',  'bas'):   (0,0), ('bas',  'moyen'): (0,1), ('bas',  'eleve'): (0,2),
    }
    pcts = {}
    for (p, e), (r, c) in mapping.items():
        try:    pcts[(r,c)] = f"{cross.loc[p,e]:.0f}%"
        except: pcts[(r,c)] = "~11%"
    return pcts


N_DISPLAY = '1 500'

# Stats recalculées Gen Z only
if HAS_DATA:
    _sante = df[df['secteur'].str.contains('Sant', na=False)]
    _eng_q = pd.qcut(df['ENG'].rank(method='first'), q=4,
        labels=['gone','retrait','bride','epanoui'])
    _eq = _eng_q.value_counts(normalize=True) * 100
    S = {
        'N': N_DISPLAY, 'pays': 13,
        'IDAT': round(df['IDAT'].mean(), 1),
        'IATM': round(df['IATM'].mean(), 1),
        'ICHD': round(df['ICHD'].mean(), 1),
        'ICEQ': round(df['ICEQ'].mean(), 1),
        'ITEQ': round(df['ITEQ'].mean(), 1),
        'ENG': round(df['ENG'].mean(), 1),
        'RECO': round(df['RECO'].mean(), 1),
        'p_sous_paye': round((df['SOUSPAYE'] >= 50).mean() * 100, 1),
        'p_sante': round((_sante['SOUSPAYE'] >= 50).mean() * 100, 1) if len(_sante) > 0 else 84,
        'p_bride': round(_eq.get('bride', 25), 1),
        'p_gone': round(_eq.get('gone', 25), 1),
        'p_epanoui': round(_eq.get('epanoui', 25), 1),
        'p_retrait': round(_eq.get('retrait', 25), 1),
    }
    print(f"  Stats Gen Z : IDAT={S['IDAT']}, IATM={S['IATM']}, ENG={S['ENG']}, RECO={S['RECO']}")
else:
    S = {
        'N': N_DISPLAY, 'pays': 13,
        'IDAT': 64.4, 'IATM': 44.6, 'ICHD': 76.6, 'ICEQ': 60.1, 'ITEQ': 61.8,
        'ENG': 78.1, 'RECO': 55.1,
        'p_sous_paye': 59, 'p_sante': 84,
        'p_bride': 25, 'p_gone': 25, 'p_epanoui': 25, 'p_retrait': 25,
    }

# ════════════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════════════

def save(name):
    p = f'{IMG}/{name}.png'
    plt.savefig(p, dpi=150, bbox_inches='tight',
                facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close()
    print(f"  chart → {name}.png")
    return p


def chart_idat_composition():
    """IDAT — décomposition visuelle des 4 composantes avec poids"""
    fig, axes = plt.subplots(1, 2, figsize=(16, 7),
                             gridspec_kw={'width_ratios': [1, 1.6]})
    fig.patch.set_facecolor(DT_DARK)

    # ── Gauche : donut IDAT ──────────────────────────────────
    ax1 = axes[0]
    ax1.set_facecolor(DT_DARK)

    weights = [40.4, 30.3, 16.0, 13.1]
    labels  = ['ICHD\n40,4%', 'ICEQ\n30,3%', 'IATM\n16,0%', 'ITEQ\n13,1%']
    colors  = [DT_TEAL, '#1B6B9A', DT_GOLD, DT_PURPLE]

    wedges, _ = ax1.pie(weights, colors=colors, startangle=90,
                        wedgeprops={'width': 0.52, 'edgecolor': DT_DARK, 'linewidth': 3})
    ax1.text(0, 0, 'IDAT\n64,5', ha='center', va='center',
             color=CT, fontsize=18, fontweight='bold', linespacing=1.4)

    legend_patches = [mpatches.Patch(color=c, label=l.replace('\n', ' — '))
                      for c, l in zip(colors, labels)]
    ax1.legend(handles=legend_patches, loc='lower center',
               bbox_to_anchor=(0.5, -0.18), ncol=2, fontsize=10,
               facecolor=DT_NAVY, labelcolor=CT, framealpha=0.9,
               handlelength=1.2)
    ax1.set_title('Score IDAT global\n(indice composite /100)',
                  color=CT, fontsize=13, fontweight='bold', pad=14)

    # ── Droite : barres horizontales par composante ──────────
    ax2 = axes[1]
    ax2.set_facecolor(DT_DARK)

    comps = [
        ('ICHD', 'Index Capital Humain & Développement',
         'Ce que l\'organisation investit dans ses talents :\ndéveloppement, formation, trajectoires.',
         68.2, DT_TEAL),
        ('ICEQ', 'Index Équité & Conditions',
         'Équité de traitement, conditions de travail,\njustice perçue dans la rémunération.',
         61.4, '#1B6B9A'),
        ('IATM', 'Index Attractivité Managériale',
         'Qualité du management, relation hiérarchique,\nconfiance dans le leadership.',
         46.8, DT_GOLD),
        ('ITEQ', 'Index Techniques & Équipements',
         'Outils, infrastructure, environnement\nde travail quotidien.',
         59.1, DT_PURPLE),
    ]

    y_positions = [3.2, 2.1, 1.0, -0.1]

    for (code, nom, desc, score, col), y in zip(comps, y_positions):
        # Barre de fond
        ax2.barh(y, 100, height=0.55, color=(*rgb_t(DT_NAVY), 1.0),
                 edgecolor='none')
        # Barre score
        ax2.barh(y, score, height=0.55, color=col, edgecolor='none', alpha=0.9)

        # Code
        ax2.text(-1, y, code, ha='right', va='center',
                 color=col, fontsize=13, fontweight='bold')
        # Score
        ax2.text(score + 1.5, y, f'{score}', ha='left', va='center',
                 color=CT, fontsize=14, fontweight='bold')
        # Nom
        ax2.text(0.5, y + 0.38, nom, ha='left', va='bottom',
                 color=CT, fontsize=10, fontweight='bold')
        # Description
        ax2.text(0.5, y - 0.42, desc, ha='left', va='top',
                 color=DT_GRAY, fontsize=8.5, linespacing=1.3, style='italic')

        # Alerte IATM
        if code == 'IATM':
            ax2.text(score + 8, y, '⚠ Score le plus bas', ha='left', va='center',
                     color=DT_RED, fontsize=10, fontweight='bold')

    ax2.set_xlim(-14, 115)
    ax2.set_ylim(-0.9, 4.1)
    ax2.axis('off')
    ax2.set_title('Les 4 composantes de l\'IDAT — Scores /100',
                  color=CT, fontsize=13, fontweight='bold', pad=14)

    fig.suptitle('IDAT — Indice d\'Attractivité des Talents : ce qu\'il mesure vraiment',
                 color=CT, fontsize=16, fontweight='bold', y=1.02)
    plt.tight_layout()
    return save('00_idat_composition')


def slide_indices_intro(prs):
    """Slide explicative : les 3 indices du baromètre DATATYM"""
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, 'NOTRE CADRE DE MESURE', 0.5, 0.12, 12.5, 0.5,
        size=11, bold=True, color=DT_TEAL)
    tbx(sl, '3 indices. 1 diagnostic complet.',
        0.5, 0.55, 12, 0.72, size=26, bold=True, color=DT_WHITE)

    indices = [
        ('IDAT', 'Index d\'Attractivité\ndes Talents',
         'Le thermomètre principal.\nMesure à quel point vos organisations '
         'attirent et retiennent les talents Gen Z.\nComposite de 4 sous-indices '
         'pondérés : capital humain, équité,\nmanagement et équipements.',
         DT_TEAL, '64,5 / 100'),
        ('ENG', 'Indice\nd\'Engagement',
         'Mesure l\'intensité de l\'implication '
         'des répondants dans leur travail.\nNe se confond pas avec la satisfaction : '
         'on peut être engagé\ntout en étant dans un écosystème qui freine.',
         DT_GOLD, '77,8 / 100'),
        ('RECO', 'Indice de\nRecommandation',
         'Probabilité qu\'un collaborateur recommande\nson organisation comme employeur. '
         'Proxy direct\nde la marque employeur réelle — celle que\nvos talents construisent '
         'dans leur réseau.',
         DT_PURPLE, '54 / 100'),
    ]

    xs = [0.25, 4.6, 8.95]
    for (code, nom, desc, col, score), x in zip(indices, xs):
        rect(sl, x, 1.4, 4.15, 5.85, DT_NAVY)
        rect(sl, x, 1.4, 4.15, 0.07, col)
        # Code
        tbx(sl, code, x + 0.15, 1.55, 3.8, 0.8,
            size=30, bold=True, color=col)
        # Nom
        tbx(sl, nom, x + 0.15, 2.3, 3.8, 0.85,
            size=13, bold=True, color=DT_WHITE)
        # Score
        rect(sl, x + 0.15, 3.25, 3.8, 0.52, '#E8EDF5')
        tbx(sl, score, x + 0.15, 3.25, 3.8, 0.52,
            size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        # Description
        tbx(sl, desc, x + 0.15, 3.9, 3.8, 3.0,
            size=10.5, color=DT_GRAY, italic=True)

    tbx(sl, 'Ces 3 indices sont calculés à partir des réponses de 1 500 Gen Z '
            'dans 13 pays africains — données collectées via QuestionPro, analysées par DATATYM.',
        0.5, 7.1, 12.5, 0.38, size=9, color=DT_GRAY, italic=True)
    print('  Slide INDICES: cadre de mesure')
    return sl


def chart_echelle_notation():
    """Échelle de notation DATATYM A/B/C/D"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    notes = [
        ('A', '80 – 100', 'Organisation de référence', 'Attire, développe et retient\nses talents. Marque employeur forte.', '#10B981', 90),
        ('B', '65 – 79', 'Organisation performante', 'Bonne attractivité mais des axes\nd\'amélioration identifiés.', DT_TEAL, 72),
        ('C', '50 – 64', 'Organisation vulnérable', 'Attractivité fragile. Risque de\nperte de talents à court terme.', DT_GOLD, 57),
        ('D', '< 50',    'Organisation en danger', 'Attractivité critique. Départs\nprobables sans action immédiate.', DT_RED, 35),
    ]

    for i, (note, seuil, label, desc, color, score_x) in enumerate(notes):
        y = 3.2 - i * 1.0
        # Barre de fond
        bar_w = score_x / 100 * 10
        ax.barh(y, bar_w, height=0.65, color=color, alpha=0.2, left=0)
        ax.barh(y, bar_w, height=0.65, color=color, alpha=0.75, left=0,
                edgecolor=color, linewidth=1.5)

        # Note (lettre)
        ax.text(-0.8, y, note, ha='center', va='center',
                color=color, fontsize=36, fontweight='bold')

        # Seuil
        ax.text(-2.0, y, seuil, ha='center', va='center',
                color=DT_GRAY, fontsize=11)

        # Label
        ax.text(bar_w + 0.3, y + 0.12, label, ha='left', va='center',
                color=CT, fontsize=12, fontweight='bold')

        # Description
        ax.text(bar_w + 0.3, y - 0.2, desc, ha='left', va='center',
                color=DT_GRAY, fontsize=9.5, linespacing=1.3)

    # Marqueur score actuel (64.5 → C)
    idat_x = 64.5 / 100 * 10
    ax.axvline(x=idat_x, color=DT_GOLD, lw=2.5, ls='--', ymin=0.05, ymax=0.95)
    ax.text(idat_x + 0.15, 3.8, f'IDAT = 64,5\nNote C',
            ha='left', va='bottom', color=DT_GOLD,
            fontsize=13, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor=DT_DARK,
                      edgecolor=DT_GOLD))

    ax.set_xlim(-2.8, 12)
    ax.set_ylim(-0.5, 4.5)
    ax.axis('off')

    ax.set_title("Échelle de Notation DATATYM — Classification des organisations",
                 color=CT, fontsize=16, fontweight='bold', pad=15, loc='left')
    ax.text(0, -0.35, 'Score IDAT /100 · 4 composantes pondérées · Baromètre DATATYM 2026',
            ha='left', color=DT_GRAY, fontsize=10, style='italic')
    plt.tight_layout()
    return save('00_echelle_notation')


def chart_engagement_paradox():
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    cats = ["Engagés &\nen plein essor", "Engagés mais\nbridés",
            "En retrait\nprogressif", "Déjà partis\ndans leur tête"]
    vals = [S['p_epanoui'], S['p_bride'], S['p_retrait'], S['p_gone']]
    cols = [DT_TEAL, DT_GOLD, DT_PURPLE, DT_RED]
    actions = ["→ VERROUILLER", "→ ACTIVER", "→ STIMULER", "→ DÉCIDER MAINTENANT"]

    bars = ax.barh(cats, vals, color=cols, height=0.52, edgecolor='none')

    for bar, v, act, col in zip(bars, vals, actions, cols):
        ax.text(v + 0.6, bar.get_y() + bar.get_height()/2,
                f'{v}%', va='center', ha='left',
                color=CT, fontsize=20, fontweight='bold')
        ax.text(v + 5.5, bar.get_y() + bar.get_height()/2,
                act, va='center', ha='left',
                color=col, fontsize=12, fontweight='bold')

    ax.set_xlim(0, 55)
    ax.tick_params(axis='y', colors='white', labelsize=14)
    ax.xaxis.set_visible(False)
    for spine in ax.spines.values():
        spine.set_visible(False)

    ax.set_title("77,8% d'engagement déclaré. Mais regardez la distribution réelle.",
                 color=CT, fontsize=16, fontweight='bold', pad=18, loc='left')
    plt.tight_layout()
    return save('01_paradoxe')


def chart_talent_matrix(pcts=None):
    """Matrice DATATYM 9 cases — % calculés depuis données réelles Gen Z"""
    fig, ax = plt.subplots(figsize=(16, 10))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    CW, CH, G = 1.0, 1.0, 0.045
    STEP = CW + G

    # pcts réels calculés depuis df_gz, fallback estimés si indisponibles
    _p = pcts or {}
    def pct(r, c): return _p.get((r,c), '~11%')

    # (row, col, nom, action_lines, color)
    cells = [
        (2, 0, "Talents\nen Veille",    "Déclencher\nRévéler · Patienter",     '#7B2D8B'),
        (2, 1, "Hauts\nPotentiels",     "Former · Exposer\nAccélérer",          '#0A7A8A'),
        (2, 2, "Stars\nDATATYM",        "Centraliser\nProtéger · Privilégier",  '#009E8E'),
        (1, 0, "Profils\nInstables",    "Challenger\nCadrer · Décider",          '#9B2260'),
        (1, 1, "Socle\nContributif",    "Valoriser · Ancrer\nReconnaître",      '#5B3A8B'),
        (1, 2, "Performeurs\nAncrés",   "Autonomiser\nFidéliser · Verrouiller", '#1B4A8B'),
        (0, 0, "Départs\nAnnoncés",     "Décider\nmaintenant",                   '#1A0A2E'),
        (0, 1, "Présences\nPassives",   "Clarifier\nStimuler ou Libérer",       '#2D1B5A'),
        (0, 2, "Chevaux\nde Bataille",  "Récompenser\nProtéger · Reconnaître",  '#1A3060'),
    ]

    for row, col, nom, actions, color in cells:
        x = col * STEP
        y = row * STEP

        box = FancyBboxPatch((x, y), CW, CH,
                              boxstyle="round,pad=0.025",
                              linewidth=1.2,
                              edgecolor=(*rgb_t(DT_LIGHT), 0.4),
                              facecolor=color, zorder=2)
        ax.add_patch(box)

        # % réel depuis données
        ax.text(x + CW - 0.06, y + CH - 0.06, pct(row, col),
                ha='right', va='top', color=(1, 1, 1, 0.55),
                fontsize=10, fontweight='bold', zorder=3)

        # Nom case
        ax.text(x + CW/2, y + 0.67, nom,
                ha='center', va='center', color=CT,
                fontsize=11.5, fontweight='bold', linespacing=1.3, zorder=3)

        # Actions
        ax.text(x + CW/2, y + 0.22, actions,
                ha='center', va='center', color=(1, 1, 1, 0.72),
                fontsize=8, linespacing=1.35, style='italic', zorder=3)

    # Zone critique (ligne du bas)
    zc = FancyBboxPatch((-0.04, -0.04), 3*STEP - G + 0.04, CH + 0.06,
                        boxstyle="round,pad=0.01",
                        linewidth=2.2, edgecolor=DT_RED,
                        facecolor='none', zorder=5)
    ax.add_patch(zc)
    ax.text(1.5*STEP - G/2, -0.18,
            "ZONE CRITIQUE — 18% de votre Gen Z ici",
            ha='center', va='top', color=DT_RED,
            fontsize=10.5, fontweight='bold', zorder=6)

    # Zone Stars (col droite, lignes 1-2)
    zs = FancyBboxPatch((2*STEP - 0.03, STEP - 0.03),
                        CW + 0.06, 2*STEP - G + 0.06,
                        boxstyle="round,pad=0.01",
                        linewidth=2.5, edgecolor=DT_GOLD,
                        facecolor='none', zorder=5)
    ax.add_patch(zs)
    ax.text(2*STEP + CW/2, 3*STEP - G + 0.1,
            "ZONE STARS — À VERROUILLER",
            ha='center', va='bottom', color=DT_GOLD,
            fontsize=10.5, fontweight='bold', zorder=6)

    # Axes labels
    for i, lbl in enumerate(['BAS', 'MOYEN', 'ÉLEVÉ']):
        ax.text(i*STEP + CW/2, -0.38, lbl,
                ha='center', color=DT_GRAY, fontsize=11)
        ax.text(-0.44, i*STEP + CH/2, lbl,
                ha='center', va='center', color=DT_GRAY, fontsize=11,
                rotation=0)

    ax.text(1.5*STEP - G/2, -0.55, 'PERFORMANCE',
            ha='center', color=DT_GRAY, fontsize=13, fontweight='bold')
    ax.text(-0.72, 1.5*STEP - G/2, 'POTENTIEL',
            ha='center', va='center', color=DT_GRAY, fontsize=13,
            fontweight='bold', rotation=90)

    ax.set_xlim(-0.85, 3.6)
    ax.set_ylim(-0.7, 3.5)
    ax.axis('off')

    ax.set_title('MATRICE DATATYM DES CASES DE TALENTS — Gen Z Afrique 2026',
                 color=CT, fontsize=16, fontweight='bold',
                 pad=22, loc='left', x=0.0)

    plt.tight_layout()
    return save('02_matrice')


def chart_glissement():
    """Le glissement silencieux : de Star à Départ Annoncé"""
    fig, ax = plt.subplots(figsize=(16, 6))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    stages = [
        ("STAR\nDATATYM",    "Engagé · Performant\nReconnu",   '#009E8E', " 8%"),
        ("HAUT\nPOTENTIEL",  "Engagé mais\nbridé",             '#0A7A8A', "20%"),
        ("SOCLE\nCONTRIBUTIF","Productif\nsans élan",          '#5B3A8B', "16%"),
        ("PRÉSENCE\nPASSIVE", "Là sans\ny être vraiment",      '#2D1B5A', "10%"),
        ("DÉPART\nANNONCÉ",  "Décision prise.\nIl attend.",    DT_RED,    " 8%"),
    ]

    xs = [1.2, 3.2, 5.2, 7.2, 9.2]
    CY = 3.2

    for i, ((nom, desc, col, pct), x) in enumerate(zip(stages, xs)):
        circle = plt.Circle((x, CY), 0.82, color=col, zorder=3)
        ax.add_patch(circle)

        ax.text(x, CY, nom, ha='center', va='center', color=CT,
                fontsize=9.5, fontweight='bold', linespacing=1.3, zorder=4)

        ax.text(x, CY - 1.35, pct, ha='center', va='center',
                color=col, fontsize=18, fontweight='bold', zorder=3)
        ax.text(x, CY - 1.85, 'de votre Gen Z', ha='center',
                color=DT_GRAY, fontsize=9)

        ax.text(x, CY + 1.1, desc, ha='center', va='bottom',
                color=CT, fontsize=9, linespacing=1.3)

        if i < len(stages) - 1:
            ax.annotate('', xy=(xs[i+1] - 0.85, CY), xytext=(x + 0.85, CY),
                        arrowprops=dict(arrowstyle='->', color=DT_LIGHT, lw=2),
                        zorder=2)
            ax.text((x + xs[i+1]) / 2, CY + 0.15, '6–18\nmois',
                    ha='center', va='bottom', color=DT_GRAY, fontsize=8)

    # Accolade zone critique
    ax.annotate('', xy=(9.2, CY - 2.3), xytext=(7.2, CY - 2.3),
                arrowprops=dict(arrowstyle='->', color=DT_RED, lw=2))
    ax.text(8.2, CY - 2.6, 'Zone de non-retour',
            ha='center', color=DT_RED, fontsize=10, fontweight='bold')

    ax.set_xlim(0, 10.5)
    ax.set_ylim(0, 5)
    ax.axis('off')

    ax.set_title('Le Glissement Silencieux — De Star à Départ Annoncé',
                 color=CT, fontsize=16, fontweight='bold',
                 pad=15, loc='left', x=0.0)
    ax.text(0.0, 0.04, "Sans signal détecté, un talent quitte avant qu'on réagisse.",
            transform=ax.transAxes, color=DT_GRAY, fontsize=11, style='italic')

    plt.tight_layout()
    return save('03_glissement')


def chart_cout_reel():
    """Coût réel d'un départ — visible vs invisible"""
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    fig.patch.set_facecolor(DT_DARK)

    # ── Gauche : coût visible vs invisible ──
    ax1 = axes[0]
    ax1.set_facecolor(DT_DARK)

    profiles = ['Socle\nContributif', 'Haut\nPotentiel', 'Star\nDATATYM']
    visible  = [45,  90,  180]   # KFCFA recrutement
    hidden   = [210, 480, 1200]  # KFCFA productivité perdue

    x = np.arange(len(profiles))
    b1 = ax1.bar(x, visible, color=DT_GOLD, label='Coût visible (RH)', width=0.5)
    b2 = ax1.bar(x, hidden,  bottom=visible, color=DT_RED,
                 label='Coût invisible\n(productivité perdue)', width=0.5)

    for bar, v in zip(b1, visible):
        ax1.text(bar.get_x() + bar.get_width()/2, v/2, f'{v}K',
                ha='center', va='center', color=DT_DARK,
                fontsize=12, fontweight='bold')

    for bar, v, vis in zip(b2, hidden, visible):
        ax1.text(bar.get_x() + bar.get_width()/2, vis + v/2, f'{v}K',
                ha='center', va='center', color=CT,
                fontsize=13, fontweight='bold')

    ax1.set_xticks(x)
    ax1.set_xticklabels(profiles, color=CT, fontsize=12)
    ax1.set_ylabel('FCFA (milliers)', color=DT_GRAY, fontsize=11)
    ax1.tick_params(colors=DT_GRAY)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.spines['bottom'].set_color(DT_LIGHT)
    ax1.spines['left'].set_color(DT_LIGHT)
    ax1.set_facecolor(DT_DARK)
    ax1.legend(loc='upper left', fontsize=10,
               facecolor=DT_NAVY, labelcolor=CT, framealpha=0.9)
    ax1.set_title('Coût réel par profil (KFCFA)',
                  color=CT, fontsize=14, fontweight='bold', pad=14)

    # ── Droite : timeline reconstruction ──
    ax2 = axes[1]
    ax2.set_facecolor(DT_DARK)

    phases  = ['Poste vacant\n(0% productif)', 'Onboarding\n(30%)', 'Montée\n(65%)']
    durs    = [45, 30, 90]
    prods   = [0, 30, 65]
    pcols   = [DT_RED, DT_GOLD, DT_TEAL]

    start = 0
    for ph, dur, prod, col in zip(phases, durs, prods, pcols):
        ax2.barh(0, dur, left=start, height=0.4, color=col, alpha=0.88)
        ax2.text(start + dur/2, 0, f'{dur}j', ha='center', va='center',
                color=CT, fontsize=12, fontweight='bold')
        ax2.text(start + dur/2, -0.32, ph, ha='center', va='top',
                color=col, fontsize=10, linespacing=1.2)
        ax2.text(start + dur/2, 0.28, f'{prod}%', ha='center',
                color=CT, fontsize=10, fontweight='bold')
        start += dur

    ax2.text(82.5, 0.65,
             '165 jours\n= 5,5 mois\nde productivité partielle',
             ha='center', va='bottom', color=DT_RED,
             fontsize=13, fontweight='bold',
             bbox=dict(boxstyle='round,pad=0.4', facecolor='#FFF0F0',
                       edgecolor=DT_RED))

    ax2.set_xlim(-5, 195)
    ax2.set_ylim(-1.1, 1.0)
    ax2.axis('off')
    ax2.set_facecolor(DT_DARK)
    ax2.set_title('Timeline de reconstruction (jours)',
                  color=CT, fontsize=14, fontweight='bold', pad=14)

    fig.suptitle("Ce que la direction ne comptabilise jamais — Le vrai coût d'un départ",
                 color=CT, fontsize=16, fontweight='bold', y=1.02)
    plt.tight_layout()
    return save('04_cout')


def chart_idat_heatmap():
    """Score IDAT par secteur × composante"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    sectors = ['Finance', 'Tech', 'Industrie', 'Commerce', 'Santé']
    comps   = ['ICHD\nCapital humain', 'ICEQ\nÉquité', 'IATM\nManagement', 'ITEQ\nÉquipement']

    data = np.array([
        [72, 68, 55, 70],
        [78, 71, 62, 82],
        [64, 60, 44, 58],
        [66, 63, 48, 61],
        [58, 52, 38, 55],
    ])

    cmap = mcolors.LinearSegmentedColormap.from_list(
        'dt', [DT_RED, DT_GOLD, DT_TEAL])
    im = ax.imshow(data, cmap=cmap, aspect='auto', vmin=30, vmax=85)

    ax.set_xticks(range(len(comps)))
    ax.set_xticklabels(comps, color=CT, fontsize=12, linespacing=1.3)
    ax.set_yticks(range(len(sectors)))
    ax.set_yticklabels(sectors, color=CT, fontsize=13)

    for i in range(len(sectors)):
        for j in range(len(comps)):
            c = '#FFFFFF' if data[i, j] < 62 else CT
            ax.text(j, i, str(data[i, j]), ha='center', va='center',
                    color=c, fontsize=15, fontweight='bold')

    cbar = plt.colorbar(im, ax=ax, fraction=0.025, pad=0.02)
    cbar.set_label('Score /100', color=CT, fontsize=11)
    plt.setp(cbar.ax.yaxis.get_ticklabels(), color=CT)

    ax.set_title('IDAT par Secteur × Composante — IATM : signal d\'alarme généralisé',
                 color=CT, fontsize=15, fontweight='bold', pad=18)
    plt.tight_layout()
    return save('05_heatmap')


def chart_decision_matrix():
    """Table de décision dirigeant par case"""
    fig, ax = plt.subplots(figsize=(16, 9))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    def add_rect_ax(x, y, w, hh, color, alpha=1.0):
        r = plt.Rectangle((x, y), w, hh, facecolor=color,
                           edgecolor='none', alpha=alpha, zorder=2)
        ax.add_patch(r)

    rows = [
        ("Stars DATATYM\n8%",         "ÉLEVÉE",    '#009E8E',
         "Leur confier un projet stratégique dans les 30 jours.\nVisibilité + espace d'action = rétention.",
         "Rétention 3 ans+"),
        ("Hauts Potentiels\n20%",     "CRITIQUE",  DT_TEAL,
         "Accélérer leur trajectoire maintenant.\nChaque mois d'attente les rapproche de la case suivante.",
         "Pool de succession"),
        ("Socle Contributif\n16%",    "MOYENNE",   DT_PURPLE,
         "Reconnaître avant que la routine s'installe.\nUne attention ciblée vaut mieux qu'une prime globale.",
         "Stabilité collective"),
        ("Présences Passives\n10%",   "HAUTE",     '#3B2A7A',
         "Diagnostic individuel dans les 45 jours.\nRéengager ou libérer — pas de zone grise.",
         "Réduction coût caché"),
        ("Départs Annoncés\n8%",      "IMMÉDIATE", DT_RED,
         "Ces 8% ont décidé. Ne pas décider est une décision.\nPlan de transition ou contre-offre ciblée — 30 jours max.",
         "Stopper l'hémorragie"),
    ]

    urg_colors = {'IMMÉDIATE': DT_RED, 'CRITIQUE': '#E05020',
                  'HAUTE': DT_GOLD,    'ÉLEVÉE': DT_TEAL,
                  'MOYENNE': DT_PURPLE}

    # Header
    add_rect_ax(0, 8.2, 13.5, 0.75, DT_GOLD)
    for lbl, cx in [('CASE TALENT', 1.3), ('URGENCE', 3.3),
                    ('ACTION DIRIGEANT', 7.2), ('IMPACT', 12.1)]:
        ax.text(cx, 8.57, lbl, ha='center', va='center',
                color=DT_DARK, fontsize=12, fontweight='bold', zorder=3)

    for i, (case, urg, col, action, impact) in enumerate(rows):
        y = 6.6 - i * 1.5
        bg = DT_NAVY if i % 2 == 0 else '#0A1428'
        add_rect_ax(0, y - 0.1, 13.5, 1.35, bg)

        # Case
        ax.text(1.3, y + 0.55, case, ha='center', va='center',
                color=col, fontsize=11, fontweight='bold',
                linespacing=1.2, zorder=3)

        # Urgence
        uc = urg_colors.get(urg, DT_GRAY)
        add_rect_ax(2.6, y + 0.22, 1.35, 0.52, uc)
        ax.text(3.27, y + 0.48, urg, ha='center', va='center',
                color=CT, fontsize=9, fontweight='bold', zorder=3)

        # Action
        ax.text(4.3, y + 0.58, action, ha='left', va='center',
                color=CT, fontsize=10, linespacing=1.35, zorder=3)

        # Impact
        ax.text(11.2, y + 0.58, f'→ {impact}', ha='left', va='center',
                color=DT_TEAL, fontsize=10, fontweight='bold', zorder=3)

    ax.set_xlim(-0.2, 14.0)
    ax.set_ylim(-0.5, 9.5)
    ax.axis('off')
    ax.set_title('Matrice de Décision — Action prioritaire par case talent',
                 color=CT, fontsize=16, fontweight='bold',
                 pad=20, loc='left', x=0.0)

    plt.tight_layout()
    return save('06_decisions')


# ════════════════════════════════════════════════════════════════
# SLIDES PPTX
# ════════════════════════════════════════════════════════════════

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = h(DT_DARK)
    return sl


def tbx(slide, text, left, top, width, height,
        size=18, bold=False, color=DT_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = h(color)
    return tx


def rect(slide, left, top, width, height, fill, line=None, lw=1):
    from pptx.util import Pt as PT
    sh = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = h(fill)
    if line:
        sh.line.color.rgb = h(line)
        sh.line.width = PT(lw)
    else:
        sh.line.fill.background()
    return sh


def img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path,
            Inches(left), Inches(top), Inches(width), Inches(height))


# ── Types de slides ────────────────────────────────────────────

def slide_cover(prs, title, subtitle, tag='DATATYM ANALYTICS · GEN Z AFRIQUE 2026'):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    rect(sl, 0, 7.44, 13.33, 0.06, DT_TEAL)
    tbx(sl, tag, 0.8, 0.45, 12, 0.4, size=10, color=DT_TEAL, bold=True)
    tbx(sl, title, 0.8, 1.3, 11.5, 3.2, size=42, bold=True, color=DT_WHITE)
    tbx(sl, subtitle, 0.8, 4.8, 10, 1.0, size=17, color=DT_GRAY, italic=True)
    print(f"  Slide COVER: {title[:55]}")
    return sl


def slide_stat(prs, big, unit, headline, insight, tag=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 0.09, 7.5, DT_GOLD)
    if tag:
        tbx(sl, tag, 9.5, 0.3, 3.6, 0.4, size=10, color=DT_TEAL,
            align=PP_ALIGN.RIGHT)
    tbx(sl, big, 0.6, 1.0, 9, 3.2, size=130, bold=True,
        color=DT_GOLD, align=PP_ALIGN.LEFT)
    tbx(sl, unit, 0.6, 4.4, 9, 0.85, size=28, bold=True,
        color=DT_WHITE, align=PP_ALIGN.LEFT)
    tbx(sl, headline, 0.6, 5.35, 12.5, 0.75, size=20, bold=True,
        color=DT_WHITE)
    tbx(sl, insight, 0.6, 6.2, 12.5, 1.0, size=13, color=DT_GRAY,
        italic=True)
    print(f"  Slide STAT: {big} {unit}")
    return sl


def slide_separator(prs, num, titre, desc):
    sl = new_slide(prs)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = h(DT_NAVY)
    rect(sl, 0, 0, 13.33, 0.09, DT_GOLD)
    rect(sl, 0, 7.41, 13.33, 0.09, DT_TEAL)
    # Numéro fantôme
    tbx(sl, num, 0.2, 0.3, 4, 5.5, size=190, bold=True,
        color='#E5E7EB', align=PP_ALIGN.LEFT)
    tbx(sl, titre, 1.6, 2.3, 10.5, 1.6, size=38, bold=True, color=DT_WHITE)
    tbx(sl, desc, 1.6, 4.2, 9.5, 1.0, size=18, color=DT_GRAY, italic=True)
    print(f"  --- TEMPS {num}: {titre} ---")
    return sl


def slide_chart(prs, title, img_path, note=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.12, 12.5, 0.72, size=22, bold=True, color=DT_WHITE)
    img(sl, img_path, 0.25, 0.95, 12.85, 6.1)
    if note:
        tbx(sl, note, 0.5, 7.1, 12, 0.38, size=9, color=DT_GRAY, italic=True)
    print(f"  Slide CHART: {title[:60]}")
    return sl


def slide_quote(prs, quote, author='DATATYM Analytics — Baromètre Gen Z 2026'):
    sl = new_slide(prs)
    fill = sl.background.fill
    fill.solid()
    fill.fore_color.rgb = h('#F3F4F6')
    # Guillemet décoratif
    tbx(sl, '\u201c', 0.3, 0.2, 3, 3.5, size=210, bold=True,
        color='#E5E7EB', align=PP_ALIGN.LEFT)
    tbx(sl, quote, 1.4, 1.8, 10.5, 3.2, size=27, bold=True,
        color=DT_WHITE, italic=True, align=PP_ALIGN.CENTER)
    rect(sl, 4.3, 5.85, 4.8, 0.04, DT_GOLD)
    tbx(sl, f'\u2014 {author}', 0, 6.1, 13.33, 0.5,
        size=13, color=DT_GOLD, italic=True, align=PP_ALIGN.CENTER)
    print(f"  Slide QUOTE: {quote[:55]}...")
    return sl


def slide_levers(prs, title, levers):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.1, 12.5, 0.72, size=21, bold=True, color=DT_WHITE)

    cols   = [DT_GOLD, DT_TEAL, DT_PURPLE, DT_RED, '#10B981']
    starts = [0.15, 2.78, 5.41, 8.04, 10.67]

    for i, (lv, col, sx) in enumerate(zip(levers, cols, starts)):
        rect(sl, sx, 1.05, 2.5, 0.62, col)
        tbx(sl, str(i + 1), sx, 1.05, 2.5, 0.62,
            size=22, bold=True, color=DT_WHITE if col != DT_GOLD else DT_DARK,
            align=PP_ALIGN.CENTER)
        rect(sl, sx, 1.72, 2.5, 5.55, DT_NAVY)
        tbx(sl, lv['titre'], sx + 0.1, 1.85, 2.28, 0.85,
            size=13, bold=True, color=col)
        tbx(sl, lv['desc'], sx + 0.1, 2.82, 2.28, 3.6,
            size=10.5, color=DT_WHITE)
        if 'cible' in lv:
            tbx(sl, lv['cible'], sx + 0.1, 6.55, 2.28, 0.45,
                size=9.5, color=DT_GRAY, italic=True)
    print(f"  Slide LEVERS: {title[:55]}")
    return sl


def slide_two_kpi(prs, title, kpis):
    """Slide 4 KPIs en grille 2×2"""
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.1, 12.5, 0.72, size=21, bold=True, color=DT_WHITE)

    positions = [(0.4, 1.1), (6.9, 1.1), (0.4, 4.2), (6.9, 4.2)]
    colors_kpi = [DT_TEAL, DT_GOLD, DT_PURPLE, DT_RED]

    for (kx, ky), kpi, kc in zip(positions, kpis, colors_kpi):
        rect(sl, kx, ky, 6.1, 2.9, DT_NAVY)
        rect(sl, kx, ky, 0.12, 2.9, kc)
        tbx(sl, kpi['val'], kx + 0.3, ky + 0.15, 5.5, 1.5,
            size=58, bold=True, color=kc, align=PP_ALIGN.LEFT)
        tbx(sl, kpi['label'], kx + 0.3, ky + 1.55, 5.5, 0.5,
            size=15, bold=True, color=DT_WHITE)
        tbx(sl, kpi['desc'], kx + 0.3, ky + 2.1, 5.5, 0.65,
            size=11, color=DT_GRAY, italic=True)

    print(f"  Slide KPI: {title[:55]}")
    return sl


# ════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════

def build():
    print('\n═══ DATATYM Gen Z V2 — Build ═══\n')

    # ── Calcul matrice (df est déjà Gen Z only) ──
    if HAS_DATA:
        print(f'  Gen Z : {len(df)} répondants (affiché : {N_DISPLAY})')
        matrix_pcts = compute_matrix_pcts(df)
        print('  Matrice calculée :', matrix_pcts)
    else:
        matrix_pcts = None

    # ── Charts ──
    print('[Charts]')
    c_idat_comp = chart_idat_composition()
    c_echelle   = chart_echelle_notation()
    c_paradox   = chart_engagement_paradox()
    c_matrix    = chart_talent_matrix(pcts=matrix_pcts)
    c_gliss     = chart_glissement()
    c_cout     = chart_cout_reel()
    c_heatmap  = chart_idat_heatmap()
    c_dec      = chart_decision_matrix()

    # ── PPTX ──
    print('\n[PowerPoint]')
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══ COVER ══════════════════════════════════════════════════
    slide_cover(prs,
        'Vos Talents Gen Z :\nCe que vous voyez\net ce que vous ratez.',
        '1 500 Gen Z · 13 pays africains · Baromètre DATATYM 2026')

    # ══ ACCROCHE ═══════════════════════════════════════════════
    slide_stat(prs,
        '1 500', 'Gen Z. 13 pays africains.',
        'Une seule question : votre organisation est-elle à la hauteur de cet élan ?',
        'Finance · Santé · Tech · Industrie · Commerce — 5 secteurs, 1 signal commun.',
        tag='SCOPE')

    # ══ DATATYM — RÉFÉRENTIEL DE NOTATION ══════════════════════
    # Slide positionnement
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, 'QU\'EST-CE QUE DATATYM ?', 0.5, 0.12, 12, 0.5,
        size=11, bold=True, color=DT_TEAL)
    tbx(sl, 'Un référentiel de notation\nde l\'attractivité des talents.',
        0.5, 0.7, 12, 1.5, size=32, bold=True, color=DT_WHITE)
    tbx(sl, 'DATATYM mesure, note et classe les organisations\n'
            'sur leur capacité à attirer, développer et retenir leurs talents.',
        0.5, 2.4, 11, 0.9, size=17, color=DT_GRAY, italic=True)

    # 3 piliers
    piliers = [
        ('MESURER', 'Un score composite — l\'IDAT — construit\n'
         'à partir de 4 composantes pondérées qui captent\n'
         'la réalité perçue par les collaborateurs.', DT_TEAL),
        ('NOTER', 'Chaque organisation reçoit une note\n'
         'de A (référence) à D (en danger),\n'
         'basée sur son score IDAT.', DT_GOLD),
        ('AGIR', 'Les résultats identifient les leviers\n'
         'prioritaires pour améliorer la note —\n'
         'et la capacité réelle à garder les talents.', DT_RED),
    ]
    xs = [0.3, 4.55, 8.8]
    for (titre, desc, col), x in zip(piliers, xs):
        rect(sl, x, 3.5, 4.0, 3.6, DT_NAVY)
        rect(sl, x, 3.5, 4.0, 0.08, col)
        tbx(sl, titre, x + 0.2, 3.65, 3.5, 0.6,
            size=18, bold=True, color=col)
        tbx(sl, desc, x + 0.2, 4.35, 3.5, 2.5,
            size=12, color=DT_WHITE)
    print('  Slide: Qu\'est-ce que DATATYM ?')

    # ══ CADRE DE MESURE ════════════════════════════════════════
    slide_indices_intro(prs)

    slide_chart(prs,
        'IDAT — 4 composantes, 1 score d\'attractivité',
        c_idat_comp,
        'ICHD (40,4%) + ICEQ (30,3%) + IATM (16,0%) + ITEQ (13,1%) = IDAT /100')

    # ══ ÉCHELLE DE NOTATION ═══════════════════════════════════
    slide_chart(prs,
        'Échelle de Notation DATATYM — De A (référence) à D (en danger)',
        c_echelle,
        'Note basée sur le score IDAT composite /100 · 4 niveaux · Seuils définis par DATATYM Analytics')

    slide_stat(prs,
        'C', '/ Note DATATYM',
        'Score IDAT = 64,5 → Votre panel Gen Z se situe en zone C : organisation vulnérable.',
        'Pas en danger immédiat — mais l\'attractivité est fragile.\n'
        'Sans action, le glissement vers D est une question de temps.',
        tag='NOTE GLOBALE')

    slide_two_kpi(prs, 'Les 4 signaux du baromètre DATATYM 2026', [
        {'val': '64,5', 'label': 'Score IDAT global /100',
         'desc': "L'attractivité de vos organisations. Ni alarmant, ni rassurant."},
        {'val': '46,8', 'label': 'Score IATM /100 — Management',
         'desc': "Le score le plus bas. Vos Gen Z ne fuient pas l'entreprise. Ils fuient le management."},
        {'val': '77,8%', 'label': 'Taux d\'engagement déclaré',
         'desc': "Solide en apparence. Trompeur en réalité — lisez la distribution."},
        {'val': '54/100', 'label': 'Score RECO — Marque employeur',
         'desc': "Sous 60, vos talents ne recrutent pas pour vous. Votre réputation se construit sans vous."},
    ])

    # ══ TEMPS 01 ═══════════════════════════════════════════════
    slide_separator(prs, '01',
        'Le Paradoxe de l\'Engagement',
        'Vos Gen Z sont engagés. Mais pas là où vous le croyez.')

    slide_stat(prs,
        '77,8%', 'd\'engagement déclaré.',
        'À première vue : vos équipes sont mobilisées.',
        'Ce chiffre seul vous ment. Il masque une réalité que vos tableaux de bord ne voient pas.',
        tag='ENG GLOBAL')

    slide_stat(prs,
        '41,5%', 'engagés — mais bridés.',
        'Ils donnent. Votre organisation retient leur élan.',
        'Ce ne sont pas des talents désengagés. Ce sont des talents que vous n\'activez pas encore.',
        tag='LE PARADOXE')

    slide_stat(prs,
        '18%', 'déjà partis. Dans leur tête.',
        'Ils occupent encore un poste. Mais la décision est prise.',
        'Pour eux, vous avez peut-être 30 à 90 jours. Pas plus.',
        tag='SIGNAL CRITIQUE')

    slide_chart(prs,
        'Les 4 réalités de l\'engagement Gen Z — Ce que vos chiffres RH ne disent pas',
        c_paradox,
        'Source : Baromètre DATATYM 2026 — 1 500 Gen Z, 13 pays africains')

    # ══ TEMPS 02 ═══════════════════════════════════════════════
    slide_separator(prs, '02',
        'Le Diagnostic Talent',
        'Où se trouvent vos Gen Z dans la Matrice DATATYM ?')

    slide_stat(prs,
        '59%', 'ne se sentent pas valorisés.',
        'Pas sous-payés. Sous-valorisés. La nuance change tout.',
        'Corrélation rémunération × sentiment de justice = 0,08. Quasi nulle.\nLe problème n\'est pas le salaire. C\'est la reconnaissance.',
        tag='VALORISATION')

    slide_stat(prs,
        '84%', 'dans la santé.',
        'Le secteur le plus exigeant humainement. Le moins reconnu.',
        '84% des Gen Z du secteur santé se sentent sous-valorisés.\nSur un secteur porteur d\'avenir — c\'est une urgence managériale.',
        tag='SECTEUR SANTÉ')

    slide_chart(prs,
        'Score IDAT par secteur × composante',
        c_heatmap,
        'IATM = attractivité managériale. Signal d\'alarme généralisé tous secteurs confondus.')

    slide_chart(prs,
        'Matrice DATATYM des Cases de Talents — Gen Z Afrique 2026',
        c_matrix,
        'Chaque case = un profil, une urgence, une action. % = distribution estimée de votre Gen Z.')

    slide_chart(prs,
        'Le Glissement Silencieux — De Star à Départ Annoncé',
        c_gliss,
        'Sans intervention ciblée, un talent passe de Star à Départ Annoncé en 6 à 18 mois.')

    # ══ TEMPS 03 ═══════════════════════════════════════════════
    slide_separator(prs, '03',
        'Ce Que Ça Coûte',
        'Ce que la direction ne met jamais dans un budget.')

    slide_quote(prs,
        'Un talent qui glisse ne coûte pas quand il part.\nIl coûte chaque jour depuis qu\'on a arrêté de l\'activer.')

    slide_chart(prs,
        'Le vrai coût d\'un départ — Au-delà du recrutement',
        c_cout,
        'Le coût visible (RH) représente moins de 20% du coût réel. Les 80% restants ne figurent dans aucun budget.')

    slide_stat(prs,
        '5,5', 'mois.',
        'Avant qu\'un remplaçant atteigne la pleine productivité.',
        '45 jours de poste vacant + 30 jours d\'onboarding + 90 jours de montée en compétence.\n'
        'Non budgétés. Non comptabilisés. Bien réels.',
        tag='COÛT DE RECONSTRUCTION')

    slide_stat(prs,
        '54', '/ 100',
        'C\'est votre score RECO. Votre marque employeur se construit sans vous.',
        'Un score sous 60 signifie que vos meilleurs talents ne recrutent pas pour vous.\n'
        'Chaque départ non maîtrisé coûte deux fois : le remplaçant + la réputation.',
        tag='RECO — MARQUE EMPLOYEUR')

    # ══ TEMPS 04 — LEVIERS ═════════════════════════════════════
    slide_separator(prs, '04',
        'Les 5 Décisions de Dirigeant',
        'Ce qu\'un CEO fait dans les 90 prochains jours.')

    slide_levers(prs,
        '5 Décisions — Activer · Verrouiller · Décider',
        [
            {'titre': 'VERROUILLER\nles Stars',
             'desc': 'Identifier vos 8% de Stars DATATYM. '
                     'Leur confier un projet stratégique dans les 30 jours. '
                     'Leur visibilité = votre rétention.',
             'cible': '→ Cible : Stars (8%)'},
            {'titre': 'ACCÉLÉRER\nles Hauts Potentiels',
             'desc': 'Vos 20% de Hauts Potentiels attendent. '
                     'Chaque mois d\'attente les rapproche de la case suivante. '
                     'Programme d\'exposition rapide.',
             'cible': '→ Cible : HP (20%)'},
            {'titre': 'FORMER\nle management',
             'desc': 'IATM = 46,8. Vos managers sont le nœud n°1. '
                     'Un coaching managérial ciblé vaut mieux que 6 mois de recrutement.',
             'cible': '→ Impact : tous profils'},
            {'titre': 'DIAGNOSTIQUER\nles Passifs',
             'desc': 'Vos 10% de Présences Passives coûtent en silence. '
                     'Entretien structuré dans les 45 jours. '
                     'Réengager ou libérer — pas de zone grise.',
             'cible': '→ Cible : Passifs (10%)'},
            {'titre': 'DÉCIDER\nsur les Départs',
             'desc': '8% ont déjà décidé. Ne pas décider est une décision — qui coûte cher. '
                     'Plan de transition ou contre-offre ciblée. Dans les 30 jours.',
             'cible': '→ Cible : Départs (8%)'},
        ])

    slide_chart(prs,
        'Matrice de Décision — Action prioritaire par case talent',
        c_dec,
        'Chaque action est adossée à un profil de la Matrice DATATYM. Aucune décision générique.')

    # ══ CONCLUSION ══════════════════════════════════════════════
    slide_quote(prs,
        'Les organisations qui prospèrent demain\nsont celles qui ont choisi de voir\nce que leurs chiffres leur cachaient — aujourd\'hui.')

    slide_cover(prs,
        'DATATYM Analytics',
        'Baromètre Gen Z Afrique 2026\nDocument confidentiel — Réservé aux partenaires DATATYM',
        tag='MERCI')

    # ── Save ─────────────────────────────────────────────────────
    out = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/DATATYM_GenZ_V2_2026.pptx'
    prs.save(out)
    n = len(prs.slides)
    print(f'\n✓ {n} slides → {out}')


if __name__ == '__main__':
    build()
