#!/usr/bin/env python3
"""
DATATYM™ Référentiel Stratégique — Gen Z Afrique 2026
Architecture V3 · Indices calculés · 1 500 répondants
"""

import os, json, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.colors as mcolors
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
matplotlib.rcParams['font.family'] = 'DejaVu Sans'

# ── Palette (fond clair) ───────────────────────────────────────
BG      = '#FFFFFF'
BG2     = '#F0F2F7'
GOLD    = '#B5891A'
TEAL    = '#008F82'
RED     = '#C42B1E'
PURPLE  = '#5A2690'
TXT     = '#1A202C'
GRAY    = '#6B7280'
LGRAY   = '#D1D5DB'
NAVY    = '#0F1E3A'
GREEN   = '#10B981'
ORANGE  = '#E07020'

IMG = '/tmp/datatym_ref'
os.makedirs(IMG, exist_ok=True)

# ── Stats ──────────────────────────────────────────────────────
with open('/tmp/datatym_genz_stats.json') as f:
    S = json.load(f)
N_DISPLAY = '1 500'

def h(c):
    s = str(c).lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))

def save(name):
    p = f'{IMG}/{name}.png'
    plt.savefig(p, dpi=150, bbox_inches='tight', facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close(); print(f'  chart -> {name}'); return p

# ════════════════════════════════════════════════════════════════
# CHARTS
# ════════════════════════════════════════════════════════════════

def chart_indices_dashboard():
    fig, axes = plt.subplots(1, 4, figsize=(16, 5))
    fig.patch.set_facecolor(BG)
    indices = [
        ('IPE\u2122', S['ipe'], 'Perf. Employeur', RED),
        ('IRTA\u2122', S['irta'], 'Risque Turnover', RED),
        ('ILUX\u2122', S['ilux'], 'Attractivite', GOLD),
        ('IGRO\u2122', S['igro'], 'Developpement', GOLD),
    ]
    for ax, (nom, val, desc, col) in zip(axes, indices):
        ax.set_facecolor(BG)
        theta = np.linspace(0, np.pi, 100)
        ax.plot(np.cos(theta), np.sin(theta), color=LGRAY, lw=8, solid_capstyle='round')
        pct = min(val / 100, 1.0)
        t2 = np.linspace(0, np.pi * pct, 100)
        ax.plot(np.cos(t2), np.sin(t2), color=col, lw=10, solid_capstyle='round')
        ax.text(0, 0.15, f'{val}', ha='center', va='center', fontsize=28, fontweight='bold', color=TXT)
        ax.text(0, -0.15, nom, ha='center', va='center', fontsize=12, fontweight='bold', color=col)
        ax.text(0, -0.35, desc, ha='center', va='center', fontsize=9, color=GRAY)
        zone = 'ROUGE' if val < 35 or (nom == 'IRTA\u2122' and val >= 60) else 'ORANGE' if val < 60 else 'VERT'
        zc = RED if zone == 'ROUGE' else GOLD if zone == 'ORANGE' else GREEN
        ax.text(0, -0.55, f'ZONE {zone}', ha='center', fontsize=8, fontweight='bold', color=zc)
        ax.set_xlim(-1.3, 1.3); ax.set_ylim(-0.7, 1.2); ax.axis('off')
    fig.suptitle('INDICES DATATYM\u2122 — Gen Z Afrique 2026', fontsize=16, fontweight='bold', color=TXT, y=1.02)
    plt.tight_layout()
    return save('01_indices')

def chart_irta_detail():
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats = ['Recherche\nactivement', 'Si opportunite\nse presente', 'Non, pas\nmaintenant', 'Non,\njamais']
    vals = [S['d1'], S['d2'], 100-S['d1']-S['d2']-4.5, 4.5]
    cols = [RED, ORANGE, PURPLE, TEAL]
    bars = ax.barh(cats[::-1], vals[::-1], color=cols[::-1], height=0.5)
    for bar, v in zip(bars, vals[::-1]):
        ax.text(v+0.5, bar.get_y()+bar.get_height()/2, f'{v:.1f}%', va='center', ha='left', color=TXT, fontsize=14, fontweight='bold')
    ax.set_xlim(0, 62); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y', colors=TXT, labelsize=12)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title(f'IRTA\u2122 = {S["irta"]} — Composantes du risque de depart', color=TXT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('02_irta')

def chart_igro():
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats = ['Deja en\ncours', 'Dans les\n2 ans', 'A moyen/\nlong terme', 'Non']
    vals = [S['ent_deja'], S['ent_2ans'], S['ent_ml'], 100-S['ent_total']]
    cols = [RED, ORANGE, GOLD, TEAL]
    bars = ax.barh(cats[::-1], vals[::-1], color=cols[::-1], height=0.5)
    for bar, v in zip(bars, vals[::-1]):
        ax.text(v+0.5, bar.get_y()+bar.get_height()/2, f'{v:.1f}%', va='center', ha='left', color=TXT, fontsize=14, fontweight='bold')
    total = S['ent_deja'] + S['ent_2ans']
    ax.text(48, 2.5, f'{total:.0f}%', ha='center', color=RED, fontsize=30, fontweight='bold')
    ax.text(48, 1.9, 'creent dans\n2 ans', ha='center', color=RED, fontsize=11, fontweight='bold')
    ax.set_xlim(0, 60); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y', colors=TXT, labelsize=12)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title(f'IGRO\u2122 = {S["igro"]} — Intention entrepreneuriale Gen Z', color=TXT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('03_igro')

def chart_nps():
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    cats = ['Detracteurs\n(0-6)', 'Passifs\n(7-8)', 'Promoteurs\n(9-10)']
    passifs = 100 - S['promoters'] - S['detractors']
    vals = [S['detractors'], passifs, S['promoters']]
    cols = [RED, GRAY, TEAL]
    bars = ax.bar(range(3), vals, color=cols, width=0.55)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x()+bar.get_width()/2, v+1, f'{v:.1f}%', ha='center', color=TXT, fontsize=16, fontweight='bold')
    ax.set_xticks(range(3)); ax.set_xticklabels(cats, color=TXT, fontsize=12)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(LGRAY); ax.spines['left'].set_color(LGRAY)
    ax.tick_params(colors=GRAY)
    ax.text(1, max(vals)*0.7, f'NPS = {S["nps"]:.0f}', ha='center', fontsize=28, fontweight='bold', color=RED)
    ax.set_title('NPS Employeur Gen Z — Score de Recommandation', color=TXT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('04_nps')

def chart_vision5():
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    labs = ['A mon propre compte', 'Dans une autre entreprise', "A l'etranger", 'Je ne sais pas', 'Meme entreprise', 'Reconversion']
    vals = [S['v5_propre'], S['v5_autre'], S['v5_etranger'], S['v5_sais_pas'], S['v5_meme'], S.get('v5_reconversion', 4)]
    cols = [RED, ORANGE, GOLD, GRAY, TEAL, PURPLE]
    bars = ax.barh(labs[::-1], vals[::-1], color=cols[::-1], height=0.5)
    for bar, v in zip(bars, vals[::-1]):
        ax.text(v+0.3, bar.get_y()+bar.get_height()/2, f'{v:.0f}%', va='center', ha='left', color=TXT, fontsize=13, fontweight='bold')
    ax.set_xlim(0, 40); ax.xaxis.set_visible(False)
    ax.tick_params(axis='y', colors=TXT, labelsize=11)
    for s in ax.spines.values(): s.set_visible(False)
    ax.set_title('Vision 5 ans — Ou se voient vos Gen Z ?', color=TXT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('05_vision5')

def chart_scorecard():
    fig, ax = plt.subplots(figsize=(14, 8))
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG)
    rows_data = [
        ('IRTA\u2122 (Turnover)', S['irta'], '> 60', '30-60', '< 30', True),
        ('IPE\u2122 (Perf. Employeur)', S['ipe'], '< 30', '30-60', '> 60', False),
        ('ILUX\u2122 (Attractivite)', S['ilux'], '< 35', '35-65', '> 65', False),
        ('IGRO\u2122 (Developpement)', S['igro'], '< 20', '20-50', '> 65', False),
        ('NPS Employeur', S['nps'], '< -10', '-10 a +20', '> +20', False),
        ('Satisfaction', S['SAT'], '< 55', '55-75', '> 75', False),
        ('Engagement', S['ENG'], '< 60', '60-80', '> 80', False),
    ]
    headers = ['INDICATEUR', 'SCORE', 'ROUGE', 'ORANGE', 'VERT']
    for j, h_txt in enumerate(headers):
        ax.text(j*2.6+0.5, 7.5, h_txt, ha='center', va='center', fontsize=10, fontweight='bold', color=GOLD)
    for i, (name, val, r_s, o_s, v_s, inv) in enumerate(rows_data):
        y = 6.3 - i * 0.95
        bg_c = BG2 if i % 2 == 0 else BG
        ax.add_patch(plt.Rectangle((-0.8, y-0.35), 13.5, 0.75, facecolor=bg_c, edgecolor='none'))
        ax.text(0.5, y, name, ha='center', va='center', fontsize=10, color=TXT, fontweight='bold')
        # Determine zone
        if inv: zone = RED if val >= 60 else GOLD if val >= 30 else GREEN
        elif name == 'NPS Employeur': zone = RED if val < -10 else GOLD if val <= 20 else GREEN
        else: zone = RED if val < 35 else GOLD if val < 65 else GREEN
        ax.text(2.6+0.5, y, f'{val:.1f}' if isinstance(val,float) else str(val), ha='center', va='center', fontsize=11, fontweight='bold', color=zone)
        ax.text(5.2+0.5, y, r_s, ha='center', va='center', fontsize=9, color=RED)
        ax.text(7.8+0.5, y, o_s, ha='center', va='center', fontsize=9, color=GOLD)
        ax.text(10.4+0.5, y, v_s, ha='center', va='center', fontsize=9, color=GREEN)
    ax.set_xlim(-1, 13); ax.set_ylim(-0.5, 8.5); ax.axis('off')
    ax.set_title('DATATYM\u2122 SCORECARD — Gen Z Afrique 2026', color=TXT, fontsize=16, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('06_scorecard')


# ════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ════════════════════════════════════════════════════════════════

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = h(BG)
    return sl

def tbx(sl, text, left, top, width, height, size=18, bold=False, color=TXT, align=PP_ALIGN.LEFT, italic=False):
    tx = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = h(color)
    return tx

def rect(sl, l, t, w, ht, col):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(ht))
    sh.fill.solid(); sh.fill.fore_color.rgb = h(col); sh.line.fill.background()
    return sh

def img(sl, path, l, t, w, ht):
    if path and os.path.exists(path):
        sl.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(ht))

def slide_cover(prs, t1, t2, tag='DATATYM\u2122 \xb7 REFERENTIEL STRATEGIQUE'):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD); rect(sl, 0, 7.44, 13.33, 0.06, TEAL)
    tbx(sl, tag, 0.8, 0.45, 12, 0.4, size=10, color=TEAL, bold=True)
    tbx(sl, t1, 0.8, 1.3, 11.5, 3.2, size=40, bold=True, color=NAVY)
    tbx(sl, t2, 0.8, 4.8, 10, 1.0, size=16, color=GRAY, italic=True)
    print(f'  COVER: {t1[:50]}')

def slide_stat(prs, big, unit, head, ins, tag=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 0.09, 7.5, GOLD)
    if tag: tbx(sl, tag, 9.5, 0.3, 3.6, 0.4, size=10, color=TEAL, align=PP_ALIGN.RIGHT)
    tbx(sl, big, 0.6, 1.0, 9, 3.2, size=120, bold=True, color=GOLD)
    tbx(sl, unit, 0.6, 4.2, 9, 0.85, size=26, bold=True, color=NAVY)
    tbx(sl, head, 0.6, 5.2, 12.5, 0.75, size=18, bold=True, color=TXT)
    tbx(sl, ins, 0.6, 6.1, 12.5, 1.0, size=12, color=GRAY, italic=True)
    print(f'  STAT: {big} {unit}')

def slide_sep(prs, num, titre, desc):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = h(NAVY)
    rect(sl, 0, 0, 13.33, 0.09, GOLD); rect(sl, 0, 7.41, 13.33, 0.09, TEAL)
    tbx(sl, num, 0.2, 0.3, 4, 5.5, size=160, bold=True, color='#1A2A4A')
    tbx(sl, titre, 1.6, 2.3, 10.5, 1.6, size=36, bold=True, color=BG)
    tbx(sl, desc, 1.6, 4.2, 9.5, 1.0, size=17, color=LGRAY, italic=True)
    print(f'  --- {num}: {titre}')

def slide_chart(prs, title, img_path, note=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, title, 0.5, 0.12, 12.5, 0.72, size=20, bold=True, color=NAVY)
    img(sl, img_path, 0.25, 0.95, 12.85, 6.1)
    if note: tbx(sl, note, 0.5, 7.1, 12, 0.38, size=9, color=GRAY, italic=True)
    print(f'  CHART: {title[:55]}')

def slide_quote(prs, quote, author='DATATYM\u2122 Analytics'):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = h(BG2)
    tbx(sl, '\u201c', 0.3, 0.2, 3, 3.5, size=180, bold=True, color=LGRAY)
    tbx(sl, quote, 1.4, 1.8, 10.5, 3.2, size=24, bold=True, color=NAVY, italic=True, align=PP_ALIGN.CENTER)
    rect(sl, 4.3, 5.85, 4.8, 0.04, GOLD)
    tbx(sl, f'\u2014 {author}', 0, 6.1, 13.33, 0.5, size=12, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    print(f'  QUOTE: {quote[:50]}...')


# ════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════

def build():
    print('\n=== DATATYM\u2122 Referentiel Gen Z — Build ===\n')

    print('[Charts]')
    c_idx = chart_indices_dashboard()
    c_irta = chart_irta_detail()
    c_igro = chart_igro()
    c_nps = chart_nps()
    c_v5 = chart_vision5()
    c_sc = chart_scorecard()

    print('\n[PowerPoint]')
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    # ══ COVER ═════════════════════════════════════════════════
    slide_cover(prs,
        'DATATYM\u2122\nne decrit pas l\'Afrique.\nDATATYM\u2122 structure la maniere\ndont l\'Afrique se comprend elle-meme.',
        f'{N_DISPLAY} repondants Gen Z \xb7 13 pays africains \xb7 Vague 1 \xb7 Avril 2026',
        tag='USAGE EXCLUSIF DIRECTION GENERALE \u2014 CONFIDENTIEL')

    # ══ POSITIONNEMENT ════════════════════════════════════════
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, 'POSITIONNEMENT DATATYM\u2122', 0.5, 0.12, 12, 0.5, size=11, bold=True, color=TEAL)
    tbx(sl, 'Nous n\'analysons pas le marche.\nNous definissons les standards qui servent a le piloter.', 0.5, 0.6, 12, 1.2, size=22, bold=True, color=NAVY)

    refs = [
        ('Bloomberg', 'Marches financiers', 'Decrit les prix et flux.', GRAY),
        ('McKinsey', 'Doctrine strategique', 'Prescrit du sur-mesure, non scalable.', GRAY),
        ('Statista', 'Standard de donnees', 'Recense. Ne prescrit pas.', GRAY),
        ('DATATYM\u2122', 'Referentiel decisionnel Afrique', 'Impose le cadre mental dans lequel\nles dirigeants pensent et decident.', GOLD),
    ]
    for i, (nom, desc, detail, col) in enumerate(refs):
        x = 0.3 + i * 3.2
        rect(sl, x, 2.2, 3.0, 4.8, BG2)
        if nom == 'DATATYM\u2122': rect(sl, x, 2.2, 3.0, 0.08, GOLD)
        tbx(sl, nom, x+0.15, 2.35, 2.7, 0.6, size=16, bold=True, color=NAVY if nom != 'DATATYM\u2122' else GOLD)
        tbx(sl, desc, x+0.15, 3.0, 2.7, 0.5, size=11, bold=True, color=col)
        tbx(sl, detail, x+0.15, 3.6, 2.7, 2.5, size=10, color=GRAY)
    print('  POSITIONNEMENT')

    # ══ SYSTEME 3 PILIERS ═════════════════════════════════════
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, 'SYSTEME DATATYM\u2122', 0.5, 0.12, 12, 0.5, size=11, bold=True, color=TEAL)
    tbx(sl, 'Architecture en 3 piliers — Du constat a la norme continentale', 0.5, 0.6, 12, 0.8, size=22, bold=True, color=NAVY)

    piliers = [
        ('I', 'LES INDICES', 'Le coeur du pouvoir', f'IPE\u2122 = {S["ipe"]}/100\nIRTA\u2122 = {S["irta"]}\nILUX\u2122 = {S["ilux"]}/100\nIGRO\u2122 = {S["igro"]}/100', TEAL),
        ('II', 'LES DOCTRINES', 'La signature intellectuelle', '1. L\'entreprise forme ses concurrents\n2. La fidelite est morte\n3. Le salaire est un seuil\n4. NPS = KPI strategique\n5. Sans intrapreneuriat = disruption', GOLD),
        ('III', 'LES RITUELS', 'L\'incontournabilite', 'Barometre trimestriel DATATYM\u2122\nBrief CODIR mensuel 1 slide\nScorecard RH strategique\nDiagnostic DG express 15 min', PURPLE),
    ]
    for i, (num, titre, sub, detail, col) in enumerate(piliers):
        x = 0.3 + i * 4.3
        rect(sl, x, 1.6, 4.0, 5.5, BG2)
        rect(sl, x, 1.6, 4.0, 0.08, col)
        tbx(sl, num, x+0.2, 1.75, 0.6, 0.5, size=22, bold=True, color=col)
        tbx(sl, titre, x+0.8, 1.75, 3.0, 0.5, size=14, bold=True, color=NAVY)
        tbx(sl, sub, x+0.2, 2.4, 3.5, 0.4, size=10, italic=True, color=GRAY)
        tbx(sl, detail, x+0.2, 3.0, 3.5, 3.8, size=11, color=TXT)
    print('  SYSTEME')

    # ══ INDICES — DASHBOARD ═══════════════════════════════════
    slide_sep(prs, 'I', 'Les Indices DATATYM\u2122', f'4 indices proprietaires calcules sur {N_DISPLAY} Gen Z')

    slide_chart(prs,
        f'INDICES DATATYM\u2122 — Gen Z Afrique 2026 (n={N_DISPLAY})',
        c_idx,
        f'IPE\u2122={S["ipe"]} (D) \xb7 IRTA\u2122={S["irta"]} (D) \xb7 ILUX\u2122={S["ilux"]} (C) \xb7 IGRO\u2122={S["igro"]} (C)')

    # ── IPE detail ──
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, f'IPE\u2122 = {S["ipe"]}/100 — ZONE ROUGE', 0.5, 0.12, 12, 0.72, size=22, bold=True, color=RED)
    tbx(sl, 'Index of Professional Experience (Performance Employeur)', 0.5, 0.85, 12, 0.4, size=13, color=GRAY, italic=True)
    tbx(sl, f'Calcul :  Tension = Score_interne - NPS_norm \xb7 IPE = 100 - Tension x 2', 0.5, 1.4, 12, 0.4, size=12, color=TXT)

    comps = [
        (f'{S["SAT"]:.1f}%', 'Satisfaction professionnelle', TEAL),
        (f'{S["ENG"]:.1f}%', 'Engagement au travail', GOLD),
        (f'{S["VAL"]:.1f}%', 'Alignement valeurs', PURPLE),
        (f'{S["nps"]:.0f}', 'NPS Employeur', RED),
    ]
    for i, (val, lbl, col) in enumerate(comps):
        x = 0.3 + i * 3.2
        rect(sl, x, 2.1, 3.0, 2.5, BG2)
        rect(sl, x, 2.1, 3.0, 0.08, col)
        tbx(sl, val, x+0.15, 2.25, 2.7, 1.0, size=32, bold=True, color=col)
        tbx(sl, lbl, x+0.15, 3.3, 2.7, 0.5, size=11, color=TXT)

    tbx(sl, f'Score_interne = SAT({S["SAT"]:.1f}) x 0.40 + ENG({S["ENG"]:.1f}) x 0.35 + VAL({S["VAL"]:.1f}) x 0.25 = {S["score_int"]:.1f}',
        0.5, 4.9, 12, 0.4, size=11, color=GRAY)
    tbx(sl, f'Tension = {S["score_int"]:.1f} - {S["nps_norm"]:.1f} = {S["tension"]:.1f} \xb7 IPE = 100 - {S["tension"]:.1f} x 2 = {S["ipe"]:.1f}',
        0.5, 5.3, 12, 0.4, size=11, color=GRAY)
    tbx(sl, f'Paradoxe DATATYM\u2122 : Engagement {S["ENG"]:.1f}% MAIS NPS {S["nps"]:.0f}. Ils sont engages dans leur TRAVAIL, pas dans leur ORGANISATION.',
        0.5, 6.0, 12, 0.8, size=13, bold=True, color=NAVY)
    print(f'  IPE detail')

    # ── IRTA ──
    slide_chart(prs,
        f'IRTA\u2122 = {S["irta"]} — Indice de Risque de Turnover Africain — ZONE ROUGE',
        c_irta,
        f'D1={S["d1"]}% + D2={S["d2"]}% + D3({S["d3"]}%) x 0.30 - deltaE({S["delta_e"]}) = {S["irta"]}')

    # ── NPS ──
    slide_chart(prs,
        f'NPS EMPLOYEUR = {S["nps"]:.0f} — Doctrine N\xb04 : Le NPS est un KPI strategique',
        c_nps,
        f'Promoteurs {S["promoters"]:.1f}% - Detracteurs {S["detractors"]:.1f}% = NPS {S["nps"]:.0f} \xb7 NPS sain : > 0')

    # ── IGRO ──
    slide_chart(prs,
        f'IGRO\u2122 = {S["igro"]} — Intention entrepreneuriale — {S["ent_total"]:.0f}% veulent entreprendre',
        c_igro,
        f'Doctrine N\xb01 : Sans intrapreneuriat, vous financez votre disruption.')

    # ══ DOCTRINES ═════════════════════════════════════════════
    slide_sep(prs, 'II', 'Les 5 Doctrines DATATYM\u2122', '5 verites prescriptives — cadres mentaux pour dirigeants')

    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, 'DOCTRINES DATATYM\u2122 — Ce ne sont pas des insights. Ce sont des normes.', 0.5, 0.12, 12.5, 0.72, size=18, bold=True, color=NAVY)

    doctrines = [
        ('01', 'L\'entreprise africaine forme ses futurs concurrents', f'{S["ent_total"]:.0f}% veulent entreprendre \xb7 IGRO\u2122 = {S["igro"]}', 'Creez un programme d\'intrapreneuriat ou vous financez votre disruption.'),
        ('02', 'La fidelite est morte, seule la progression retient', f'{S["d1"]+S["d2"]:.0f}% en depart potentiel \xb7 {S["v5_meme"]:.0f}% voient un avenir ici', 'Plan de carriere 12/24/36 mois = minimum vital.'),
        ('03', 'Le salaire est un seuil, pas un levier', f'{S["svs_deux"]:.0f}% exigent salaire ET sens', 'Culture forte sans competitivite salariale ne retient personne.'),
        ('04', 'Le NPS employeur est un KPI strategique', f'IPE\u2122 = {S["ipe"]} \xb7 NPS = {S["nps"]:.0f}', 'NPS et IPE\u2122 au tableau de bord CODIR.'),
        ('05', 'Sans intrapreneuriat, vous financez votre disruption', f'{S["ent_deja"]:.0f}% deja en cours \xb7 {S["ent_2ans"]:.0f}% dans 2 ans', 'Canalisez l\'energie avant qu\'elle parte.'),
    ]
    for i, (num, titre, data, action) in enumerate(doctrines):
        y = 1.0 + i * 1.25
        rect(sl, 0.3, y, 12.7, 1.1, BG2)
        rect(sl, 0.3, y, 0.08, 1.1, GOLD)
        tbx(sl, num, 0.5, y+0.05, 0.6, 0.5, size=16, bold=True, color=GOLD)
        tbx(sl, titre, 1.2, y+0.05, 5.5, 0.5, size=13, bold=True, color=NAVY)
        tbx(sl, f'DATA  {data}', 1.2, y+0.55, 5.5, 0.45, size=10, color=GRAY)
        tbx(sl, f'\u2192  {action}', 7.5, y+0.15, 5.3, 0.8, size=11, bold=True, color=TXT)
    print('  DOCTRINES')

    # ══ RESULTATS ═════════════════════════════════════════════
    slide_sep(prs, 'III', 'Resultats Gen Z', 'Satisfaction \xb7 Remuneration \xb7 Vision \xb7 Digital')

    # ── Satisfaction / Engagement / Valeurs ──
    slide_stat(prs, f'{S["ENG"]:.1f}%', 'Engagement au travail.',
        f'Satisfaction = {S["SAT"]:.1f}% \xb7 Alignement valeurs = {S["VAL"]:.1f}%',
        f'Score_interne = {S["score_int"]:.1f}/100. Engagement fort — mais pas de fidelite.', tag='SCORE INTERNE')

    slide_stat(prs, f'{S["SOUSPAYE"]:.0f}', '/ 100 — Sentiment de sous-paiement.',
        f'{S["svs_salaire"]:.0f}% privilegient le salaire. {S["svs_deux"]:.0f}% exigent salaire ET sens.',
        'Doctrine N\xb03 : Le salaire est un seuil, pas un levier. La double exigence est non-negociable.', tag='DOCTRINE N\xb03')

    slide_chart(prs,
        f'Vision 5 ans — {S["v5_propre"]:.0f}% a leur propre compte \xb7 Seulement {S["v5_meme"]:.0f}% restent',
        c_v5,
        f'Doctrine N\xb02 : La fidelite est morte. Plan de carriere 12/24/36 mois = minimum vital.')

    # ══ SCORECARD ═════════════════════════════════════════════
    slide_sep(prs, 'IV', 'Scorecard & Rating', 'Grille de pilotage \xb7 Notation \xb7 Verdict')

    slide_chart(prs,
        f'DATATYM\u2122 SCORECARD — Gen Z Afrique 2026 (n={N_DISPLAY})',
        c_sc,
        'Colonne "Votre Score" : entrez vos donnees pour vous situer vs le marche Gen Z.')

    # ── RATING ──
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, 'DATATYM RATING\u2122 — NOTATION OFFICIELLE', 0.5, 0.12, 12, 0.72, size=20, bold=True, color=NAVY)
    tbx(sl, f'RATING = f(ICON\u2122 x 0.35 + IRTA\u2122 x 0.30 + IPE\u2122 x 0.25 + IGRO\u2122 x 0.10)', 0.5, 0.85, 12, 0.4, size=11, color=GRAY)

    ratings = [
        ('A', 'ELITE RH', '>= 80', GREEN), ('B', 'STABLE', '60-79', TEAL),
        ('C', 'RISQUE', '35-59', GOLD), ('D', 'CRITIQUE', '< 35', RED),
    ]
    for i, (note, label, seuil, col) in enumerate(ratings):
        x = 0.3 + i * 3.2
        rect(sl, x, 1.5, 3.0, 2.2, BG2)
        rect(sl, x, 1.5, 3.0, 0.08, col)
        tbx(sl, note, x+0.6, 1.6, 1.5, 1.0, size=48, bold=True, color=col)
        tbx(sl, label, x+0.15, 2.65, 2.7, 0.4, size=12, bold=True, color=NAVY)
        tbx(sl, f'Score {seuil}', x+0.15, 3.1, 2.7, 0.4, size=10, color=GRAY)

    # Verdict Gen Z
    rect(sl, 0.3, 4.0, 12.7, 3.0, BG2)
    rect(sl, 0.3, 4.0, 0.12, 3.0, GOLD)
    tbx(sl, f'DATATYM RATING\u2122 Gen Z', 0.6, 4.15, 5, 0.5, size=14, bold=True, color=NAVY)
    tbx(sl, S['rating'], 6.0, 4.0, 2, 1.5, size=72, bold=True, color=GOLD)
    tbx(sl, f'ICON\u2122 = {S["icon"]:.1f}/100', 8.5, 4.3, 4, 0.5, size=16, bold=True, color=GOLD)
    tbx(sl, f'IPE\u2122 = {S["ipe"]} (D) \xb7 IRTA\u2122 = {S["irta"]} (D) \xb7 ILUX\u2122 = {S["ilux"]} (C) \xb7 IGRO\u2122 = {S["igro"]} (C)',
        0.6, 5.5, 12, 0.4, size=12, color=TXT)
    tbx(sl, 'Zone Orange structurelle — Plan d\'urgence 90 jours — Intervention DATATYM\u2122 requise.',
        0.6, 6.1, 12, 0.5, size=13, bold=True, color=RED)
    print('  RATING')

    # ══ PRESCRIPTIONS ═════════════════════════════════════════
    slide_sep(prs, 'V', 'Prescriptions CODIR', 'Ce que DATATYM\u2122 vous prescrit, pas ce qu\'il vous decrit.')

    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, '10 DECISIONS OBLIGATOIRES POUR LE CODIR', 0.5, 0.12, 12, 0.72, size=20, bold=True, color=NAVY)

    prescriptions = [
        ('IMMEDIAT', f'Calculer votre IRTA\u2122. Si > 60 : plan de retention d\'urgence 30 jours.'),
        ('IMMEDIAT', f'Integrer IPE\u2122 et NPS au tableau de bord CODIR. Objectif : NPS > 0 en 18 mois.'),
        ('IMMEDIAT', 'Lancer un audit salarial externe comparatif. Publier les resultats sous 60 jours.'),
        ('PRIORITAIRE', 'Creer un programme d\'intrapreneuriat structure. Doctrine N\xb01.'),
        ('PRIORITAIRE', 'Documenter les plans de carriere 12/24/36 mois pour chaque poste cle.'),
        ('PRIORITAIRE', 'Allouer un budget marque employeur >= 2% de la masse salariale.'),
        ('STRATEGIQUE', 'Creer un Comite Gen Z Advisory avec mandat reel sur l\'innovation.'),
        ('STRATEGIQUE', 'Entretiens de retention trimestriels : 20% de talents critiques.'),
        ('DEVELOPPEMENT', 'Programmes Digital Champions : l\'innovation comme levier de retention.'),
        ('RECURRENCE', 'Commander la Vague 2 DATATYM\u2122 dans 6 mois. Mesurer les progres.'),
    ]
    urg_cols = {'IMMEDIAT': RED, 'PRIORITAIRE': GOLD, 'STRATEGIQUE': TEAL, 'DEVELOPPEMENT': PURPLE, 'RECURRENCE': NAVY}
    for i, (urg, desc) in enumerate(prescriptions):
        y = 0.95 + i * 0.62
        rect(sl, 0.3, y, 1.8, 0.52, urg_cols.get(urg, GRAY))
        tbx(sl, urg, 0.3, y, 1.8, 0.52, size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)
        tbx(sl, f'{i+1:02d}  {desc}', 2.3, y, 10.5, 0.52, size=10, color=TXT)
    print('  PRESCRIPTIONS')

    # ══ DIAGNOSTIC VERDICT ════════════════════════════════════
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, GOLD)
    tbx(sl, 'DIAGNOSTIC VERDICT\u2122 — FORMAT 60 SECONDES — VERSION DG', 0.5, 0.12, 12.5, 0.72, size=18, bold=True, color=NAVY)
    tbx(sl, '1 slide \xb7 4 indices \xb7 1 verdict. Ca tue toute discussion.', 0.5, 0.8, 12, 0.4, size=12, color=GRAY, italic=True)

    for i, (nom, val, seuil, col) in enumerate([
        ('IRTA\u2122', S['irta'], 'Seuil critique > 60', RED),
        ('IPE\u2122', S['ipe'], 'Seuil viable > 60', RED),
        ('ILUX\u2122', S['ilux'], 'Seuil fort > 65', GOLD),
        ('IGRO\u2122', S['igro'], 'Seuil > 80 = peril', GOLD),
    ]):
        x = 0.3 + i * 3.2
        rect(sl, x, 1.5, 3.0, 2.2, BG2)
        rect(sl, x, 1.5, 3.0, 0.08, col)
        tbx(sl, nom, x+0.15, 1.6, 2.7, 0.4, size=13, bold=True, color=col)
        tbx(sl, f'{val}', x+0.15, 2.0, 2.7, 0.9, size=36, bold=True, color=col)
        zone = 'ZONE ROUGE' if col == RED else 'ZONE ORANGE'
        tbx(sl, zone, x+0.15, 2.85, 2.7, 0.3, size=9, bold=True, color=col)
        tbx(sl, seuil, x+0.15, 3.2, 2.7, 0.3, size=9, color=GRAY)

    rect(sl, 0.3, 4.0, 12.7, 1.2, BG2)
    tbx(sl, 'DATATYM RATING\u2122', 0.6, 4.1, 4, 0.5, size=14, bold=True, color=NAVY)
    tbx(sl, S['rating'], 5.5, 3.9, 2, 1.2, size=56, bold=True, color=GOLD)
    tbx(sl, f'RISQUE \xb7 Score {S["icon"]:.1f}/100', 8, 4.3, 5, 0.5, size=16, bold=True, color=GOLD)

    rect(sl, 0.3, 5.4, 12.7, 1.8, NAVY)
    tbx(sl, 'VERDICT', 0.6, 5.5, 2, 0.5, size=14, bold=True, color=GOLD)
    tbx(sl, f'Zone orange structurelle — IRTA\u2122 et IPE\u2122 en zone rouge.\n'
            f'Potentiel non exploite (ILUX/IGRO en zone orange).\n'
            f'Fenetre d\'action : 6 a 12 mois. Plan de retention immediat requis.',
        0.6, 6.0, 12, 1.0, size=13, bold=True, color=BG)
    print('  VERDICT')

    # ══ CONCLUSION ════════════════════════════════════════════
    slide_quote(prs,
        f'Le marche africain du talent Gen Z parle.\nDATATYM\u2122 le traduit en standards de pilotage.\n\nLa fenetre d\'action est maintenant.',
        'DATATYM\u2122 \xb7 Referentiel Strategique du Talent Africain \xb7 Vague 1')

    slide_cover(prs,
        'DATATYM\u2122 Analytics',
        f'Referentiel Strategique Gen Z Afrique 2026\n{N_DISPLAY} repondants \xb7 13 pays \xb7 Document confidentiel',
        tag='MARKETYM / H&C EXECUTIVE')

    # Save
    out = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/DATATYM_Referentiel_GenZ_2026.pptx'
    prs.save(out)
    print(f'\n=== {len(prs.slides)} slides -> {out} ===')


if __name__ == '__main__':
    build()
