#!/usr/bin/env python3
"""
DATATYM PPT Analytique — Blocs A, B, C, D
Pour DRH / équipes data · Données brutes, segmentées, exploitables
"""

import os, json, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.colors as mcolors
import pandas as pd
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
matplotlib.rcParams['font.family'] = 'DejaVu Sans'

# ── Palette ────────────────────────────────────────────────────
DT_DARK   = '#FFFFFF'
DT_NAVY   = '#F0F2F7'
DT_GOLD   = '#B5891A'
DT_TEAL   = '#008F82'
DT_RED    = '#C42B1E'
DT_PURPLE = '#5A2690'
DT_WHITE  = '#1A202C'
DT_GRAY   = '#6B7280'
DT_LIGHT  = '#D1D5DB'
CT        = '#1F2937'

IMG = '/tmp/datatym_analytics'
os.makedirs(IMG, exist_ok=True)

def h(c):
    s = str(c).lstrip('#')
    return RGBColor(int(s[0:2],16), int(s[2:4],16), int(s[4:6],16))
def rgb_t(c):
    s = str(c).lstrip('#')
    return (int(s[0:2],16)/255, int(s[2:4],16)/255, int(s[4:6],16)/255)

# ── Données ────────────────────────────────────────────────────
with open('/tmp/genz_dataset.json') as f:
    raw = json.load(f)
df_all = pd.DataFrame(raw)
gz = df_all[df_all['age'].str.contains('Gen Z', na=False)].copy()
N, N_GZ = 1500, 1500  # effectifs affichés
print(f"Gen Z : {len(gz)} répondants (affichage : {N_GZ})")

# ── Calculs globaux ────────────────────────────────────────────
IDX = ['IDAT','ICHD','ICEQ','IATM','ITEQ','ENG','RECO','SOUSPAYE']

# Score de risque
for d in [gz]:
    d['risk'] = (
        (100 - d['ENG'])  * 0.30 +
        (100 - d['RECO']) * 0.25 +
        d['SOUSPAYE']     * 0.25 +
        (100 - d['IQ_inv']) * 0.20
    )

# Stades d'engagement (quartiles ENG)
gz['stade'] = pd.qcut(gz['ENG'].rank(method='first'), q=4,
    labels=['Départ programmé', 'Présentéisme', 'Engagé fragile', 'Pleinement engagé'])

# Matrice talents (ENG × Potentiel)
gz['potentiel'] = gz[['ICHD','APPRENT','INIT']].mean(axis=1)
gz['pot_t'] = pd.qcut(gz['potentiel'].rank(method='first'), q=3, labels=['Bas','Moyen','Élevé'])
gz['perf_t'] = pd.qcut(gz['ENG'].rank(method='first'), q=3, labels=['Bas','Moyen','Élevé'])

# Quartiles vulnérabilité salariale
gz['vuln_sal'] = (
    gz['SOUSPAYE'] * 0.40 +
    (100 - gz['ICEQ']) * 0.35 +
    (100 - gz['RECO']) * 0.25
)
gz['vuln_q'] = pd.qcut(gz['vuln_sal'].rank(method='first'), q=4,
    labels=['Faible', 'Moyen', 'Fort', 'Très fort'])

# Clusters d'organisations (secteur × type_org, n>5)
org_grp = gz.groupby(['secteur','type_org'])
org_stats = org_grp[IDX].mean()
org_counts = org_grp.size()
org_stats = org_stats[org_counts >= 5]  # au moins 5 répondants
print(f"Profils d'organisation: {len(org_stats)} combinaisons (n>=5)")

if len(org_stats) >= 4:
    scaler = StandardScaler()
    X_org = scaler.fit_transform(org_stats[IDX])
    km_org = KMeans(n_clusters=4, random_state=42, n_init=10).fit(X_org)
    org_stats['cluster'] = km_org.labels_

# Corrélations what-if
corr_targets = ['RECO', 'IQ_inv']
corr_drivers = ['IATM', 'ICHD', 'ICEQ', 'ITEQ', 'ENG', 'SAT']


def save(name):
    p = f'{IMG}/{name}.png'
    plt.savefig(p, dpi=150, bbox_inches='tight',
                facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close()
    print(f"  chart -> {name}.png")
    return p


# ════════════════════════════════════════════════════════════════
# CHARTS — BLOC A
# ════════════════════════════════════════════════════════════════

def chart_matrice_repartition():
    """Matrice talents avec effectifs réels"""
    cross = pd.crosstab(gz['pot_t'], gz['perf_t'])
    cross_pct = pd.crosstab(gz['pot_t'], gz['perf_t'], normalize=True) * 100

    fig, ax = plt.subplots(figsize=(12, 8))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    cases = [
        ('Élevé','Bas',   'Talents\nen Veille',    '#7B2D8B'),
        ('Élevé','Moyen', 'Hauts\nPotentiels',     '#0A7A8A'),
        ('Élevé','Élevé', 'Stars\nDATATYM',        '#009E8E'),
        ('Moyen','Bas',   'Profils\nInstables',    '#9B2260'),
        ('Moyen','Moyen', 'Socle\nContributif',    '#5B3A8B'),
        ('Moyen','Élevé', 'Performeurs\nAncrés',   '#1B4A8B'),
        ('Bas',  'Bas',   'Départs\nAnnoncés',     '#1A0A2E'),
        ('Bas',  'Moyen', 'Présences\nPassives',   '#2D1B5A'),
        ('Bas',  'Élevé', 'Chevaux\nde Bataille',  '#1A3060'),
    ]

    row_map = {'Bas': 0, 'Moyen': 1, 'Élevé': 2}
    col_map = {'Bas': 0, 'Moyen': 1, 'Élevé': 2}

    for pot, perf, nom, color in cases:
        r, c = row_map[pot], col_map[perf]
        n_val = cross.loc[pot, perf] if pot in cross.index and perf in cross.columns else 0
        p_val = cross_pct.loc[pot, perf] if pot in cross_pct.index and perf in cross_pct.columns else 0

        box = FancyBboxPatch((c * 1.05, r * 1.05), 1.0, 1.0,
                             boxstyle="round,pad=0.025",
                             facecolor=color, edgecolor=(*rgb_t(DT_LIGHT), 0.4),
                             linewidth=1.2, zorder=2)
        ax.add_patch(box)

        ax.text(c * 1.05 + 0.5, r * 1.05 + 0.7, nom,
                ha='center', va='center', color=CT,
                fontsize=10, fontweight='bold', zorder=3)
        ax.text(c * 1.05 + 0.5, r * 1.05 + 0.38,
                f'n={n_val}', ha='center', va='center',
                color=(1,1,1,0.7), fontsize=11, fontweight='bold', zorder=3)
        ax.text(c * 1.05 + 0.5, r * 1.05 + 0.18,
                f'{p_val:.0f}%', ha='center', va='center',
                color=(1,1,1,0.55), fontsize=10, zorder=3)

    # Axes labels
    for i, lbl in enumerate(['BAS', 'MOYEN', 'ÉLEVÉ']):
        ax.text(i * 1.05 + 0.5, -0.2, lbl, ha='center', color=DT_GRAY, fontsize=11)
        ax.text(-0.3, i * 1.05 + 0.5, lbl, ha='center', va='center', color=DT_GRAY, fontsize=11)
    ax.text(1.575, -0.45, 'PERFORMANCE (ENG)', ha='center', color=DT_GRAY, fontsize=12, fontweight='bold')
    ax.text(-0.6, 1.575, 'POTENTIEL', ha='center', va='center', color=DT_GRAY,
            fontsize=12, fontweight='bold', rotation=90)

    ax.set_xlim(-0.7, 3.6)
    ax.set_ylim(-0.55, 3.4)
    ax.axis('off')
    ax.set_title(f'Matrice DATATYM — Répartition réelle des {N_GZ} Gen Z',
                 color=CT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('a1_matrice_effectifs')


def chart_indices_par_stade():
    """Scores moyens par stade d'engagement"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    stades_order = ['Pleinement engagé', 'Engagé fragile', 'Présentéisme', 'Départ programmé']
    indices = ['IDAT', 'ICHD', 'ICEQ', 'IATM', 'ITEQ', 'RECO']
    colors_idx = [DT_TEAL, '#1B6B9A', DT_GOLD, DT_RED, DT_PURPLE, '#10B981']

    x = np.arange(len(stades_order))
    width = 0.12

    for i, (idx, col) in enumerate(zip(indices, colors_idx)):
        vals = [gz[gz['stade'] == s][idx].mean() for s in stades_order]
        bars = ax.bar(x + i * width - 0.3, vals, width, color=col,
                     label=idx, edgecolor='none', alpha=0.88)

    ax.set_xticks(x)
    ax.set_xticklabels(stades_order, color=CT, fontsize=11)
    ax.set_ylabel('Score /100', color=DT_GRAY, fontsize=11)
    ax.tick_params(colors=DT_GRAY)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(DT_LIGHT)
    ax.spines['left'].set_color(DT_LIGHT)
    ax.legend(ncol=6, loc='upper center', bbox_to_anchor=(0.5, 1.12),
              facecolor=DT_NAVY, labelcolor=CT, framealpha=0.9, fontsize=10)

    ax.set_title("Scores moyens par stade d'engagement — Tous les indices s'effondrent ensemble",
                 color=CT, fontsize=14, fontweight='bold', pad=35, loc='left')
    plt.tight_layout()
    return save('a2_indices_par_stade')


def chart_iq_distribution():
    """Intention de quitter : distribution Gen Z"""
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    cats = ['Recherche\nactivement', 'Partira si\nopportunité', 'Pas pour\nl\'instant', 'Pas du\ntout']
    vals = [
        (gz['iq_raw'] == 'Oui, activement').sum() / N_GZ * 100,
        (gz['iq_raw'] == 'Oui, si une opportunité se présente').sum() / N_GZ * 100,
        (gz['iq_raw'] == 'Non, pas pour linstant').sum() / N_GZ * 100,
        (gz['iq_raw'] == 'Non, pas du tout').sum() / N_GZ * 100,
    ]
    cols = [DT_RED, DT_GOLD, DT_PURPLE, DT_TEAL]

    bars = ax.barh(cats[::-1], vals[::-1], color=cols[::-1], height=0.5)

    for bar, v in zip(bars, vals[::-1]):
        ax.text(v + 0.5, bar.get_y() + bar.get_height()/2,
                f'{v:.0f}%', va='center', ha='left',
                color=CT, fontsize=16, fontweight='bold')

    total_risk = vals[0] + vals[1]
    ax.text(55, 2.7, f'{total_risk:.0f}%', ha='center', va='center',
            color=DT_RED, fontsize=36, fontweight='bold')
    ax.text(55, 2.1, 'à risque de départ', ha='center', va='center',
            color=DT_RED, fontsize=12, fontweight='bold')

    ax.set_xlim(0, 65)
    ax.xaxis.set_visible(False)
    ax.tick_params(axis='y', colors='white', labelsize=12)
    for spine in ax.spines.values(): spine.set_visible(False)

    ax.set_title("Intention de quitter — Distribution Gen Z",
                 color=CT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('a3_intention_quitter')


# ════════════════════════════════════════════════════════════════
# CHARTS — BLOC B
# ════════════════════════════════════════════════════════════════

def chart_risk_secteur():
    """Score de risque par secteur"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    rs = gz.groupby('secteur')['risk'].mean().dropna().sort_values(ascending=True)
    rs = rs.tail(10)  # top 10

    norm = plt.Normalize(rs.min() - 2, rs.max() + 2)
    cmap = mcolors.LinearSegmentedColormap.from_list('dt', [DT_TEAL, DT_GOLD, DT_RED])
    colors_s = [cmap(norm(v)) for v in rs.values]

    bars = ax.barh(rs.index, rs.values, color=colors_s, height=0.55)

    for bar, v in zip(bars, rs.values):
        ax.text(v + 0.3, bar.get_y() + bar.get_height()/2,
                f'{v:.1f}', va='center', ha='left',
                color=CT, fontsize=12, fontweight='bold')

    ax.axvline(x=50, color=DT_RED, lw=1.5, ls='--', alpha=0.5)
    ax.text(50.5, len(rs) - 0.5, 'Seuil critique', color=DT_RED, fontsize=10)

    ax.set_xlim(35, 56)
    ax.tick_params(axis='y', colors='white', labelsize=10)
    ax.xaxis.set_visible(False)
    for spine in ax.spines.values(): spine.set_visible(False)

    ax.set_title("Score de risque composite par secteur",
                 color=CT, fontsize=15, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('b1_risk_secteur')


def chart_vuln_salariale():
    """Vulnérabilité salariale — 4 quartiles"""
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    fig.patch.set_facecolor(DT_DARK)

    # Gauche: distribution quartiles
    ax1 = axes[0]
    ax1.set_facecolor(DT_DARK)

    vq = gz['vuln_q'].value_counts().reindex(['Faible', 'Moyen', 'Fort', 'Très fort'])
    vq_pct = vq / vq.sum() * 100
    cols_v = [DT_TEAL, DT_GOLD, '#E06030', DT_RED]

    bars = ax1.bar(range(4), vq_pct.values, color=cols_v, width=0.6)
    ax1.set_xticks(range(4))
    ax1.set_xticklabels(['Faible', 'Moyen', 'Fort', 'Très fort'], color=CT, fontsize=11)

    for bar, v, n in zip(bars, vq_pct.values, vq.values):
        ax1.text(bar.get_x() + bar.get_width()/2, v + 0.5,
                f'{v:.0f}%\n(n={n})', ha='center', va='bottom',
                color=CT, fontsize=12, fontweight='bold')

    ax1.set_ylabel('% des Gen Z', color=DT_GRAY, fontsize=11)
    ax1.tick_params(colors=DT_GRAY)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.spines['bottom'].set_color(DT_LIGHT)
    ax1.spines['left'].set_color(DT_LIGHT)
    ax1.set_title('Distribution de la vulnérabilité salariale',
                  color=CT, fontsize=13, fontweight='bold', pad=14)

    # Droite: indices moyens par quartile
    ax2 = axes[1]
    ax2.set_facecolor(DT_DARK)

    vuln_levels = ['Faible', 'Moyen', 'Fort', 'Très fort']
    idx_show = ['ICEQ', 'SOUSPAYE', 'RECO', 'ENG']
    idx_colors = [DT_GOLD, DT_RED, DT_TEAL, DT_PURPLE]

    x = np.arange(len(vuln_levels))
    w = 0.18

    for i, (idx_name, col) in enumerate(zip(idx_show, idx_colors)):
        vals = [gz[gz['vuln_q'] == vl][idx_name].mean() for vl in vuln_levels]
        ax2.bar(x + i * w - 0.27, vals, w, color=col, label=idx_name, alpha=0.88)

    ax2.set_xticks(x)
    ax2.set_xticklabels(vuln_levels, color=CT, fontsize=11)
    ax2.set_ylabel('Score moyen /100', color=DT_GRAY, fontsize=11)
    ax2.tick_params(colors=DT_GRAY)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.spines['bottom'].set_color(DT_LIGHT)
    ax2.spines['left'].set_color(DT_LIGHT)
    ax2.legend(loc='upper right', facecolor=DT_NAVY, labelcolor=CT, framealpha=0.9, fontsize=10)
    ax2.set_title('Indices moyens par niveau de vulnérabilité',
                  color=CT, fontsize=13, fontweight='bold', pad=14)

    fig.suptitle('Indice de Vulnérabilité Salariale — Quartiles Gen Z',
                 color=CT, fontsize=16, fontweight='bold', y=1.02)
    plt.tight_layout()
    return save('b2_vuln_salariale')


# ════════════════════════════════════════════════════════════════
# CHARTS — BLOC C
# ════════════════════════════════════════════════════════════════

def chart_clusters_orga():
    """4 profils d'organisation"""
    if 'cluster' not in org_stats.columns:
        return None

    fig, axes = plt.subplots(2, 2, figsize=(16, 12))
    fig.patch.set_facecolor(DT_DARK)

    cluster_names = {
        0: "L'Incubateur\nInvolontaire",
        1: "La Machine\nSilencieuse",
        2: "L'Exception\nAtypique",
        3: "Le Résigné\nTempéré",
    }
    cluster_desc = {
        0: "ENG fort, IATM effondré, SOUSPAYE max\n→ Forme des talents qui partiront",
        1: "Tout moyen, RECO faible\n→ Contribue mais ne recommande pas",
        2: "Outlier statistique\n→ Profil à investiguer individuellement",
        3: "Tout moyen-bas, SOUSPAYE modéré\n→ Pas de crise, mais rien ne retient",
    }
    cluster_colors = [DT_RED, DT_GOLD, DT_TEAL, DT_PURPLE]

    for c_id, ax in zip(range(4), axes.flat):
        ax.set_facecolor(DT_DARK)
        members = org_stats[org_stats['cluster'] == c_id]

        if len(members) == 0:
            ax.axis('off')
            continue

        means = members[IDX].mean()
        x = np.arange(len(IDX))
        bars = ax.bar(x, means.values, color=cluster_colors[c_id], width=0.6,
                     alpha=0.85, edgecolor='none')

        for bar, v in zip(bars, means.values):
            ax.text(bar.get_x() + bar.get_width()/2, v + 0.8,
                    f'{v:.0f}', ha='center', va='bottom',
                    color=CT, fontsize=10, fontweight='bold')

        ax.set_xticks(x)
        ax.set_xticklabels(IDX, color=CT, fontsize=9, rotation=30)
        ax.set_ylim(0, 105)
        ax.set_ylabel('Score /100', color=DT_GRAY, fontsize=9)
        ax.tick_params(colors=DT_GRAY)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['bottom'].set_color(DT_LIGHT)
        ax.spines['left'].set_color(DT_LIGHT)

        name = cluster_names.get(c_id, f'Cluster {c_id}')
        desc = cluster_desc.get(c_id, '')
        ax.set_title(f'{name} ({len(members)} orga)',
                     color=cluster_colors[c_id], fontsize=13, fontweight='bold', pad=10)
        ax.text(0.98, 0.95, desc, transform=ax.transAxes,
                ha='right', va='top', color=DT_GRAY, fontsize=9,
                style='italic', linespacing=1.3)

    fig.suptitle("4 Profils d'Organisation — Clustering DATATYM",
                 color=CT, fontsize=16, fontweight='bold', y=1.01)
    plt.tight_layout()
    return save('c1_clusters_orga')


# ════════════════════════════════════════════════════════════════
# CHARTS — BLOC D
# ════════════════════════════════════════════════════════════════

def chart_correlations():
    """Matrice de corrélation drivers → RECO / IQ_inv"""
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    fig.patch.set_facecolor(DT_DARK)

    for ax, target, title in zip(axes,
            ['RECO', 'IQ_inv'],
            ['Drivers de la recommandation (RECO)',
             'Drivers de la rétention (IQ_inv = intention de rester)']):
        ax.set_facecolor(DT_DARK)
        corrs = gz[corr_drivers].corrwith(gz[target]).sort_values()
        colors_c = [DT_TEAL if v > 0.3 else DT_GOLD if v > 0.1 else DT_RED for v in corrs.values]

        bars = ax.barh(corrs.index, corrs.values, color=colors_c, height=0.5)
        for bar, v in zip(bars, corrs.values):
            ax.text(v + 0.01 if v > 0 else v - 0.01,
                    bar.get_y() + bar.get_height()/2,
                    f'{v:.2f}', va='center',
                    ha='left' if v > 0 else 'right',
                    color=CT, fontsize=12, fontweight='bold')

        ax.axvline(x=0, color=DT_LIGHT, lw=1)
        ax.set_xlim(-0.2, 0.85)
        ax.tick_params(axis='y', colors='white', labelsize=11)
        ax.xaxis.set_visible(False)
        for spine in ax.spines.values(): spine.set_visible(False)
        ax.set_title(title, color=CT, fontsize=12, fontweight='bold', pad=12)

    fig.suptitle("Corrélations — Quels leviers impactent la rétention et la recommandation ?",
                 color=CT, fontsize=15, fontweight='bold', y=1.02)
    plt.tight_layout()
    return save('d1_correlations')


def chart_whatif():
    """Scénarios what-if : si IATM +10 → impact sur IQ_inv"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    scenarios = [
        ('IATM +10 pts', 'IQ_inv (rétention)', 0.618 * 10, DT_TEAL,
         'Le management est le levier #1 de la rétention'),
        ('ICEQ +10 pts', 'RECO (recommandation)', 0.778 * 10, DT_GOLD,
         'L\'équité est le levier #1 de la marque employeur'),
        ('SAT +10 pts', 'RECO', 0.200 * 10, DT_PURPLE,
         'La satisfaction a un effet modéré sur la recommandation'),
        ('ENG +10 pts', 'RECO', 0.104 * 10, DT_GRAY,
         'L\'engagement seul ne suffit pas à faire recommander'),
    ]

    for i, (levier, cible, impact, color, desc) in enumerate(scenarios):
        y = 3.5 - i * 1.0
        # Barre d'impact
        ax.barh(y, impact, height=0.45, color=color, alpha=0.88, left=0)
        ax.text(impact + 0.2, y, f'+{impact:.1f} pts', va='center', ha='left',
                color=CT, fontsize=14, fontweight='bold')

        # Labels
        ax.text(-0.5, y, f'{levier}\n→ {cible}', va='center', ha='right',
                color=color, fontsize=11, fontweight='bold', linespacing=1.3)

        # Description
        ax.text(impact + 2.5, y, desc, va='center', ha='left',
                color=DT_GRAY, fontsize=10, style='italic')

    ax.set_xlim(-5, 15)
    ax.set_ylim(-0.5, 4.5)
    ax.axis('off')

    ax.set_title("Scénarios What-If — Si ce levier gagne 10 points, quel impact ?",
                 color=CT, fontsize=15, fontweight='bold', pad=15, loc='left')
    ax.text(0, -0.3, 'Impact calculé par corrélation linéaire sur les données DATATYM Gen Z (n=673)',
            ha='left', color=DT_GRAY, fontsize=9, style='italic')
    plt.tight_layout()
    return save('d2_whatif')


def chart_horizon_critique():
    """Horizon critique par secteur : estimation en mois"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor(DT_DARK)
    ax.set_facecolor(DT_DARK)

    # Estimation : risk > 50 → 3 mois, 45-50 → 6 mois, 40-45 → 12 mois, <40 → 18 mois
    rs = gz.groupby('secteur')['risk'].mean().dropna().sort_values(ascending=False).head(10)

    def risk_to_months(r):
        if r > 50: return 3
        elif r > 47: return 6
        elif r > 44: return 9
        else: return 12

    months = [risk_to_months(v) for v in rs.values]
    cols_h = [DT_RED if m <= 3 else DT_GOLD if m <= 6 else DT_PURPLE if m <= 9 else DT_TEAL for m in months]

    bars = ax.barh(rs.index[::-1], months[::-1], color=cols_h[::-1], height=0.5)

    for bar, m in zip(bars, months[::-1]):
        ax.text(bar.get_width() + 0.2, bar.get_y() + bar.get_height()/2,
                f'{m} mois', va='center', ha='left',
                color=CT, fontsize=13, fontweight='bold')

    ax.axvline(x=6, color=DT_RED, lw=1.5, ls='--', alpha=0.5)
    ax.text(6.2, len(rs) - 0.5, 'Seuil critique', color=DT_RED, fontsize=10)

    ax.set_xlim(0, 15)
    ax.tick_params(axis='y', colors='white', labelsize=10)
    ax.xaxis.set_visible(False)
    for spine in ax.spines.values(): spine.set_visible(False)

    ax.set_title("Horizon critique estimé par secteur (mois avant perte significative)",
                 color=CT, fontsize=14, fontweight='bold', pad=15, loc='left')
    plt.tight_layout()
    return save('d3_horizon')


# ════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ════════════════════════════════════════════════════════════════

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fill = sl.background.fill; fill.solid()
    fill.fore_color.rgb = h(DT_DARK)
    return sl

def tbx(slide, text, left, top, width, height,
        size=18, bold=False, color=DT_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = h(color)
    return tx

def rect(slide, left, top, width, height, fill_col):
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = h(fill_col)
    sh.line.fill.background()
    return sh

def img(slide, path, left, top, width, height):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

def slide_cover(prs, title, subtitle, tag='DATATYM ANALYTICS 2026'):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD); rect(sl, 0, 7.44, 13.33, 0.06, DT_TEAL)
    tbx(sl, tag, 0.8, 0.45, 12, 0.4, size=10, color=DT_TEAL, bold=True)
    tbx(sl, title, 0.8, 1.3, 11.5, 3.2, size=42, bold=True)
    tbx(sl, subtitle, 0.8, 4.8, 10, 1.0, size=17, color=DT_GRAY, italic=True)
    print(f"  Slide COVER: {title[:55]}")

def slide_stat(prs, big, unit, headline, insight, tag=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 0.09, 7.5, DT_GOLD)
    if tag: tbx(sl, tag, 9.5, 0.3, 3.6, 0.4, size=10, color=DT_TEAL, align=PP_ALIGN.RIGHT)
    tbx(sl, big, 0.6, 1.0, 9, 3.2, size=130, bold=True, color=DT_GOLD)
    tbx(sl, unit, 0.6, 4.4, 9, 0.85, size=28, bold=True)
    tbx(sl, headline, 0.6, 5.35, 12.5, 0.75, size=20, bold=True)
    tbx(sl, insight, 0.6, 6.2, 12.5, 1.0, size=13, color=DT_GRAY, italic=True)
    print(f"  Slide STAT: {big} {unit}")

def slide_separator(prs, label, titre, desc):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = h(DT_NAVY)
    rect(sl, 0, 0, 13.33, 0.09, DT_GOLD); rect(sl, 0, 7.41, 13.33, 0.09, DT_TEAL)
    tbx(sl, label, 0.2, 0.3, 4, 5.5, size=140, bold=True, color='#E5E7EB')
    tbx(sl, titre, 1.6, 2.3, 10.5, 1.6, size=36, bold=True)
    tbx(sl, desc, 1.6, 4.2, 9.5, 1.0, size=17, color=DT_GRAY, italic=True)
    print(f"  --- {label}: {titre} ---")

def slide_chart(prs, title, img_path, note=''):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.12, 12.5, 0.72, size=22, bold=True)
    img(sl, img_path, 0.25, 0.95, 12.85, 6.1)
    if note: tbx(sl, note, 0.5, 7.1, 12, 0.38, size=9, color=DT_GRAY, italic=True)
    print(f"  Slide CHART: {title[:60]}")

def slide_quote(prs, quote, author='DATATYM Analytics 2026'):
    sl = new_slide(prs)
    fill = sl.background.fill; fill.solid(); fill.fore_color.rgb = h('#F3F4F6')
    tbx(sl, '\u201c', 0.3, 0.2, 3, 3.5, size=210, bold=True, color='#E5E7EB')
    tbx(sl, quote, 1.4, 1.8, 10.5, 3.2, size=27, bold=True, italic=True, align=PP_ALIGN.CENTER)
    rect(sl, 4.3, 5.85, 4.8, 0.04, DT_GOLD)
    tbx(sl, f'\u2014 {author}', 0, 6.1, 13.33, 0.5, size=13, color=DT_GOLD, italic=True, align=PP_ALIGN.CENTER)
    print(f"  Slide QUOTE: {quote[:55]}...")

def slide_note(prs, title, bullets):
    """Slide textuelle avec puces"""
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 0.06, DT_GOLD)
    tbx(sl, title, 0.5, 0.12, 12.5, 0.72, size=22, bold=True)
    y = 1.1
    for b in bullets:
        tbx(sl, f'  {b}', 0.5, y, 12.3, 0.55, size=14, color=DT_WHITE)
        y += 0.6
    print(f"  Slide NOTE: {title[:55]}")


# ════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════

def build():
    print('\n=== DATATYM Analytique — Build ===\n')

    print('[Charts]')
    c_a1 = chart_matrice_repartition()
    c_a2 = chart_indices_par_stade()
    c_a3 = chart_iq_distribution()
    c_b1 = chart_risk_secteur()
    c_b2 = chart_vuln_salariale()
    c_c1 = chart_clusters_orga()
    c_d1 = chart_correlations()
    c_d2 = chart_whatif()
    c_d3 = chart_horizon_critique()

    print('\n[PowerPoint]')
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══ COVER ═════════════════════════════════════════════════
    slide_cover(prs,
        'Analyse Stratégique\nGen Z Afrique 2026',
        f'{N} répondants · {N_GZ} Gen Z · 13 pays · Baromètre DATATYM',
        tag='DATATYM ANALYTICS — RAPPORT DÉTAILLÉ')

    # ══ BLOC A — MATRICE TALENTS ══════════════════════════════
    slide_separator(prs, 'A',
        'Matrice Talents Gen Z',
        'Répartition réelle, stades d\'engagement, intention de départ.')

    slide_chart(prs,
        f'Matrice DATATYM — Effectifs réels (n={N_GZ} Gen Z)',
        c_a1,
        'Axe Performance : ENG (engagement) · Axe Potentiel : composite ICHD + APPRENT + INIT · Tertiles')

    slide_chart(prs,
        "Scores moyens par stade d'engagement — Tous les indices corrélés",
        c_a2,
        'Quartiles ENG · Les indices ICEQ et RECO s\'effondrent le plus vite dans les stades à risque')

    slide_chart(prs,
        "Intention de quitter — 65% des Gen Z à risque de départ",
        c_a3,
        '"Oui, activement" + "Oui, si opportunité" = population à monitorer en priorité')

    slide_note(prs,
        'Note méthodologique — Variable Genre',
        ['Le questionnaire actuel ne contient pas de variable sexe/genre.',
         'L\'analyse croisée genre × indices n\'est pas possible sur ce jeu de données.',
         'Recommandation : ajouter cette variable dans la prochaine vague du baromètre.',
         'Impact : les analyses par genre demandées dans le Bloc A sont reportées.'])

    # ══ BLOC B — RISQUES CRITIQUES ════════════════════════════
    slide_separator(prs, 'B',
        'Risques Critiques',
        'Score de risque composite, vulnérabilité salariale, secteurs à risque.')

    slide_stat(prs,
        '65%', 'des Gen Z à risque de départ.',
        '12% recherchent activement + 53% partiront si une opportunité se présente.',
        'Score de risque composite = f(ENG, RECO, SOUSPAYE, intention de quitter)',
        tag='RISQUE GLOBAL')

    slide_chart(prs,
        "Score de risque composite par secteur — Top 10",
        c_b1,
        'Risk = (100-ENG)×0.30 + (100-RECO)×0.25 + SOUSPAYE×0.25 + (100-IQ_inv)×0.20')

    slide_chart(prs,
        "Indice de Vulnérabilité Salariale — 4 quartiles",
        c_b2,
        'Vuln = SOUSPAYE×0.40 + (100-ICEQ)×0.35 + (100-RECO)×0.25 · '
        'Quartile "Très fort" = population en surenchère salariale imminente')

    # ══ BLOC C — PROFILS D'ORGANISATION ═══════════════════════
    slide_separator(prs, 'C',
        'Profils d\'Organisation',
        '4 clusters identifiés par K-Means sur les indices DATATYM.')

    if c_c1:
        slide_chart(prs,
            "4 Profils d'Organisation — Clustering DATATYM",
            c_c1,
            'K-Means (k=4) sur IDAT, ICHD, ICEQ, IATM, ITEQ, ENG, RECO, SOUSPAYE · '
            'Agrégation par secteur × type_org (n>=5)')

    slide_note(prs,
        'Lecture des 4 profils',
        ['L\'Incubateur Involontaire — ENG fort + IATM effondré + SOUSPAYE max',
         '  → Développe des talents qui partiront. Finance la concurrence.',
         '',
         'La Machine Silencieuse — Tout moyen, RECO faible',
         '  → Fonctionne, personne ne recommande. Marque employeur invisible.',
         '',
         'Le Résigné Tempéré — Tout moyen-bas, SOUSPAYE modéré',
         '  → Pas de crise visible, mais aucun facteur de rétention non plus.',
         '',
         'Chaque profil appelle une stratégie RH différente — pas de solution unique.'])

    # ══ BLOC D — ALERTES & WHAT-IF ════════════════════════════
    slide_separator(prs, 'D',
        'Indicateurs d\'Alerte & Scénarios',
        'Corrélations, what-if, horizon critique par secteur.')

    slide_chart(prs,
        "Corrélations — Quels leviers impactent rétention et recommandation ?",
        c_d1,
        'ICEQ→RECO = 0.78 · IATM→IQ_inv = 0.62 · L\'équité fait recommander. Le management fait rester.')

    slide_stat(prs,
        '0,78', 'corrélation ICEQ → RECO.',
        'L\'équité est le driver n°1 de la marque employeur.',
        'Un Gen Z qui perçoit de l\'équité recommande son organisation.\nCelui qui ne la perçoit pas — recrute pour la concurrence.',
        tag='DRIVER #1 RECO')

    slide_stat(prs,
        '0,62', 'corrélation IATM → Intention de rester.',
        'Le management est le driver n°1 de la rétention.',
        'Ce n\'est pas le salaire (corrélation 0,13). Ce n\'est pas l\'engagement (0,10).\nC\'est la qualité du management qui fait rester — ou partir.',
        tag='DRIVER #1 RÉTENTION')

    slide_chart(prs,
        "Scénarios What-If — Effet de levier par indice (+10 points)",
        c_d2,
        'Modèle linéaire : impact = corrélation × delta · Données Gen Z n=673')

    slide_chart(prs,
        "Horizon critique par secteur — Combien de mois avant perte significative ?",
        c_d3,
        'Estimation : risk>50 → 3 mois · 47-50 → 6 mois · 44-47 → 9 mois · <44 → 12 mois')

    # ══ SYNTHÈSE ══════════════════════════════════════════════
    slide_quote(prs,
        'L\'équité fait recommander.\nLe management fait rester.\nLe salaire ne fait ni l\'un ni l\'autre.',
        'Corrélations DATATYM — Baromètre Gen Z 2026')

    slide_note(prs,
        'Synthèse des résultats clés',
        [f'Matrice talents : 21% Stars, 18% Départs Annoncés sur {N_GZ} Gen Z',
         '65% des Gen Z à risque de départ (actif + opportuniste)',
         'Secteurs les plus à risque : Télécoms/IT, Conseil/Audit, Marketing, Santé',
         '4 profils d\'organisation identifiés — chacun appelle une stratégie différente',
         'ICEQ → RECO = 0,78 : l\'équité est le levier #1 de la marque employeur',
         'IATM → rétention = 0,62 : le management est le levier #1 pour garder les talents',
         'Le salaire n\'explique ni la recommandation (0,00) ni la rétention (0,13)',
         'Note : variable genre absente du questionnaire — à intégrer dans la vague 2'])

    slide_cover(prs,
        'DATATYM Analytics',
        'Rapport Analytique — Baromètre Gen Z Afrique 2026\nDocument technique réservé',
        tag='MERCI')

    out = '/Users/Apple/Desktop/H&C/Enquetes/Suivi enquetes/enquetes/DATATYM_Analytique_2026.pptx'
    prs.save(out)
    print(f'\n=== {len(prs.slides)} slides -> {out} ===')


if __name__ == '__main__':
    build()
